# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

src = '/Users/nfilho/Downloads/URT_UNIFIED_ROM_TEMPLATE_PT (3).pptx'
out = '/Users/nfilho/claude/CLARO_CPQ_ROM.pptx'

prs = Presentation(src)

# ── Helper: replace text in a shape preserving formatting ────────────────────
def replace_text(shape, old, new):
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)

def replace_all_text(slide, mapping):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for old, new in mapping.items():
                replace_text(shape, old, new)

def set_cell(tbl, row, col, text, bold=False, size=None, color=None):
    cell = tbl.cell(row, col)
    cell.text = ''
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    if size:
        run.font.size = Pt(size)
    if color:
        run.font.color.rgb = color

def add_bullet(tf, text, bold=False, size=9):
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)

def set_textbox(shape, text, bold=False, size=None):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # clear existing runs in first paragraph
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ''
    p = tf.paragraphs[0]
    if not p.runs:
        run = p.add_run()
    else:
        run = p.runs[0]
    run.text = text
    if bold:
        run.font.bold = True
    if size:
        run.font.size = Pt(size)

# ── SLIDE 1 — Capa ────────────────────────────────────────────────────────────
sl = prs.slides[0]
replace_all_text(sl, {
    'Inserir_Papel_Membro_Equipe_Salesforce_EM': 'Engagement Manager',
    'Inserir_Papel_Membro_Equipe_Salesforce_SBS': 'Solution Architect',
    'Inserir_Papel_Membro_Equipe_Salesforce_AP': 'Technical Architect',
    'Inserir_Papel_Membro_Equipe_Salesforce_Sponsor': 'Account Executive',
    'INSERIR MÊS E ANO': 'Junho 2026',
    '(Inserir Nome do Cliente)': 'Claro Brasil',
    'XXX Cloud': 'Communications Cloud',
    'XXX + Implementação YYYY': 'CPQ Design — Catálogo Unificado',
})

# ── SLIDE 3 — Disclaimer (keep as-is, replace logo placeholder only) ─────────
# no action needed

# ── SLIDE 4 — Agenda (keep structure) ────────────────────────────────────────
# no action needed

# ── SLIDE 7 — Missão do Cliente ───────────────────────────────────────────────
sl = prs.slides[6]
replace_all_text(sl, {
    'Entendemos que sua missão principal é (Inserir Missão do Cliente)...':
    'Entendemos que sua missão principal é posicionar a Claro Brasil como a operadora digital de vanguarda e líder em eficiência, por meio de um catálogo de produtos unificado e inteligente que elimina a complexidade sistêmica das matrizes legadas e automatiza o fluxo Lead-to-Cash — do produto ao dinheiro — em todas as cinco verticais de negócio.'
})

# ── SLIDE 8 — Motores Principais ─────────────────────────────────────────────
sl = prs.slides[7]
replace_all_text(sl, {
    '(Inserir Motor Principal #1)': '🚀  Redução drástica do Time-to-Market (TTM): lançar e atualizar combos e promoções com agilidade competitiva, eliminando a rigidez das matrizes customizadas legadas.',
    '(Inserir Motor Principal #2)': '🔗  Unificação do catálogo multi-LOB: consolidar Móvel, TV, Banda Larga, Fixo e Aparelhos em um único Shared Catalog aderente ao padrão internacional TM Forum SID.',
    '(Inserir Motor Principal #3)': '⚙️  Eliminação de débitos técnicos: substituir matrizes de preço e elegibilidade em código rígido pelo Business Rules Engine (BRE) declarativo nativo da Communications Cloud.',
})

# ── SLIDE 9 — Objetivos de Negócios (tabela) ──────────────────────────────────
sl = prs.slides[8]
tbl = None
for shape in sl.shapes:
    if shape.has_table:
        tbl = shape.table
        break
if tbl:
    data = [
        ('Reduzir o TTM de lançamento de ofertas em 45% em 12 meses',
         'Implementar o EPC com Object Types e herança de atributos, permitindo configuração de novas ofertas sem desenvolvimento de código.'),
        ('Atingir 100% das ofertas (~100 produtos) mapeadas no EPC em 3 meses',
         'Modelar as 5 LOBs no Shared Catalog TM Forum SID durante a Fase de Modelagem (semanas 4–7).'),
        ('Eliminar 100% das matrizes customizadas de precificação',
         'Migrar toda lógica tarifária para Expression Sets e Tabelas de Decisão declarativas do BRE.'),
        ('Reduzir erros no carrinho e abandono por incompatibilidade em 30%',
         'Implantar Advanced CPQ Validation Rules para elegibilidade por CEP e compatibilidade entre tecnologias em tempo real.'),
    ]
    for r, (obj, how) in enumerate(data, 1):
        if r < len(tbl.rows):
            set_cell(tbl, r, 0, obj, bold=True, size=9)
            set_cell(tbl, r, 1, how, size=9)

# ── SLIDE 10 — Desafios ───────────────────────────────────────────────────────
sl = prs.slides[9]
for shape in sl.shapes:
    if shape.has_text_frame:
        t = shape.text_frame.text.strip()
        if 'Dados de clientes fragmentados' in t:
            tf = shape.text_frame
            paragraphs_text = [
                '🔴  Silos organizacionais severos entre as 5 LOBs (Móvel, TV, Banda Larga, Fixo, Aparelhos) impedindo lançamento ágil de combos e sincronização de regras.',
                '🔴  Matrizes de elegibilidade e precificação customizadas em código rígido, fora do padrão internacional — altíssima complexidade de manutenção e débito técnico acumulado.',
                '🔴  Coexistência crítica com legados massivos de Billing e OSS/BSS sem arquitetura de integração documentada (Swagger/WSDL não fornecidos).',
                '🟡  Falta de separação nativa entre portfólio comercial (visível ao cliente) e especificações técnicas de rede (CFS), gerando inconsistências no provisionamento.',
                '🟡  Validação de elegibilidade regional por CEP realizada via chamadas síncronas a matrizes externas — sem SLA de resposta e risco de abandono de carrinho.',
            ]
            for i, para in enumerate(tf.paragraphs):
                if i < len(paragraphs_text):
                    for run in para.runs:
                        run.text = ''
                    if para.runs:
                        para.runs[0].text = paragraphs_text[i]
                    else:
                        run = para.add_run()
                        run.text = paragraphs_text[i]
            break

# ── SLIDE 13 — Visão do Futuro ────────────────────────────────────────────────
sl = prs.slides[12]
replace_all_text(sl, {
    'Imagine um futuro onde cada equipe é empoderada\ncom os dados certos no momento certo…':
    'Imagine um futuro onde um Analista de Negócios da Claro configura um novo combo multi-LOB em minutos, sem abrir um ticket para TI — apenas ajustando parâmetros no BRE.',
    'Onde cada interação com o cliente é inteligente,\npersonalizada e fluida...':
    'Onde o carrinho de compras valida automaticamente a viabilidade de Fibra pelo CEP, bloqueia incompatibilidades tecnológicas e calcula descontos de fidelidade em menos de 2 segundos.',
    'Onde a tecnologia não é mais uma barreira,\nmas o motor para seu crescimento e inovação…':
    'Onde cada pedido de combo é automaticamente decomposto em ordens técnicas para Billing e OSS/BSS — sem integração síncrona frágil, sem inconsistência de faturamento.',
})

# ── SLIDE 14 — Valor (tabela) ─────────────────────────────────────────────────
sl = prs.slides[13]
tbl = None
for shape in sl.shapes:
    if shape.has_table:
        tbl = shape.table
        break
if tbl:
    data = [
        ('Enterprise Product Catalog (EPC) — Shared Catalog TM Forum SID',
         'Unificação de 100% do portfólio das 5 LOBs em catálogo mestre centralizado',
         'Redução de 45% no TTM · 1 catálogo master (vs. múltiplos silos)'),
        ('Business Rules Engine (BRE) — Expression Sets Declarativos',
         'Eliminação das matrizes customizadas em código; autonomia para Negócios atualizar tarifas',
         '100% das regras via BRE · Redução de 50% no esforço de manutenção tarifária'),
        ('Advanced CPQ Validation Rules — Elegibilidade por CEP',
         'Validação em tempo real no carrinho: cobertura de fibra, compatibilidade de tecnologias',
         'Resposta < 2s · Redução de 30% no abandono de carrinho por incompatibilidade'),
        ('Order Decomposition Engine (OM) — Assíncrono',
         'Transformação automática de combos comerciais em ordens técnicas para Billing e OSS/BSS',
         'Taxa de erro de payload → 0% · Integridade transacional total'),
    ]
    for r, (cap, impact, kpi) in enumerate(data, 1):
        if r < len(tbl.rows):
            set_cell(tbl, r, 0, cap, bold=True, size=8)
            set_cell(tbl, r, 1, impact, size=8)
            set_cell(tbl, r, 2, kpi, size=8)

# ── SLIDE 17 — Arquitetura (mapa textual) ────────────────────────────────────
sl = prs.slides[16]
replace_all_text(sl, {
    'Inserir MAPA da Arquitetura':
    'CAMADA DE EXPERIÊNCIA\nCanais de Venda & Atendimento Claro BR (Físico + Digital)\n\n↓\n\nCAMADA DE PROCESSO\nCPQ Engine (Carrinho Inteligente) + Advanced Validation Rules\n\n↓\n\nCAMADA DE DADOS\nEnterprise Product Catalog (EPC) — Product2 + Object Types + CFS\n\n↓\n\nCAMADA DE INTELIGÊNCIA\nBusiness Rules Engine (BRE) — Expression Sets + Tabelas de Decisão\n\n↓\n\nCAMADA DE INTEGRAÇÃO\nMiddleware / MuleSoft — Barramento Corporativo (API REST Assíncrona)\n\n↓  ↓\n\nSISTEMAS EXTERNOS CORE\nBilling Core  |  OSS/BSS (Provisionamento de Rede)  |  Georreferenciamento (CEP)'
})

# ── SLIDE 18 — Capacidades Salesforce ────────────────────────────────────────
sl = prs.slides[17]
replace_all_text(sl, {
    'Exemplo: Data Cloud: Para unificar todos os dados de clientes em uma única fonte de verdade.\nExemplo: Sales Cloud: Para automatizar o processo de lead a fechamento.\nExemplo: Service Cloud: Para oferecer um atendimento ao cliente personalizado e multicanal.\nExemplo: MuleSoft: Para integrar perfeitamente com seus sistemas existentes como.':
    '📦  Communications Cloud — EPC: Modelagem do Shared Catalog com até 100 produtos comerciais (20/LOB) em duas camadas — Produto Comercial e Produto Técnico (CFS) — aderente ao padrão TM Forum SID.\n\n⚙️  Business Rules Engine (BRE): Substituição integral das matrizes customizadas legadas por Expression Sets e Tabelas de Decisão declarativas para cálculo de mensalidades, instalação e descontos de fidelidade de 12 meses.\n\n🛒  CPQ Engine — Advanced Validation Rules: Validação em tempo real de elegibilidade por CEP (viabilidade técnica de Fibra), compatibilidade entre tecnologias e planos antigos vs. novos.\n\n🔀  Order Decomposition Engine (OM): Transformação assíncrona de pedidos comerciais em Fulfillment Requests técnicos (CFS) para Billing e OSS/BSS via barramento corporativo MuleSoft.\n\n🔗  MuleSoft Anypoint Platform: Orquestração da Camada Anticorrupção (ACL) para integração com Billing, OSS/BSS e sistema de Georreferenciamento por CEP, implementando o padrão Strangler Fig de modernização incremental.'
})

# ── SLIDE 19 — Decisões Arquitetônicas ───────────────────────────────────────
sl = prs.slides[18]
replace_all_text(sl, {
    'Decisão: Usar o MuleSoft como uma plataforma de integração estratégica em vez de desenvolver integrações ponto a ponto. Justificativa: A complexidade de conectar-se a 16 DMSs heterogêneos mais um ecossistema de sistemas satélites torna a abordagem de integração ponto a ponto insustentável, frágil e cara de manter. Uma abordagem de conectividade liderada por API com o MuleSoft é um requisito técnico inegociável para gerenciar esse nível de complexidade, garantir governança de dados centralizada e possibilitar a agilidade futura requerida pela visão "Evolutiva" do projeto. Esta decisão transforma integrações de um passivo técnico para um ativo estratégico reutilizável.':
    '🔵  Decisão 1 — Shared Catalog TM Forum SID (Duas Camadas)\nJustificativa: Desacoplar o portfólio comercial das especificações técnicas de rede é o principal enabler de redução de TTM. O modelo monolítico legado foi descartado por sua rigidez e impossibilidade de evolução sem impacto em Billing e OSS/BSS.\n\n🔵  Decisão 2 — BRE Declarativo substituindo Matrizes Customizadas\nJustificativa: Matrizes em código Apex geram débito técnico permanente e dependência de TI para qualquer alteração tarifária. O BRE entrega autonomia ao Negócio e elimina deploys para atualizações de preço. Alternativa descartada: manutenção de código customizado Apex.\n\n🔵  Decisão 3 — OM Assíncrono via MuleSoft (Strangler Fig Pattern)\nJustificativa: Integrações síncronas ponto a ponto com Billing e OSS/BSS foram descartadas por fragibilidade transacional. A arquitetura assíncrona via MuleSoft como ACL garante integridade, escalabilidade e modernização incremental sem big bang. Sistema de Registro (SoR) migra gradualmente: Legado → Híbrido → Salesforce EPC (100%).\n\n🔵  Decisão 4 — SSO via Okta + Provisionamento SailPoint\nJustificativa: Conformidade com LGPD e ANATEL exige controle granular de acesso por LOB. OWD Privado + Permission Sets por vertical garantem isolamento de regras de precificação entre analistas de linhas concorrentes.'
})

# ── SLIDE 20 — Princípios Orientadores ───────────────────────────────────────
sl = prs.slides[19]
replace_all_text(sl, {
    'Exemplos:\nCliques, Não Código: Daremos prioridade à configuração declarativa para acelerar a entrega, reduzir custos de manutenção e capacitar sua equipe a fazer atualizações futuras.\nAbordagem Primeiro o Padrão: Ao aderir à funcionalidade pronta para uso, garantimos que sua plataforma seja escalável, segura e pronta para se beneficiar de 3 atualizações do Salesforce por ano.\nArquitetura Liderada por API: Projetaremos integrações para serem ativos reutilizáveis, prevenindo dívidas técnicas e criando uma base ágil para projetos futuros.':
    '1️⃣  Padrões Primeiro, Configurar antes de Personalizar: Maximização total de capacidades nativas (EPC, BRE, OM). Customizações em código rígido são terminantemente excluídas do escopo deste projeto.\n\n2️⃣  Aderência ao TM Forum SID: Cada decisão de modelagem de catálogo deve obedecer ao padrão internacional de telecomunicações — garantia de escalabilidade, interoperabilidade e longevidade da arquitetura.\n\n3️⃣  Arquitetura Desacoplada Orientada por API: Integrações via MuleSoft como Camada Anticorrupção (ACL) — cada fluxo é um ativo reutilizável, impedindo a reintrodução de débitos técnicos ponto a ponto.\n\n4️⃣  Entregar Valor Incrementalmente (Strangler Fig): Modernização em 4 fases controladas. Matrizes legadas são desativadas apenas após execução paralela (Parallel Run) validada — zero risco de big bang.\n\n5️⃣  Governar para o Futuro (CoE): Centro de Excelência híbrido com releases quinzenais para catálogo e trimestrais para mudanças estruturais. Modelo de suporte L1/L2/L3 com SLAs P1 < 4 horas.'
})

# ── SLIDE 22 — Resumo do Escopo ──────────────────────────────────────────────
sl = prs.slides[21]
replace_all_text(sl, {
    'Resumo do Escopo do Projeto...':
    'Resumo do Escopo do Projeto\n\nIniciativa: Design de Catálogo Unificado e Estrutura de CPQ — Communications Cloud\nCliente: Claro Brasil  |  Duração: 3 meses (12 semanas)  |  Modelo: T&M  |  Sizing: XL\n\nEsta iniciativa cobre o design técnico completo, a parametrização da PoC funcional e o Knowledge Transfer para as equipes de TI e Negócios da Claro Brasil. O escopo está restrito à fase de design e modelagem — não inclui implementação de produção, redesenho de orquestração pós-OM ou customizações em código.\n\nProdutos Salesforce: Communications Cloud (EPC + OM) · Business Rules Engine (BRE)\n\nCobertura: 5 LOBs (Móvel, TV, Banda Larga, Fixo, Aparelhos) · até 100 produtos comerciais · até 5 promoções/descontos · 1 PoC por LOB (5 total)'
})

# ── SLIDE 23 — Atividades e Entregáveis ──────────────────────────────────────
sl = prs.slides[22]
for shape in sl.shapes:
    if shape.has_text_frame:
        t = shape.text_frame.text.strip()
        if t == 'Atividades':
            tf = shape.text_frame
            for p in tf.paragraphs:
                for r in p.runs: r.text = ''
            if tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = 'Atividades por Fase'
            acts = [
                'F1 — Imersão (3 sem): workshops com TI e Negócios, mapeamento de requisitos, detalhamento funcional das ofertas vigentes por LOB',
                'F2 — Modelagem (4 sem): design EPC (Object Types, CFS, atributos, cardinalidade), modelagem BRE (Expression Sets, tabelas de decisão), desenho regras CPQ (elegibilidade CEP, compatibilidade)',
                'F3 — PoC Build (3 sem): configuração de 1 oferta real por LOB em Sandbox (prateleira + carrinho + decomposição OM)',
                'F4 — Knowledge Transfer (2 sem): workshops técnicos e de negócios, demonstração da PoC, entrega do Relatório Final',
            ]
            for act in acts:
                p2 = tf.add_paragraph()
                run = p2.add_run()
                run.text = act
                run.font.size = Pt(9)
        elif t == 'Entregáveis':
            tf = shape.text_frame
            for p in tf.paragraphs:
                for r in p.runs: r.text = ''
            if tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = 'Entregáveis'
            delivs = [
                '📄  Product Model Blueprint (PMB): diagrama visual completo da árvore de produtos — planos, serviços e descontos no Salesforce EPC',
                '📋  Documento de Recomendações de Regras de Negócio: especificação de elegibilidade, compatibilidade e precificação via BRE',
                '💻  Sandbox Configurada (PoC Funcional): 1 oferta por LOB com carrinho CPQ, validações e payloads de decomposição OM funcionais',
                '📊  Relatório Final do Projeto + Gravações dos Workshops de KT',
            ]
            for d in delivs:
                p2 = tf.add_paragraph()
                run = p2.add_run()
                run.text = d
                run.font.size = Pt(9)

# ── SLIDE 24 — Dentro/Fora do Escopo ─────────────────────────────────────────
sl = prs.slides[23]
for shape in sl.shapes:
    if shape.has_text_frame:
        t = shape.text_frame.text.strip()
        if t == 'Dentro do Escopo':
            tf = shape.text_frame
            for p in tf.paragraphs:
                for r in p.runs: r.text = ''
            if tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = '✅  Dentro do Escopo'
            items_in = [
                '• Modelagem EPC Shared Catalog — até 100 produtos comerciais (20/LOB) em duas camadas (Comercial + CFS Técnico)',
                '• Configuração de Bundles e Promoções — até 5 promoções/descontos: Fidelidade, Node, Dia das Mães',
                '• Definição de templates de produto (Object Type / Product Specification)',
                '• Modelagem de Atributos e Cardinalidade (velocidade, franquia, obrigatório/opcional)',
                '• BRE — Expression Sets e Tabelas de Decisão para mensalidade, instalação e fidelidade 12 meses',
                '• Regras CPQ — Elegibilidade por CEP (viabilidade de Fibra) e Compatibilidade entre tecnologias',
                '• Desenho da Decomposição de Pedidos (OM) — payloads CFS para Billing e OSS/BSS',
                '• PoC funcional: 1 oferta por LOB (prateleira, carrinho, decomposição) em Sandbox',
                '• Knowledge Transfer: workshops formais para TI e Negócios',
                '• Product Model Blueprint (PMB) e Documento de Recomendações',
            ]
            for item in items_in:
                p2 = tf.add_paragraph()
                run = p2.add_run()
                run.text = item
                run.font.size = Pt(8)
        elif t == 'Fora do Escopo':
            tf = shape.text_frame
            for p in tf.paragraphs:
                for r in p.runs: r.text = ''
            if tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = '❌  Fora do Escopo'
            items_out = [
                '• Implementação em produção (go-live)',
                '• Redesenho completo do fluxo de orquestração pós-decomposição OM',
                '• Desenvolvimento ou implementação de customizações em código Apex / hooks em métodos nativos',
                '• Testes de customizações legadas existentes baseadas em matrizes antigas',
                '• Saneamento, limpeza e carga massiva de dados históricos de ativos contratuais',
                '• Aquisição ou provisionamento de hardware/infraestrutura física',
                '• AMS (Application Managed Services) pós-projeto',
                '• Integração de canais de venda adicionais (além do escopo de validação CPQ)',
            ]
            for item in items_out:
                p2 = tf.add_paragraph()
                run = p2.add_run()
                run.text = item
                run.font.size = Pt(8)

# ── SLIDE 25 — Integrações ────────────────────────────────────────────────────
sl = prs.slides[24]
for shape in sl.shapes:
    if shape.has_table:
        tbl = shape.table
        integrations = [
            ('Barramento Corporativo (Middleware/MuleSoft)', 'Salesforce OM → Middleware', 'MuleSoft', 'REST Assíncrona', 'Fulfillment Request / Payload CFS'),
            ('Sistemas Core de Billing', 'Bidirecional (Contrato/Ativos)', 'MuleSoft', 'REST Síncrona', 'Installed Products / Contratos Vigentes'),
            ('OSS/BSS (Provisionamento de Rede)', 'Salesforce OM → OSS', 'MuleSoft', 'REST Assíncrona', 'Ordens de Ativação Técnica'),
            ('Sistema de Georreferenciamento (CEP)', 'Salesforce CPQ → Geo', 'MuleSoft', 'REST Síncrona', 'CEP / Indicadores de Cobertura Fibra'),
        ]
        for r, (sys, dir_, mw, api, obj) in enumerate(integrations, 1):
            if r < len(tbl.rows):
                set_cell(tbl, r, 0, sys, bold=True, size=8)
                set_cell(tbl, r, 1, dir_, size=8)
                set_cell(tbl, r, 2, mw, size=8)
                set_cell(tbl, r, 3, api, size=8)
                set_cell(tbl, r, 4, obj, size=8)
        break

# ── SLIDE 26 — Requisitos e Pressupostos ─────────────────────────────────────
sl = prs.slides[25]
for shape in sl.shapes:
    if shape.has_text_frame:
        t = shape.text_frame.text.strip()
        if t == 'Requisitos':
            tf = shape.text_frame
            for p in tf.paragraphs:
                for r in p.runs: r.text = ''
            if tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = '📌  Requisitos (Dependências da Claro Brasil)'
            reqs = [
                '• Sandbox funcional configurada disponibilizada pela Claro antes do início da Fase 3 (PoC Build) — pré-requisito inegociável.',
                '• Documentação técnica (Swagger/OpenAPI/WSDL) das APIs de Billing, OSS/BSS e Georreferenciamento entregue no início da Fase 1.',
                '• Definição formal dos Product Owners com autonomia de aprovação para cada LOB (Móvel, TV, Banda Larga, Fixo, Aparelhos).',
                '• Disponibilização de amostras das planilhas/matrizes de elegibilidade e precificação vigentes para mapeamento BRE.',
                '• Participação ativa de analistas de TI e Negócios nos workshops de Imersão e KT.',
            ]
            for req in reqs:
                p2 = tf.add_paragraph()
                run = p2.add_run()
                run.text = req
                run.font.size = Pt(8)
        elif t == 'Pressupostos':
            tf = shape.text_frame
            for p in tf.paragraphs:
                for r in p.runs: r.text = ''
            if tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = '⚠️  Pressupostos e Premissas'
            presups = [
                '• O volume total de ofertas avaliadas não excederá 100 itens complexos (até 20 por LOB).',
                '• A modelagem proposta baseada em capacidades nativas declarativas eliminará a necessidade de customizações em código.',
                '• O barramento corporativo possui estabilidade e capacidade técnica para receber mensagens assíncronas de OM.',
                '• O cronograma de 3 meses é fixo e orientado a entregáveis — extensões dependem de aprovação formal.',
                '• As regras de precificação vigentes serão fornecidas em formato legível (planilha, diagrama) no início da Fase 2.',
                '• Documentação de Billing/OSS/BSS não fornecida até o início da Imersão é risco de extensão de escopo.',
            ]
            for presup in presups:
                p2 = tf.add_paragraph()
                run = p2.add_run()
                run.text = presup
                run.font.size = Pt(8)

# ── SLIDE 27 — Roteiro Visual ────────────────────────────────────────────────
sl = prs.slides[26]
replace_all_text(sl, {
    'Fase 1\nAgente de Suporte a Negócios para Análise de Viabilidade Comercial & Agente de Catálogo de Produtos para Consulta Técnica e Comercial de Produtos':
    'Fase 1 — Imersão & Discovery\nSemanas 1–3\nWorkshops com TI e Negócios · Mapeamento de requisitos · Detalhamento funcional das ofertas vigentes por LOB · Entregável: Documento de Requisitos Funcionais',
    'Fase 2\nxxx':
    'Fase 2 — Modelagem (PMB)\nSemanas 4–7\nDesign EPC (Object Types, CFS, atributos) · Modelagem BRE (Expression Sets) · Regras CPQ · Entregável: Product Model Blueprint + Documento BRE',
    'Fase 3\nxxx':
    'Fase 3 — PoC Build\nSemanas 8–10\nConfiguração da Sandbox · 1 oferta por LOB (prateleira, carrinho, decomposição OM) · Validação de payloads CFS · Entregável: Sandbox Configurada + Relatório de Testes',
    'Início': 'Jun/2026',
    'Objetivo': 'Fase 4 — KT\nSemanas 11–12\nWorkshops formais TI e Negócios · Demonstração PoC · Entregável: Relatório Final + PoC Funcional Homologada',
})

# ── SLIDE 28 — Roteiro Detalhado ─────────────────────────────────────────────
sl = prs.slides[27]
replace_all_text(sl, {
    'Fase 1: Visão 360 Fundamentada (Semanas 1-8)\nResultado: Seus agentes de serviço terão uma visão unificada do histórico do cliente, permitindo um atendimento mais informado e personalizado desde o primeiro dia.':
    'Fase 1: Imersão & Discovery (Semanas 1–3)\nResultado: Visão completa das ofertas vigentes, mapeamento de requisitos funcionais e técnicos, alinhamento de POs por LOB. Base sólida para o Product Model Blueprint.',
    'Fase 2: Automação Proativa de Vendas (Semanas 9-16)\nResultado: Sua equipe de vendas se beneficiará do roteamento automático de leads e gestão de tarefas, liberando tempo para se concentrar em vendas.':
    'Fase 2: Modelagem — Product Model Blueprint (Semanas 4–7)\nResultado: PMB completo com EPC estruturado em duas camadas (Comercial + CFS), Object Types definidos, Expression Sets BRE parametrizados e regras CPQ de elegibilidade e compatibilidade desenhadas.',
    'Fase 3: Atendimento ao Vivo & Realização de Valor (Semana 17+)\nResultado: A solução completa está ao vivo, com acompanhamento contínuo do valor em relação aos KPIs definidos no início de nosso envolvimento.':
    'Fase 3: PoC Build & Fase 4: KT (Semanas 8–12)\nResultado: Sandbox configurada com 1 oferta real por LOB demonstrando prateleira CPQ, validações em tempo real e decomposição OM funcional. Equipes de TI e Negócios capacitadas e aprovação formal da PoC assinada pelos POs das 5 verticais.',
})

# ── SLIDE 30 — Cronograma ─────────────────────────────────────────────────────
sl = prs.slides[29]
replace_all_text(sl, {
    '[Insira um gráfico de cronograma de alto nível mostrando as principais fases e suas durações estimadas.]\nFase 1: Fundação (Semanas 1-8) Descoberta, Design, Configuração da Plataforma Central\nFase 2: Automação de Vendas (Semanas 9-16) Construir, Testar, Implantar Processo de Vendas\nFase 3: Serviço & Suporte (Semanas 17-24) Construir, Testar, Implantar Consola de Serviços':
    'CRONOGRAMA — CLARO CPQ DESIGN (12 SEMANAS / 3 MESES)\n\n├── Semanas 1–3   │ FASE 1: Imersão & Discovery\n│                │ Workshops TI + Negócios · Mapeamento de requisitos · Detalhamento das ofertas vigentes\n│\n├── Semanas 4–7   │ FASE 2: Modelagem — Product Model Blueprint\n│                │ Design EPC · Object Types · Expression Sets BRE · Regras CPQ · Decomposição OM\n│\n├── Semanas 8–10  │ FASE 3: PoC Build\n│                │ Sandbox configurada · 1 oferta/LOB · Prateleira + Carrinho + OM · Testes de payloads CFS\n│\n└── Semanas 11–12 │ FASE 4: Knowledge Transfer\n                 │ Workshops formais · Demonstração PoC · Relatório Final · Aceite formal dos POs\n\nPRÉ-REQUISITO CRÍTICO: Sandbox funcional disponibilizada pela Claro antes da Semana 8.\nDEPENDÊNCIA CRÍTICA: Documentação de APIs Billing/OSS/BSS entregue até a Semana 1.'
})

# ── SLIDE 41 — Equipe ─────────────────────────────────────────────────────────
sl = prs.slides[40]
replace_all_text(sl, {
    'Nossa Equipe de Projeto Combinada': 'Nossa Equipe de Projeto Combinada',
    'Equipo Projeto': 'SALESFORCE PS\n\nEngagement Manager\n→ Coordenação geral, governança e interface com o cliente\n\nSolution Architect (Communications Cloud CMT)\n→ Design do EPC, BRE, Object Types e Decomposição OM\n\nTechnical Architect\n→ Integrações MuleSoft, ACL, contratos de APIs\n\nBusiness Analyst\n→ Workshops de Imersão, documentação funcional por LOB\n\nQA / DevOps Engineer\n→ Estratégia de testes, VBT, Copado/Gearset, ambientes\n\n─────────────────────────\nCLARO BRASIL (Cliente)\n\nBusiness Sponsor (Diretoria)\nProduct Owners por LOB (5)\nArquiteto de TI / Catálogo\nAnalistas de Negócios por Vertical\nGestor de Transição / Release'
})

# ── SLIDE 42 — RACI ───────────────────────────────────────────────────────────
sl = prs.slides[41]
for shape in sl.shapes:
    if shape.has_table:
        tbl = shape.table
        raci_data = [
            ('Definição de Requisitos Funcionais (por LOB)', 'Facilita & Documenta', 'Fornece Especialistas & Aprova'),
            ('Design do Catálogo EPC (Object Types, CFS)', 'Projeta & Configura', 'Valida & Fornece Feedback'),
            ('Parametrização do BRE (Expression Sets)', 'Projeta & Configura', 'Fornece Dados de Preços & Aprova'),
            ('Configuração da PoC em Sandbox', 'Constrói & Testa', 'Executa UAT & Assina Aceite'),
            ('Disponibilização da Sandbox', 'Orienta Requisitos', 'Provisiona & Entrega'),
        ]
        for r, (act, sf, client) in enumerate(raci_data, 1):
            if r < len(tbl.rows):
                set_cell(tbl, r, 0, act, size=8)
                set_cell(tbl, r, 1, sf, size=8)
                set_cell(tbl, r, 2, client, size=8)
        break

# ── SLIDE 50 — Investimento ───────────────────────────────────────────────────
sl = prs.slides[49]
for shape in sl.shapes:
    if shape.has_table:
        tbl = shape.table
        set_cell(tbl, 1, 0, 'Honorários Estimados de Serviços Profissionais', bold=True, size=9)
        set_cell(tbl, 1, 1, 'A confirmar via SOW formal  |  Faixa ROM: nível XL, ~3 meses T&M', size=9)
        set_cell(tbl, 2, 0, 'Modelo de Contratação', bold=True, size=9)
        set_cell(tbl, 2, 1, 'Tempo & Materiais (T&M)', size=9)
        set_cell(tbl, 3, 0, 'Despesas de Viagem & Deslocamento (T&E)', bold=True, size=9)
        set_cell(tbl, 3, 1, 'Não incluídas nos honorários — a definir conforme necessidade de presença on-site', size=9)
        break

replace_all_text(sl, {
    'Validade: Esta Estimativa Bruta de Magnitude é válida até.':
    'Validade: Esta Estimativa Bruta de Magnitude é válida por 30 dias a partir da data de entrega.',
})

# ── SLIDE 52 — Próximos Passos ────────────────────────────────────────────────
sl = prs.slides[51]
replace_all_text(sl, {
    '1. Revisão e Alinhamento do ROM\nAtividade: Revisar em conjunto esta apresentação para alinhar escopo, abordagem e investimento.\nResultado: Acordo mútuo para prosseguir.':
    '1. Revisão e Alinhamento do ROM\nAtividade: Revisar esta apresentação em conjunto com as lideranças de TI e Negócios da Claro para alinhar escopo das 5 LOBs, abordagem e investimento.\nResultado: Acordo mútuo para prosseguir e identificação dos POs por vertical.',
    '2. Finalização do Escopo & SOW\nAtividade: Refinar os detalhes do escopo e desenvolver a Declaração Formal de Trabalho (SOW).\nResultado: Um SOW assinado.':
    '2. Finalização do Escopo & SOW\nAtividade: Refinar escopo com base no alinhamento — em especial: volumetria de APIs Billing/OSS/BSS, disponibilidade da Sandbox e formalização dos POs por LOB.\nResultado: SOW assinado.',
    '3. Início do Projeto\nAtividade: Início formal do projeto com a equipe do projeto combinada.\nResultado: Nossa jornada de transformação começa.':
    '3. Início do Projeto\nAtividade: Kick-off formal com a equipe combinada Salesforce + Claro Brasil. Início da Fase 1 — Imersão & Discovery.\nResultado: Primeira sessão de workshops com TI e Negócios em até 2 semanas após assinatura da SOW.',
})

# ── Replace all remaining logo placeholders ───────────────────────────────────
for sl in prs.slides:
    replace_all_text(sl, {'Inserir LOGO do Cliente': 'Claro Brasil'})

prs.save(out)
print(f"Salvo: {out}")
print(f"Total slides: {len(prs.slides)}")
