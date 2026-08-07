from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── Paleta Salesforce ───────────────────────────────────────────────────────
SF_BLUE       = RGBColor(0x00, 0x96, 0xFF)   # #0096FF
SF_DARK_BLUE  = RGBColor(0x03, 0x2D, 0x60)   # #032D60
SF_LIGHT_BLUE = RGBColor(0xD4, 0xEE, 0xFF)   # #D4EEFF
SF_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
SF_GRAY       = RGBColor(0x7F, 0x8C, 0x8D)
SF_LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF9)
SF_DARK_GRAY  = RGBColor(0x32, 0x3E, 0x48)
SF_GREEN      = RGBColor(0x2E, 0x7D, 0x32)
SF_ORANGE     = RGBColor(0xFF, 0x6B, 0x00)
SF_TEAL       = RGBColor(0x00, 0x6D, 0x9A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # blank

# ─── Helpers ─────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        if not fill:
            shape.line.fill.background()
        else:
            shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=14, bold=False, color=SF_DARK_GRAY,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_para(tf, text, size=12, bold=False, color=SF_DARK_GRAY,
             align=PP_ALIGN.LEFT, space_before=None, bullet=False):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = _Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = _Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p

def slide_header(slide, title, subtitle=None, accent=SF_BLUE):
    # topo azul
    add_rect(slide, 0, 0, 13.33, 0.08, fill=accent)
    add_text(slide, title, 0.4, 0.15, 10, 0.55,
             size=22, bold=True, color=SF_DARK_BLUE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.65, 10, 0.35,
                 size=13, color=SF_GRAY)
    # linha divisória
    add_rect(slide, 0.4, 1.05, 12.53, 0.03, fill=accent)

def slide_footer(slide, num, total):
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill=SF_DARK_BLUE)
    add_text(slide, "DATAPREV  |  SISDIP / DFT  |  Proposta Técnica — Salesforce Professional Services",
             0.3, 7.12, 11, 0.3, size=9, color=SF_WHITE)
    add_text(slide, f"{num} / {total}", 12.2, 7.12, 1, 0.3,
             size=9, color=SF_WHITE, align=PP_ALIGN.RIGHT)

TOTAL_SLIDES = 14

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=SF_DARK_BLUE)
add_rect(s, 0, 0, 13.33, 0.12, fill=SF_BLUE)
add_rect(s, 0, 7.38, 13.33, 0.12, fill=SF_BLUE)

# marca lateral decorativa
add_rect(s, 0, 0.12, 0.35, 7.26, fill=SF_BLUE)

# título
add_text(s, "SISDIP / DFT", 0.7, 1.6, 11, 0.7,
         size=36, bold=True, color=SF_WHITE, align=PP_ALIGN.LEFT)
add_text(s, "Sistema de Dimensionamento da Força de Trabalho", 0.7, 2.3, 10, 0.55,
         size=20, color=SF_LIGHT_BLUE, align=PP_ALIGN.LEFT)
add_text(s, "Proposta Técnica — Salesforce Professional Services", 0.7, 3.0, 10, 0.45,
         size=15, bold=True, color=SF_BLUE, align=PP_ALIGN.LEFT)

add_rect(s, 0.7, 3.6, 5.5, 0.04, fill=SF_BLUE)

add_text(s, "Cliente: DATAPREV / MGI — Ministério da Gestão e Inovação", 0.7, 3.8, 11, 0.38,
         size=13, color=SF_LIGHT_BLUE)
add_text(s, "Julho 2026", 0.7, 4.2, 4, 0.35,
         size=13, color=SF_GRAY)
add_text(s, "Confidencial — uso exclusivo DATAPREV / MGI", 0.7, 6.6, 10, 0.35,
         size=10, italic=True, color=SF_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=SF_LIGHT_GRAY)
slide_header(s, "Agenda", "O que vamos abordar nesta proposta")
slide_footer(s, 2, TOTAL_SLIDES)

items = [
    ("01", "Contexto e Desafio",          "Situação atual do SISDIP e oportunidade de melhoria"),
    ("02", "Solução Proposta",             "Arquitetura Salesforce end-to-end para o DFT"),
    ("03", "Integrações e Ecossistema",    "Conexão com 8 sistemas federais prioritários"),
    ("04", "Mapeamento de Requisitos",     "20 requisitos × produto Salesforce"),
    ("05", "Arquitetura de Agentes IA",    "Agentforce — automação inteligente do DFT"),
    ("06", "Abordagem de Entrega",         "Metodologia PS Salesforce — fases e governança"),
    ("07", "Por que Salesforce",           "Diferenciais vs. abordagem alternativa"),
    ("08", "Próximos Passos",              "Alinhamentos necessários para avançar"),
]

cols = [
    (0.4, 5.5),
    (6.9, 5.5),
]
for i, (num, title, desc) in enumerate(items):
    col = i % 2
    row = i // 2
    lx = cols[col][0]
    ly = 1.3 + row * 1.3
    add_rect(s, lx, ly, 0.55, 0.55, fill=SF_BLUE)
    add_text(s, num, lx, ly + 0.04, 0.55, 0.5, size=16, bold=True,
             color=SF_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, lx + 0.65, ly, 4.7, 0.35,
             size=13, bold=True, color=SF_DARK_BLUE)
    add_text(s, desc, lx + 0.65, ly + 0.33, 4.7, 0.4,
             size=10, color=SF_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CONTEXTO E DESAFIO
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=SF_LIGHT_GRAY)
slide_header(s, "Contexto e Desafio", "Situação atual do SISDIP e a oportunidade")
slide_footer(s, 3, TOTAL_SLIDES)

# box esquerda — situação atual
add_rect(s, 0.4, 1.25, 5.9, 5.5, fill=SF_WHITE)
add_rect(s, 0.4, 1.25, 5.9, 0.45, fill=SF_ORANGE)
add_text(s, "⚠  Situação Atual", 0.5, 1.28, 5.7, 0.38,
         size=13, bold=True, color=SF_WHITE)

pontos_atual = [
    "Dados de força de trabalho fragmentados em 8+ plataformas federais",
    "SISDIP opera de forma isolada, sem integração em tempo real",
    "Processos manuais: análise de perfis, mapeamento de cargos, movimentação",
    "Alta latência decisória: consolidar informações leva horas ou dias",
    "Risco de inconsistência e retrabalho entre sistemas (PGD, SEI, Recruta, PEI…)",
    "Sem rastreabilidade integrada das decisões de dimensionamento",
]
for i, pt in enumerate(pontos_atual):
    add_text(s, f"•  {pt}", 0.55, 1.85 + i * 0.73, 5.6, 0.65,
             size=11, color=SF_DARK_GRAY)

# box direita — oportunidade
add_rect(s, 6.9, 1.25, 5.9, 5.5, fill=SF_WHITE)
add_rect(s, 6.9, 1.25, 5.9, 0.45, fill=SF_BLUE)
add_text(s, "✔  Oportunidade", 7.0, 1.28, 5.7, 0.38,
         size=13, bold=True, color=SF_WHITE)

pontos_oport = [
    "Centralizar orquestração do DFT no SISDIP via Salesforce",
    "Integrar PGD, SEI, PEI, Recruta e outros em tempo real via MuleSoft",
    "Automatizar análise de perfil × entrega com Agentforce IA",
    "Oferecer painéis executivos unificados com Tableau",
    "Garantir LGPD, soberania de dados e rastreabilidade com Shield",
    "Escalar para múltiplos órgãos sem replicação de dados",
]
for i, pt in enumerate(pontos_oport):
    add_text(s, f"•  {pt}", 7.05, 1.85 + i * 0.73, 5.6, 0.65,
             size=11, color=SF_DARK_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SOLUÇÃO PROPOSTA (visão geral)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=SF_LIGHT_GRAY)
slide_header(s, "Solução Proposta", "SISDIP como nó central do ecossistema DFT — powered by Salesforce")
slide_footer(s, 4, TOTAL_SLIDES)

# camadas da solução
camadas = [
    (SF_DARK_BLUE, "Experiência & Colaboração",
     "Experience Cloud  ·  Interface analista/gestor  ·  Slack (notificações e aprovações)"),
    (SF_BLUE,      "Inteligência & Automação",
     "Agentforce (2 agentes)  ·  Einstein AI  ·  Flow Orchestration  ·  Einstein Trust Layer"),
    (SF_TEAL,      "Plataforma & Dados",
     "Service Cloud  ·  Data Cloud (multi-órgão)  ·  Salesforce Shield  ·  Platform"),
    (SF_GREEN,     "Integração",
     "MuleSoft Anypoint Platform  ·  APIs somente-leitura SISDIP  ·  Conectores PGD / SEI / Recruta / PEI"),
    (SF_DARK_GRAY, "Análise & Relatórios",
     "Tableau  ·  Painéis gerenciais  ·  Relatórios exportáveis  ·  Rastreabilidade DFT"),
]

for i, (cor, titulo, desc) in enumerate(camadas):
    ly = 1.25 + i * 1.05
    add_rect(s, 0.4, ly, 12.5, 0.95, fill=cor)
    add_text(s, titulo, 0.55, ly + 0.05, 4.5, 0.38,
             size=13, bold=True, color=SF_WHITE)
    add_text(s, desc, 0.55, ly + 0.45, 12.1, 0.42,
             size=11, color=SF_LIGHT_BLUE if cor != SF_LIGHT_GRAY else SF_DARK_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — INTEGRAÇÕES E ECOSSISTEMA
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=SF_LIGHT_GRAY)
slide_header(s, "Integrações e Ecossistema", "8 sistemas federais conectados via MuleSoft Anypoint Platform")
slide_footer(s, 5, TOTAL_SLIDES)

# SISDIP central
add_rect(s, 5.2, 2.9, 2.9, 1.0, fill=SF_DARK_BLUE)
add_text(s, "SISDIP", 5.2, 2.92, 2.9, 0.5,
         size=18, bold=True, color=SF_WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Salesforce Platform", 5.2, 3.4, 2.9, 0.45,
         size=10, color=SF_LIGHT_BLUE, align=PP_ALIGN.CENTER)

integ = [
    # (label, prioridade, l, t)
    ("PGD\nEntregas e planos\nde trabalho",       "P1", 0.3,  1.1),
    ("SEI\nPortarias e processos\nadministrativos","P2", 0.3,  3.3),
    ("PEI\nPlanejamento\nestratégico",             "P3", 0.3,  5.4),
    ("Power BI\nPainéis analíticos\nexistentes",   "P4", 3.5,  1.1),
    ("Carreiras\nCompatíveis",                     "P5", 9.2,  1.1),
    ("Perfil\nProfissiográfico",                   "P6", 9.2,  3.3),
    ("Recruta\nMovimentação\nde servidores",        "P7", 9.2,  5.4),
    ("Observatório\nde Pessoal",                   "P8", 3.5,  5.4),
]

for label, prio, lx, ly in integ:
    add_rect(s, lx, ly, 2.8, 1.15, fill=SF_WHITE)
    add_rect(s, lx, ly, 2.8, 0.28, fill=SF_BLUE)
    add_text(s, prio, lx + 0.05, ly + 0.01, 0.5, 0.25,
             size=10, bold=True, color=SF_WHITE)
    lines = label.split("\n")
    add_text(s, lines[0], lx + 0.08, ly + 0.32, 2.6, 0.3,
             size=12, bold=True, color=SF_DARK_BLUE)
    if len(lines) > 1:
        add_text(s, "\n".join(lines[1:]), lx + 0.08, ly + 0.62, 2.6, 0.45,
                 size=9, color=SF_GRAY)

add_text(s, "★  Carta de Serviços (Portal)\nFront-end interativo — coleta estruturada de entregas dos gestores",
         3.5, 3.3, 5.8, 0.8, size=10, color=SF_DARK_BLUE,
         align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MAPEAMENTO DE REQUISITOS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=SF_LIGHT_GRAY)
slide_header(s, "Mapeamento de Requisitos", "20 requisitos mapeados para produtos Salesforce")
slide_footer(s, 6, TOTAL_SLIDES)

# cabeçalho tabela
cols_w = [0.65, 4.8, 3.7, 2.5]
cols_x = [0.4, 1.1, 5.95, 9.7]
headers = ["REQ", "Necessidade", "Produto Salesforce", "Observação"]
for i, (hdr, cx, cw) in enumerate(zip(headers, cols_x, cols_w)):
    add_rect(s, cx, 1.2, cw, 0.38, fill=SF_DARK_BLUE)
    add_text(s, hdr, cx + 0.05, 1.22, cw - 0.1, 0.32,
             size=10, bold=True, color=SF_WHITE, align=PP_ALIGN.CENTER)

reqs = [
    ("01", "Acesso federado SISDIP sem replicação",         "MuleSoft Anypoint Platform",            "✅ Contratado"),
    ("02", "Catálogo de cargos por órgão",                  "Salesforce Platform / App Builder",      "Licença padrão"),
    ("03", "Registro e gestão de resultados com status",    "Service Cloud – Agentforce 1 Ed.",       "✅ Contratado"),
    ("04", "Acionamento automático nova entrega",           "Flow Orchestration (nativo)",            "✅ Incluso"),
    ("05", "IA analisa entrega → sugere cargo + justif.",   "Agentforce – Einstein 1 Ed.",            "⚠ Verificar"),
    ("06", "Refinamento iterativo sugestão IA",             "Agentforce (mesmo agente)",              "⚠ Verificar"),
    ("07", "Gestão de status do processo de análise",       "Service Cloud – Cases / Flows",          "✅ Contratado"),
    ("08", "Interface revisão e aprovação do analista",     "Experience Cloud / Lightning App",        "⚠ Verificar"),
    ("09", "Resumo executivo por órgão sob demanda",        "Agentforce (2º agente)",                 "⚠ Verificar"),
    ("10", "Painel gerencial em tempo real",                "Tableau",                                "✅ Contratado"),
    ("14", "Relatórios exportáveis (Excel/PDF)",            "Tableau / Reports & Dashboards",         "✅ Contratado"),
    ("15", "Análise em até 60 segundos",                    "Agentforce + Flow",                      "✅ Incluso"),
    ("16", "Processamento simultâneo múltiplos órgãos",     "Data Cloud + Flow",                      "⚠ Verificar qtd"),
    ("17", "Rastreabilidade e versionamento",               "Salesforce Shield",                      "✅ Contratado"),
    ("18", "Controle de acesso por perfis",                 "Shield + Permission Sets",               "✅ Contratado"),
    ("19", "LGPD + soberania de dados",                     "Einstein Trust Layer + Data Cloud",      "✅ Incluso"),
    ("20", "Acesso SISDIP somente leitura",                 "MuleSoft (read-only API)",               "✅ Contratado"),
]

row_h = 0.28
for i, (req, need, prod, obs) in enumerate(reqs):
    ly = 1.62 + i * row_h
    bg = SF_WHITE if i % 2 == 0 else SF_LIGHT_GRAY
    add_rect(s, 0.4, ly, 12.55, row_h, fill=bg)
    add_text(s, req,  cols_x[0] + 0.05, ly + 0.02, cols_w[0] - 0.1, row_h - 0.04,
             size=8.5, bold=True, color=SF_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, need, cols_x[1] + 0.05, ly + 0.02, cols_w[1] - 0.1, row_h - 0.04,
             size=8.5, color=SF_DARK_GRAY)
    add_text(s, prod, cols_x[2] + 0.05, ly + 0.02, cols_w[2] - 0.1, row_h - 0.04,
             size=8.5, color=SF_DARK_GRAY)
    obs_color = SF_GREEN if "✅" in obs else (SF_ORANGE if "⚠" in obs else SF_DARK_GRAY)
    add_text(s, obs,  cols_x[3] + 0.05, ly + 0.02, cols_w[3] - 0.1, row_h - 0.04,
             size=8.5, bold=("✅" in obs or "⚠" in obs), color=obs_color)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ARQUITETURA DE AGENTES IA
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=SF_LIGHT_GRAY)
slide_header(s, "Arquitetura de Agentes IA", "Agentforce — 2 agentes para automatizar o ciclo DFT")
slide_footer(s, 7, TOTAL_SLIDES)

# Agente 1
add_rect(s, 0.4, 1.25, 5.9, 5.6, fill=SF_WHITE)
add_rect(s, 0.4, 1.25, 5.9, 0.5, fill=SF_BLUE)
add_text(s, "Agente 1 — Análise de Cargo", 0.55, 1.27, 5.6, 0.42,
         size=14, bold=True, color=SF_WHITE)

ag1 = [
    ("Trigger",    "Nova entrega registrada no SISDIP via Flow Orchestration"),
    ("Input",      "Dados da entrega + Perfil profissiográfico + Catálogo de cargos"),
    ("Processamento","Einstein AI analisa entrega e correlaciona com cargos compatíveis"),
    ("Output",     "Sugestão de cargo + justificativa fundamentada"),
    ("Refinamento","Analista revisa, ajusta e aprova via Experience Cloud"),
    ("Rastreabilidade","Decisão versionada com Shield; auditável por órgão/servidor"),
    ("SLA",        "Análise entregue em até 60 segundos (Agentforce + Flow)"),
]
for i, (lbl, txt) in enumerate(ag1):
    ly = 1.9 + i * 0.68
    add_rect(s, 0.45, ly, 1.3, 0.5, fill=SF_LIGHT_BLUE)
    add_text(s, lbl, 0.48, ly + 0.05, 1.22, 0.4,
             size=9, bold=True, color=SF_DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, txt, 1.85, ly + 0.05, 4.3, 0.5,
             size=10, color=SF_DARK_GRAY)

# Agente 2
add_rect(s, 6.9, 1.25, 5.9, 5.6, fill=SF_WHITE)
add_rect(s, 6.9, 1.25, 5.9, 0.5, fill=SF_DARK_BLUE)
add_text(s, "Agente 2 — Resumo Executivo", 7.05, 1.27, 5.6, 0.42,
         size=14, bold=True, color=SF_WHITE)

ag2 = [
    ("Trigger",    "Solicitação do gestor (sob demanda) ou schedule automático"),
    ("Escopo",     "Consolidação por órgão: totais, perfis faltantes, gaps críticos"),
    ("Fonte",      "Data Cloud agrega dados de múltiplos órgãos sem replicação"),
    ("Output",     "Relatório executivo pronto para envio a dirigentes do MGI"),
    ("Canal",      "Slack: notificação automática para gestores e tomadores de decisão"),
    ("Formato",    "Exportável em PDF/Excel via Tableau Reports"),
    ("Governança", "Einstein Trust Layer garante LGPD em todo processamento"),
]
for i, (lbl, txt) in enumerate(ag2):
    ly = 1.9 + i * 0.68
    add_rect(s, 6.95, ly, 1.3, 0.5, fill=SF_LIGHT_BLUE)
    add_text(s, lbl, 6.98, ly + 0.05, 1.22, 0.4,
             size=9, bold=True, color=SF_DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, txt, 8.35, ly + 0.05, 4.3, 0.5,
             size=10, color=SF_DARK_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — GOVERNANÇA E LGPD
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=SF_LIGHT_GRAY)
slide_header(s, "Governança, Segurança e LGPD",
             "Conformidade nativa — nenhum dado sensível sai do ambiente controlado")
slide_footer(s, 8, TOTAL_SLIDES)

pilares = [
    (SF_DARK_BLUE, "Einstein Trust Layer",
     "Camada de governança nativa do Agentforce. Todo prompt e resposta de IA passa por filtros de privacidade, mascaramento de dados sensíveis e auditoria antes de qualquer ação."),
    (SF_BLUE, "Salesforce Shield",
     "Criptografia de campo (Field Audit Trail), rastreabilidade de acesso e monitoramento de eventos. Atende LGPD Art. 46-49: medidas técnicas para proteção de dados pessoais."),
    (SF_TEAL, "Data Cloud — Zero Copy",
     "Dados do SISDIP e sistemas federais são consultados in-loco via MuleSoft, sem replicação nem armazenamento redundante. Soberania de dados preservada nos sistemas de origem."),
    (SF_GREEN, "Controle de Acesso",
     "Permission Sets granulares por perfil (analista, gestor, dirigente, TI). Integração com identidade federal via SSO. Acesso mínimo necessário por papel e órgão."),
    (SF_ORANGE, "Auditoria e Rastreabilidade",
     "Cada decisão de dimensionamento é versionada, com registro de quem fez, quando, com qual dado e qual resultado. Histórico completo para prestação de contas ao TCU/CGU."),
    (SF_DARK_GRAY, "Conformidade Contratual",
     "Salesforce mantém certificações ISO 27001, SOC 2 Type II e FedRAMP-equivalente. Contrato DTP já vigente garante SLA e responsabilidades de tratamento de dados."),
]

for i, (cor, titulo, desc) in enumerate(pilares):
    col = i % 2
    row = i // 2
    lx = 0.4 + col * 6.5
    ly = 1.25 + row * 1.95
    add_rect(s, lx, ly, 6.1, 1.75, fill=SF_WHITE)
    add_rect(s, lx, ly, 0.18, 1.75, fill=cor)
    add_text(s, titulo, lx + 0.28, ly + 0.12, 5.7, 0.38,
             size=13, bold=True, color=cor)
    add_text(s, desc, lx + 0.28, ly + 0.52, 5.7, 1.1,
             size=10, color=SF_DARK_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ABORDAGEM DE ENTREGA (fases)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=SF_LIGHT_GRAY)
slide_header(s, "Abordagem de Entrega",
             "Metodologia Salesforce PS — entrega iterativa e incremental em 5 fases")
slide_footer(s, 9, TOTAL_SLIDES)

fases = [
    ("F1", "Descoberta\ne Alinhamento",
     "Workshops de discovery\nMapeamento de integrações\nDefinição de MVP\nKick-off com stakeholders"),
    ("F2", "Arquitetura\ne Design",
     "Design da solução SISDIP\nProtótipo Agentforce\nDefinição MuleSoft APIs\nAprovação arquitetura"),
    ("F3", "Construção\nCore",
     "Impl. Service Cloud\nAgentes Agentforce\nIntegrações MuleSoft\nExperience Cloud"),
    ("F4", "Dados &\nInteligência",
     "Data Cloud multi-órgão\nTableau dashboards\nRefinamento IA\nShield & LGPD"),
    ("F5", "Testes,\nGo-Live & UAT",
     "QA / Testes integração\nUAT com DATAPREV\nTreinamento usuários\nHandover & suporte"),
]

fase_colors = [SF_BLUE, SF_TEAL, SF_DARK_BLUE, SF_GREEN, SF_ORANGE]

for i, (cod, titulo, bullets) in enumerate(fases):
    lx = 0.4 + i * 2.5
    add_rect(s, lx, 1.25, 2.35, 0.6, fill=fase_colors[i])
    add_text(s, cod, lx + 0.05, 1.27, 0.6, 0.52,
             size=18, bold=True, color=SF_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, titulo, lx + 0.7, 1.27, 1.55, 0.52,
             size=11, bold=True, color=SF_WHITE)
    add_rect(s, lx, 1.88, 2.35, 4.85, fill=SF_WHITE)
    for j, bl in enumerate(bullets.split("\n")):
        add_text(s, f"• {bl}", lx + 0.12, 2.0 + j * 0.58, 2.15, 0.5,
                 size=10, color=SF_DARK_GRAY)

# seta de progressão
add_text(s, "→", 2.7, 1.4, 0.25, 0.35, size=14, bold=True, color=SF_GRAY, align=PP_ALIGN.CENTER)
add_text(s, "→", 5.2, 1.4, 0.25, 0.35, size=14, bold=True, color=SF_GRAY, align=PP_ALIGN.CENTER)
add_text(s, "→", 7.7, 1.4, 0.25, 0.35, size=14, bold=True, color=SF_GRAY, align=PP_ALIGN.CENTER)
add_text(s, "→", 10.2, 1.4, 0.25, 0.35, size=14, bold=True, color=SF_GRAY, align=PP_ALIGN.CENTER)

add_text(s, "★  Metodologia Salesforce PS Delivery Framework — com Quality Gates ao final de cada fase",
         0.4, 6.9, 12.5, 0.35, size=10, italic=True, color=SF_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — PRODUTOS E PLATAFORMA
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=SF_LIGHT_GRAY)
slide_header(s, "Produtos e Plataforma",
             "Portfólio Salesforce utilizado na solução SISDIP / DFT")
slide_footer(s, 10, TOTAL_SLIDES)

produtos = [
    (SF_BLUE,      "Service Cloud\nAgentforce 1 Edition",
     "Gestão de casos e processos do DFT. Base operacional do SISDIP. Workflows, aprovações e automações nativas."),
    (SF_DARK_BLUE, "Agentforce\nEinstein AI",
     "Motor de inteligência. Agente de Análise de Cargo e Agente de Resumo Executivo. Einstein Trust Layer para LGPD."),
    (SF_TEAL,      "MuleSoft Anypoint\nPlatform",
     "Camada de integração. Conecta SISDIP a PGD, SEI, PEI, Recruta e outros sistemas legados. Acesso read-only sem replicação."),
    (SF_GREEN,     "Data Cloud",
     "Consolidação de dados de múltiplos órgãos sem copiar. Habilita análises cross-órgão e segmentação por carreira/perfil."),
    (SF_ORANGE,    "Experience Cloud",
     "Portal do analista e do gestor. Interface de revisão, aprovação e acompanhamento do processo de dimensionamento."),
    (SF_DARK_GRAY, "Tableau",
     "Painéis gerenciais e relatórios executivos. Indicadores de força de trabalho por órgão, carreira e perfil profissiográfico."),
    (SF_BLUE,      "Salesforce Shield",
     "Criptografia, auditoria de campo e monitoramento de eventos. Rastreabilidade para CGU/TCU. Conformidade LGPD nativa."),
    (SF_TEAL,      "Slack",
     "Canal de notificação e colaboração. Alertas de análise concluída, aprovações pendentes e resumos executivos para gestores."),
]

for i, (cor, titulo, desc) in enumerate(produtos):
    col = i % 4
    row = i // 4
    lx = 0.4 + col * 3.1
    ly = 1.25 + row * 2.7
    add_rect(s, lx, ly, 2.85, 2.5, fill=SF_WHITE)
    add_rect(s, lx, ly, 2.85, 0.5, fill=cor)
    lines = titulo.split("\n")
    add_text(s, lines[0], lx + 0.1, ly + 0.05, 2.65, 0.28,
             size=11, bold=True, color=SF_WHITE)
    if len(lines) > 1:
        add_text(s, lines[1], lx + 0.1, ly + 0.28, 2.65, 0.22,
                 size=9, color=SF_LIGHT_BLUE)
    add_text(s, desc, lx + 0.1, ly + 0.6, 2.65, 1.75,
             size=10, color=SF_DARK_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — DIFERENCIAIS SALESFORCE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=SF_LIGHT_GRAY)
slide_header(s, "Por que Salesforce",
             "Plataforma unificada com infraestrutura já contratada pela DATAPREV")
slide_footer(s, 11, TOTAL_SLIDES)

diferenciais = [
    ("Plataforma já contratada",
     "Service Cloud, Tableau, MuleSoft e Shield estão no contrato DTP vigente. Menor risco, menor tempo de procurement e aceleração de entrega."),
    ("Agentforce — IA governada",
     "Únicos a oferecer agentes de IA com controle determinístico, rastreabilidade nativa e Einstein Trust Layer para LGPD — sem expor dados em modelos públicos."),
    ("Zero Copy para dados federais",
     "Data Cloud acessa dados do SISDIP, PGD, SEI e outros sem replicar ou mover. Soberania de dados nos sistemas de origem — requisito crítico para o setor público."),
    ("Ecossistema integrado nativo",
     "Uma única plataforma para CRM, integração (MuleSoft), IA (Agentforce), dados (Data Cloud), analytics (Tableau) e colaboração (Slack). Sem pontos de falha entre sistemas."),
    ("Entrega incremental com PS",
     "Time Salesforce Professional Services entrega em fases com Quality Gates. Conhecimento do produto e do cliente em um único time — sem intermediários."),
    ("Escalabilidade multi-órgão",
     "Arquitetura desenhada desde o início para operar com múltiplos ministérios simultaneamente, com isolamento de dados por órgão e governança centralizada."),
]

for i, (titulo, desc) in enumerate(diferenciais):
    col = i % 2
    row = i // 2
    lx = 0.4 + col * 6.5
    ly = 1.25 + row * 1.95
    add_rect(s, lx, ly, 6.1, 1.75, fill=SF_WHITE)
    # ícone numerado
    add_rect(s, lx, ly, 0.5, 1.75, fill=SF_BLUE)
    add_text(s, str(i + 1), lx + 0.02, ly + 0.6, 0.46, 0.5,
             size=18, bold=True, color=SF_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, titulo, lx + 0.6, ly + 0.12, 5.4, 0.38,
             size=13, bold=True, color=SF_DARK_BLUE)
    add_text(s, desc, lx + 0.6, ly + 0.52, 5.4, 1.1,
             size=10, color=SF_DARK_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — PREMISSAS E DEPENDÊNCIAS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=SF_LIGHT_GRAY)
slide_header(s, "Premissas e Dependências",
             "Itens necessários para garantir o sucesso da entrega")
slide_footer(s, 12, TOTAL_SLIDES)

add_rect(s, 0.4, 1.25, 5.9, 5.6, fill=SF_WHITE)
add_rect(s, 0.4, 1.25, 5.9, 0.45, fill=SF_BLUE)
add_text(s, "Premissas do Projeto", 0.55, 1.27, 5.6, 0.38,
         size=13, bold=True, color=SF_WHITE)

premissas = [
    "DATAPREV indica Product Owner dedicado (mín. 50% do tempo)",
    "Ambientes Salesforce (sandbox + produção) disponíveis no início da F1",
    "Confirmação de licenças: Agentforce Builder, Experience Cloud e Data Cloud",
    "APIs dos sistemas legados (PGD, SEI, Recruta) documentadas e disponíveis em sandbox até F2",
    "Aprovação de arquitetura técnica ao final da F2 antes de iniciar construção",
    "Janelas de change freeze informadas com antecedência mínima de 4 semanas",
    "Time cliente disponível para UAT na F5 (mín. 3 semanas, 50% capacidade)",
    "Acesso ao ambiente de testes dos sistemas legados para validação das integrações",
]
for i, pt in enumerate(premissas):
    add_text(s, f"•  {pt}", 0.55, 1.85 + i * 0.6, 5.6, 0.55,
             size=10.5, color=SF_DARK_GRAY)

add_rect(s, 6.9, 1.25, 5.9, 5.6, fill=SF_WHITE)
add_rect(s, 6.9, 1.25, 5.9, 0.45, fill=SF_ORANGE)
add_text(s, "Dependências Externas", 7.05, 1.27, 5.6, 0.38,
         size=13, bold=True, color=SF_WHITE)

deps = [
    "PGD — API de consulta de planos e entregas de trabalho (leitura)",
    "SEI — Webservice de portarias e processos administrativos",
    "PEI — Endpoint de metas e objetivos estratégicos por órgão",
    "Recruta — API de movimentação e redistribuição de servidores",
    "Perfil Profissiográfico — Base de competências e perfis técnicos",
    "Carreiras Compatíveis — Tabela de equivalência de cargos (SIAPE-integrado)",
    "Observatório de Pessoal — Indicadores macro (leitura periódica/batch)",
    "SSO corporativo DATAPREV — Federação de identidade para acesso ao portal",
]
for i, pt in enumerate(deps):
    add_text(s, f"•  {pt}", 7.05, 1.85 + i * 0.6, 5.6, 0.55,
             size=10.5, color=SF_DARK_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — PRÓXIMOS PASSOS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=SF_LIGHT_GRAY)
slide_header(s, "Próximos Passos",
             "Alinhamentos necessários para avançar à proposta formal")
slide_footer(s, 13, TOTAL_SLIDES)

passos = [
    ("1", "Confirmar licenças no contrato DTP",
     "AE Fernanda confirma disponibilidade de Agentforce Builder, Experience Cloud e Data Cloud no contrato vigente.",
     "Fernanda Rodrigues (AE)", "Imediato"),
    ("2", "Validar disponibilidade de APIs legadas",
     "DATAPREV confirma documentação e sandbox de PGD, SEI, PEI e Recruta.",
     "DATAPREV / TI MGI", "Semana 1–2"),
    ("3", "Alinhar janelas de change freeze",
     "Identificar períodos de bloqueio nos próximos 4–6 meses para planejamento do cronograma.",
     "DATAPREV", "Semana 1–2"),
    ("4", "Definir Product Owner do cliente",
     "DATAPREV indica representante técnico dedicado para participar ativamente das fases de discovery e UAT.",
     "DATAPREV / MGI", "Semana 2"),
    ("5", "Workshop de Discovery — Kick-off",
     "Sessão presencial ou remota com stakeholders para detalhar requisitos, priorizar integrações e confirmar arquitetura.",
     "Salesforce PS + DATAPREV", "Semana 2–3"),
]

for i, (num, titulo, desc, resp, prazo) in enumerate(passos):
    ly = 1.25 + i * 1.1
    add_rect(s, 0.4, ly, 12.5, 1.0, fill=SF_WHITE)
    add_rect(s, 0.4, ly, 0.55, 1.0, fill=SF_BLUE)
    add_text(s, num, 0.42, ly + 0.22, 0.5, 0.5,
             size=18, bold=True, color=SF_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, titulo, 1.05, ly + 0.08, 7.5, 0.32,
             size=12, bold=True, color=SF_DARK_BLUE)
    add_text(s, desc, 1.05, ly + 0.42, 7.5, 0.5,
             size=10, color=SF_DARK_GRAY)
    add_rect(s, 8.75, ly + 0.08, 2.0, 0.38, fill=SF_LIGHT_BLUE)
    add_text(s, resp, 8.8, ly + 0.1, 1.9, 0.35,
             size=9, bold=True, color=SF_DARK_BLUE, align=PP_ALIGN.CENTER)
    add_rect(s, 10.85, ly + 0.08, 1.95, 0.38, fill=SF_BLUE)
    add_text(s, prazo, 10.9, ly + 0.1, 1.85, 0.35,
             size=9, bold=True, color=SF_WHITE, align=PP_ALIGN.CENTER)

# cabeçalho colunas
add_text(s, "Responsável", 8.75, 1.1, 2.0, 0.2, size=8.5, bold=True,
         color=SF_GRAY, align=PP_ALIGN.CENTER)
add_text(s, "Prazo", 10.85, 1.1, 1.95, 0.2, size=8.5, bold=True,
         color=SF_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — ENCERRAMENTO
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=SF_DARK_BLUE)
add_rect(s, 0, 0, 13.33, 0.12, fill=SF_BLUE)
add_rect(s, 0, 7.38, 13.33, 0.12, fill=SF_BLUE)
add_rect(s, 0, 0.12, 0.35, 7.26, fill=SF_BLUE)

add_text(s, "Transformando o DFT com", 0.7, 1.5, 11, 0.55,
         size=24, color=SF_LIGHT_BLUE, align=PP_ALIGN.LEFT)
add_text(s, "Inteligência e Integração", 0.7, 2.05, 11, 0.65,
         size=32, bold=True, color=SF_WHITE, align=PP_ALIGN.LEFT)
add_text(s, "Salesforce Professional Services", 0.7, 2.72, 11, 0.45,
         size=16, bold=True, color=SF_BLUE, align=PP_ALIGN.LEFT)

add_rect(s, 0.7, 3.35, 4.5, 0.04, fill=SF_BLUE)

mensagens = [
    "Plataforma já contratada — menor risco e menor custo de implantação",
    "Agentforce como diferencial: IA governada, auditável e compatível com LGPD",
    "Integração nativa com os 8 sistemas federais via MuleSoft já vigente",
    "Entrega incremental com Quality Gates — visibilidade total do progresso",
]
for i, msg in enumerate(mensagens):
    add_text(s, f"✔  {msg}", 0.7, 3.55 + i * 0.52, 11, 0.45,
             size=13, color=SF_WHITE)

add_text(s, "Próximo passo: Workshop de Discovery — DATAPREV + Salesforce PS",
         0.7, 5.8, 11, 0.4, size=14, bold=True, color=SF_BLUE)

add_text(s, "salesforce.com  |  Salesforce Professional Services LATAM  |  Julho 2026",
         0.7, 6.6, 11, 0.35, size=10, italic=True, color=SF_GRAY)

# ─── Salvar ──────────────────────────────────────────────────────────────────
OUTPUT = "/Users/nfilho/claude/DATAPREV_SGP_Proposta_Tecnica.pptx"
prs.save(OUTPUT)
print(f"PPTX gerado: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
