from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Cores ──────────────────────────────────────────
SF_DARK   = RGBColor(0x03,0x2D,0x60)
SF_BLUE   = RGBColor(0x00,0x70,0xD2)
SF_TEAL   = RGBColor(0x00,0xA1,0xE0)
SF_GREEN  = RGBColor(0x2E,0x84,0x4A)
SF_ORANGE = RGBColor(0xE6,0x50,0x00)
SF_YELLOW = RGBColor(0xFF,0xC8,0x49)
SF_RED    = RGBColor(0xC2,0x3B,0x22)
SF_PURPLE = RGBColor(0x6B,0x37,0xBF)
WHITE     = RGBColor(0xFF,0xFF,0xFF)
LGRAY     = RGBColor(0xF3,0xF3,0xF3)
MGRAY     = RGBColor(0xD8,0xDC,0xE2)
DGRAY     = RGBColor(0x55,0x55,0x55)
DARK      = RGBColor(0x1A,0x1A,0x1A)

# ── Helpers ────────────────────────────────────────
def rect(slide, l,t,w,h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
    s.line.width = lw
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line; s.line.width = lw or Pt(1)
    else: s.line.fill.background()
    return s

def txt(slide, text, l,t,w,h, sz=11, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size=Pt(sz); r.font.bold=bold
    r.font.italic=italic; r.font.color.rgb=color
    return tb

def header(slide, title, sub, accent=SF_BLUE, logo=""):
    rect(slide,0,0,13.33,1.18,fill=SF_DARK)
    rect(slide,0,1.18,13.33,0.05,fill=accent)
    txt(slide,"ANATEL — Inovação Digital · Defesa Técnica OTT",
        0.3,0.07,10,0.28,sz=9,color=RGBColor(0xAA,0xBB,0xCC))
    txt(slide,title,0.3,0.32,10.5,0.58,sz=21,bold=True,color=WHITE)
    if sub:
        txt(slide,sub,0.3,0.88,10.5,0.27,sz=10,
            color=RGBColor(0xAA,0xCC,0xEE),italic=True)
    if logo:
        rect(slide,12.45,0.18,0.65,0.65,fill=accent)
        txt(slide,logo,12.45,0.27,0.65,0.35,sz=9,bold=True,
            color=WHITE,align=PP_ALIGN.CENTER)

def footer(slide, pg):
    rect(slide,0,7.25,13.33,0.25,fill=SF_DARK)
    txt(slide,"Salesforce Professional Services — LATAM  |  Uso Interno  |  Confidencial",
        0.3,7.27,8,0.2,sz=8,color=RGBColor(0x77,0x88,0x99))
    txt(slide,pg,9,7.27,4,0.2,sz=8,
        color=RGBColor(0x77,0x88,0x99),align=PP_ALIGN.RIGHT)

def section_label(slide, label, color, l,t,w=1.5,h=0.26):
    rect(slide,l,t,w,h,fill=color)
    txt(slide,label,l+0.06,t+0.02,w-0.12,h-0.04,
        sz=8,bold=True,color=WHITE,align=PP_ALIGN.CENTER)

def block_card(slide, l,t,w,h, title, items, accent, bg=None):
    bg = bg or LGRAY
    rect(slide,l,t,w,h,fill=bg,line=accent,lw=Pt(1))
    rect(slide,l,t,w,0.3,fill=accent)
    txt(slide,title,l+0.1,t+0.03,w-0.15,0.24,sz=9,bold=True,color=WHITE)
    for i,item in enumerate(items):
        iy = t+0.38+i*0.27
        if iy+0.24 > t+h: break
        rect(slide,l+0.12,iy+0.05,0.07,0.07,fill=accent)
        txt(slide,item,l+0.26,iy,w-0.35,0.26,sz=8.5,color=DARK,wrap=True)

# ══════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ══════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
rect(s1,0,0,13.33,7.5,fill=SF_DARK)
rect(s1,9.2,-0.3,5,5,fill=RGBColor(0x02,0x1F,0x44))
rect(s1,10,4.2,4,4,fill=RGBColor(0x01,0x18,0x35))
rect(s1,0,6.8,13.33,0.7,fill=RGBColor(0x02,0x20,0x45))
# Faixa colorida
for i,(c,lbl) in enumerate([(SF_BLUE,"PSS"),(SF_TEAL,"SC+AF"),
                              (SF_GREEN,"DC+MS"),(SF_ORANGE,"RevCloud"),
                              (SF_PURPLE,"MC"),(SF_RED,"Shield")]):
    rect(s1,0.4+i*1.1,5.4,1.0,0.32,fill=c)
    txt(s1,lbl,0.4+i*1.1,5.44,1.0,0.24,sz=8,bold=True,
        color=WHITE,align=PP_ALIGN.CENTER)

txt(s1,"ANATEL",0.4,1.4,12,0.55,sz=13,color=RGBColor(0x66,0x99,0xCC))
txt(s1,"Inovação Digital",0.4,1.88,12,0.95,sz=40,bold=True,color=WHITE)
txt(s1,"Modernização da Regulação de Telecomunicações via Salesforce",
    0.4,2.82,11,0.45,sz=16,color=RGBColor(0xAA,0xCC,0xEE))
txt(s1,"Defesa Técnica OTT — Proposta de Solução por Cloud",
    0.4,3.28,11,0.35,sz=12,color=RGBColor(0x77,0x99,0xBB),italic=True)
rect(s1,0.4,3.82,5,0.04,fill=SF_TEAL)
txt(s1,"Salesforce Professional Services — LATAM  |  Junho 2026",
    0.4,3.98,9,0.3,sz=10,color=RGBColor(0x88,0x99,0xAA))
txt(s1,"Programa XL · 5 Fases · 18 Meses · Go-live Final Nov/2027",
    0.4,4.3,9,0.3,sz=10,color=RGBColor(0x66,0x88,0xAA),italic=True)

# ══════════════════════════════════════════════════════
# SLIDE 2 — VISÃO GERAL DO PROGRAMA
# ══════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
header(s2,"Visão Geral do Programa",
       "5 fases · 18 meses · 8 clouds · Greenfield · Prazo regulatório crítico TFF mar/2027",
       SF_BLUE)
footer(s2,"Slide 2 de 9")
rect(s2,0,1.23,13.33,5.8,fill=LGRAY)

fases = [
    ("F0","Fundação\n& Data Model","Abr–Jun\n2026","10 sem",SF_DARK,
     "Data model unificado\nDEV/QA/PROD\nDashboard 6 gerências"),
    ("F1","Omnichannel\nAgêntico","Jun–Out\n2026","16 sem",SF_BLUE,
     "Service Cloud\nAgentforce v1\nWhatsApp oficial"),
    ("F2","TFF/TFI\nArrecadação","Out/26–Mar\n2027","21 sem",SF_ORANGE,
     "MuleSoft 4 fontes\nData Cloud 10M reg.\nJornada anti-inadimp."),
    ("F3","MMAR\nLicenciamento","Jan–Ago\n2027","18 sem",SF_GREEN,
     "PSS + OmniStudio\nAgentforce v3\nPortal requerente"),
    ("F4","Ouvidoria 360°\n& Inteligência","Ago–Nov\n2027","13 sem",SF_PURPLE,
     "CRM Analytics\nAgentforce avançado\nTransparência pública"),
]

fw = 2.28
fx_start = 0.3
for i,(fn,fname,fdates,fwks,fc,fdet) in enumerate(fases):
    fx = fx_start + i*(fw+0.12)
    rect(s2,fx,1.35,fw,5.6,fill=WHITE,line=fc,lw=Pt(1.5))
    rect(s2,fx,1.35,fw,0.46,fill=fc)
    txt(s2,fn,fx+0.08,1.37,0.4,0.42,sz=14,bold=True,color=WHITE)
    txt(s2,fname,fx+0.48,1.38,fw-0.55,0.42,sz=10,bold=True,color=WHITE,wrap=True)
    fc_light = {SF_DARK:RGBColor(0x05,0x4A,0x8C), SF_BLUE:RGBColor(0x00,0x96,0xFF),
                SF_ORANGE:RGBColor(0xFF,0x72,0x22), SF_GREEN:RGBColor(0x45,0xA5,0x6B),
                SF_PURPLE:RGBColor(0x89,0x57,0xDC)}
    rect(s2,fx,1.81,fw,0.5,fill=fc_light.get(fc,SF_TEAL))
    txt(s2,fdates,fx+0.08,1.83,fw*0.55,0.22,sz=8,bold=True,color=WHITE)
    txt(s2,fwks,fx+fw*0.55,1.83,fw*0.42,0.22,sz=9,bold=True,
        color=WHITE,align=PP_ALIGN.RIGHT)
    # Detalhes
    for j,line in enumerate(fdet.split("\n")):
        ly = 2.44+j*0.38
        rect(s2,fx+0.12,ly+0.06,0.08,0.08,fill=fc)
        txt(s2,line,fx+0.26,ly,fw-0.34,0.36,sz=9,color=DARK)

# Clouds por fase
clouds_map = [
    ["Data Cloud\nShield"],
    ["Service Cloud\nAgentforce\nWhatsApp"],
    ["MuleSoft\nData Cloud\nMarketing Cloud"],
    ["PSS\nOmniStudio\nRevenue Cloud\nCLM"],
    ["CRM Analytics\nAgentforce\nShield"],
]
for i,(fc_info,(_,_,_,_,fc,_)) in enumerate(zip(clouds_map,fases)):
    fx = fx_start + i*(fw+0.12)
    rect(s2,fx,5.12,fw,0.04,fill=fc)
    for j,cl in enumerate(fc_info):
        rect(s2,fx+0.08+j*0.72,5.2,0.65,0.5,fill=fc)
        txt(s2,cl,fx+0.08+j*0.72,5.22,0.65,0.46,
            sz=7,bold=True,color=WHITE,align=PP_ALIGN.CENTER,wrap=True)

# Prazo crítico
rect(s2,3.52,6.55,3.2,0.52,fill=SF_RED)
txt(s2,"⚠  PRAZO REGULATÓRIO CRÍTICO\nTFF go-live: início Mar/2027 · 31/Mar = deadline não negociável",
    3.6,6.57,3.1,0.46,sz=8,bold=True,color=WHITE,wrap=True)

# ══════════════════════════════════════════════════════
# SLIDE 3 — SERVICE CLOUD + AGENTFORCE
# ══════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
header(s3,"Service Cloud + Agentforce",
       "Fases 1 e 4 · Atendimento omnichannel · Agentes autônomos · Ouvidoria 360°",SF_BLUE,"SC+AF")
footer(s3,"Slide 3 de 9")
rect(s3,0,1.23,13.33,5.8,fill=LGRAY)

# Contexto negócio
rect(s3,0.2,1.35,8.5,1.38,fill=WHITE,line=SF_BLUE,lw=Pt(1))
rect(s3,0.2,1.35,8.5,0.3,fill=SF_BLUE)
txt(s3,"CONTEXTO DE NEGÓCIO",0.3,1.37,8.3,0.24,sz=9,bold=True,color=WHITE)
txt(s3,"A ANATEL recebe demandas de cidadãos e entidades reguladas via múltiplos canais sem visão unificada. "
       "O atendimento hoje é baseado em chatbots de árvore fixa — sem resolução autônoma, sem contexto do histórico "
       "e sem SLA rastreável. 6 gerências operam de forma isolada, gerando retrabalho e baixa satisfação do cidadão. "
       "A Ouvidoria processa reclamações manualmente sem dados preditivos.",
    0.3,1.72,8.3,0.9,sz=9.5,color=DARK,wrap=True)

# KPIs dor
for i,(v,l,c) in enumerate([("4M+","Tickets/ano sem triagem autônoma",SF_RED),
                              ("6","Gerências isoladas sem fila unificada",SF_ORANGE),
                              ("0%","Resolução autônoma atual",SF_ORANGE),
                              ("SLA","Sem rastreamento por caso",SF_RED)]):
    bx = 9.0+i*1.0 if i<2 else 9.0+(i-2)*1.0
    by = 1.35 if i<2 else 2.08
    rect(s3,bx,by,0.9,0.62,fill=c)
    txt(s3,v,bx,by+0.03,0.9,0.32,sz=14,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(s3,l,bx+0.04,by+0.35,0.84,0.25,sz=7,color=WHITE,wrap=True,align=PP_ALIGN.CENTER)

# Como atenderemos
rect(s3,0.2,2.82,5.9,4.06,fill=WHITE,line=SF_BLUE,lw=Pt(1))
rect(s3,0.2,2.82,5.9,0.3,fill=SF_BLUE)
txt(s3,"COMO ATENDEREMOS — SOLUÇÃO SALESFORCE",0.3,2.84,5.7,0.24,sz=9,bold=True,color=WHITE)
sol_items = [
    ("Service Cloud","Fila unificada de atendimento para as 6 gerências · Vista 360° do cidadão · "
     "SLA automático por tipo de caso · Gestão de casos da Ouvidoria"),
    ("Agentforce v1 (Fase 1)","Triagem autônoma via WhatsApp: status de licença, 2ª via de boleto, prazo de análise · "
     "Transbordo humanizado com resumo automático do contexto (Einstein Conversation Insights)"),
    ("Agentforce v2 (Fase 2)","Reemissão autônoma de boleto TFF via agente · Integração com MuleSoft para consulta em tempo real"),
    ("Agentforce v3 (Fase 3)","Status e pendências MMAR · Fluxo guiado de solicitação de outorga via WhatsApp"),
    ("Ouvidoria 360° (Fase 4)","Visão consolidada de todos os módulos · SLA rigoroso · Relatório de transparência pública · "
     "Integração com CRM Analytics para análise preditiva"),
]
for i,(cloud,desc) in enumerate(sol_items):
    iy = 3.2+i*0.72
    rect(s3,0.28,iy,5.74,0.64,fill=RGBColor(0xEB,0xF4,0xFF),line=SF_BLUE,lw=Pt(0.5))
    rect(s3,0.28,iy,0.08,0.64,fill=SF_BLUE)
    txt(s3,cloud,0.44,iy+0.04,5.4,0.22,sz=9,bold=True,color=SF_BLUE)
    txt(s3,desc,0.44,iy+0.26,5.38,0.35,sz=8,color=DGRAY,wrap=True)

# Entregas e fora do escopo
rect(s3,6.3,2.82,3.3,4.06,fill=WHITE,line=SF_GREEN,lw=Pt(1))
rect(s3,6.3,2.82,3.3,0.3,fill=SF_GREEN)
txt(s3,"PRINCIPAIS ENTREGAS",6.4,2.84,3.1,0.24,sz=9,bold=True,color=WHITE)
entregas = ["Fila unificada 6 gerências","WhatsApp oficial integrado",
            "Agentforce v1 → v3 progressivo","Transbordo humanizado com resumo",
            "Portal do cidadão (Experience)","SLA tracking por caso",
            "Ouvidoria com visão 360°","Dashboard executivo tempo real"]
for i,e in enumerate(entregas):
    iy = 3.2+i*0.41
    rect(s3,6.38,iy+0.07,0.08,0.08,fill=SF_GREEN)
    txt(s3,e,6.52,iy,3.0,0.38,sz=8.5,color=DARK)

rect(s3,9.8,2.82,3.33,4.06,fill=WHITE,line=SF_RED,lw=Pt(1))
rect(s3,9.8,2.82,3.33,0.3,fill=SF_RED)
txt(s3,"FORA DO ESCOPO",9.9,2.84,3.1,0.24,sz=9,bold=True,color=WHITE)
fora = ["CTI / PABX / telefonia legada","Integração com SEI / SIPAC",
        "Histórico de atendimentos anteriores","Automações financeiras (Fase 1)",
        "Criação de conteúdo / scripts","Treinamento de modelos de IA"]
for i,f in enumerate(fora):
    iy = 3.2+i*0.48
    rect(s3,9.88,iy+0.08,0.08,0.08,fill=SF_RED)
    txt(s3,f,10.02,iy,3.0,0.42,sz=8.5,color=DARK)

# ══════════════════════════════════════════════════════
# SLIDE 4 — PUBLIC SECTOR SOLUTIONS (PSS) — MMAR
# ══════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
header(s4,"Public Sector Solutions (PSS) — MMAR",
       "Fase 3 · Licenciamento marítimo e aeronáutico · Licensing & Permitting + BRE · OmniStudio",
       SF_GREEN,"PSS")
footer(s4,"Slide 4 de 9")
rect(s4,0,1.23,13.33,5.8,fill=LGRAY)

rect(s4,0.2,1.35,8.5,1.42,fill=WHITE,line=SF_GREEN,lw=Pt(1))
rect(s4,0.2,1.35,8.5,0.3,fill=SF_GREEN)
txt(s4,"CONTEXTO DE NEGÓCIO",0.3,1.37,8.3,0.24,sz=9,bold=True,color=WHITE)
txt(s4,"O sistema Mosaico gerencia o licenciamento de estações de rádio para aeronaves e embarcações — processo hoje "
       "manual, fragmentado e dependente de análise jurídica humana em todas as etapas. São 6 tipos de estação "
       "(Embarcação, Radiobaliza, Costeira, Portuária, Móvel, Aeronave), cada um com dados técnicos específicos "
       "(MMSI, frequências, DSC, especificações de antena). Tempo de emissão atual: semanas. Impacto: segurança "
       "marítima e aérea comprometida por atrasos burocráticos.",
    0.3,1.72,8.3,0.95,sz=9.5,color=DARK,wrap=True)

for i,(v,l,c) in enumerate([("6","Tipos de estação no escopo",SF_GREEN),
                              ("Semanas","Tempo atual de emissão",SF_ORANGE),
                              ("Marinha\n+ DECEA","Integrações externas críticas",SF_BLUE),
                              ("XL","Sizing complexidade",SF_RED)]):
    bx = 9.0+i*1.0 if i<2 else 9.0+(i-2)*1.0
    by = 1.35 if i<2 else 2.1
    rect(s4,bx,by,0.9,0.62,fill=c)
    txt(s4,v,bx,by+0.03,0.9,0.32,sz=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(s4,l,bx+0.04,by+0.35,0.84,0.25,sz=7,color=WHITE,wrap=True,align=PP_ALIGN.CENTER)

rect(s4,0.2,2.86,5.9,4.02,fill=WHITE,line=SF_GREEN,lw=Pt(1))
rect(s4,0.2,2.86,5.9,0.3,fill=SF_GREEN)
txt(s4,"COMO ATENDEREMOS — SOLUÇÃO SALESFORCE",0.3,2.88,5.7,0.24,sz=9,bold=True,color=WHITE)
sol4 = [
    ("PSS Licensing & Permitting","Fluxo digital completo de solicitação de outorga · Automação "
     "multinível com checkpoints jurídicos e técnicos por tipo de estação"),
    ("Business Rules Engine (BRE)","Tradução das regras de análise legal do Mosaico para motor nativo Salesforce · "
     "Dispensa análise humana onde a lei permitir · Trilha de auditoria jurídica completa"),
    ("OmniStudio","Fluxo guiado de solicitação adaptativo por tipo de estação · "
     "Captura de dados técnicos específicos (MMSI, frequências, DSC, antenas)"),
    ("CLM — Gestão de Contratos","5 templates dinâmicos de outorga · Aprovação multinível (gerente + jurídico + diretoria) · "
     "Histórico de contratos 7 anos via Shield · Vinculado ao Mosaico via MuleSoft"),
    ("MuleSoft — Integração Mosaico","API para Mosaico legado · Integrações bidirecionais com Marinha do Brasil e DECEA · "
     "Sincronização de dados técnicos de estações"),
]
for i,(cloud,desc) in enumerate(sol4):
    iy = 3.24+i*0.72
    rect(s4,0.28,iy,5.74,0.64,fill=RGBColor(0xED,0xF8,0xF0),line=SF_GREEN,lw=Pt(0.5))
    rect(s4,0.28,iy,0.08,0.64,fill=SF_GREEN)
    txt(s4,cloud,0.44,iy+0.04,5.4,0.22,sz=9,bold=True,color=SF_GREEN)
    txt(s4,desc,0.44,iy+0.26,5.38,0.35,sz=8,color=DGRAY,wrap=True)

rect(s4,6.3,2.86,3.3,4.02,fill=WHITE,line=SF_GREEN,lw=Pt(1))
rect(s4,6.3,2.86,3.3,0.3,fill=SF_GREEN)
txt(s4,"PRINCIPAIS ENTREGAS",6.4,2.88,3.1,0.24,sz=9,bold=True,color=WHITE)
for i,e in enumerate(["Fluxo digital de solicitação MMAR","BRE com regras legais do Mosaico",
                       "Integração Marinha + DECEA","Portal do requerente (Experience)",
                       "CLM — 5 templates de outorga","Trilha de auditoria jurídica",
                       "Agentforce v3 para pendências","Histórico 7 anos (Shield)"]):
    iy = 3.24+i*0.41
    rect(s4,6.38,iy+0.07,0.08,0.08,fill=SF_GREEN)
    txt(s4,e,6.52,iy,3.0,0.38,sz=8.5,color=DARK)

rect(s4,9.8,2.86,3.33,4.02,fill=WHITE,line=SF_RED,lw=Pt(1))
rect(s4,9.8,2.86,3.33,0.3,fill=SF_RED)
txt(s4,"FORA DO ESCOPO",9.9,2.88,3.1,0.24,sz=9,bold=True,color=WHITE)
for i,f in enumerate(["Integração ANAC / Receita Federal","Assinatura digital ICP-Brasil",
                       "Track changes com contrapartes","Módulo aeronáutico (regras divergentes)",
                       "Mais de 5 templates CLM","Mais de 2 fluxos de aprovação CLM"]):
    iy = 3.24+i*0.48
    rect(s4,9.88,iy+0.08,0.08,0.08,fill=SF_RED)
    txt(s4,f,10.02,iy,3.0,0.42,sz=8.5,color=DARK)

# ══════════════════════════════════════════════════════
# SLIDE 5 — MULESOFT + DATA CLOUD (TFF/TFI)
# ══════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
header(s5,"MuleSoft + Data Cloud — TFF/TFI Arrecadação",
       "Fase 2 · 10M registros · 4 fontes legadas · Batch 2x/mês · Prazo regulatório 31/Mar/2027",
       SF_TEAL,"DC+MS")
footer(s5,"Slide 5 de 9")
rect(s5,0,1.23,13.33,5.8,fill=LGRAY)

rect(s5,0.2,1.35,8.5,1.42,fill=WHITE,line=SF_TEAL,lw=Pt(1))
rect(s5,0.2,1.35,8.5,0.3,fill=SF_TEAL)
txt(s5,"CONTEXTO DE NEGÓCIO",0.3,1.37,8.3,0.24,sz=9,bold=True,color=WHITE)
txt(s5,"A ANATEL arrecada anualmente a Taxa de Fiscalização do Funcionamento (TFF) de ~10M de entidades reguladas. "
       "Hoje, os dados estão fragmentados em 4 sistemas legados sem chave de identidade unificada, gerando inadimplência "
       "por falta de notificação efetiva. O prazo regulatório de 31/março para geração de boletos é não negociável — "
       "qualquer atraso no go-live implica em risco jurídico para a ANATEL. Retenção obrigatória de 7 anos para auditoria TCU/CGU.",
    0.3,1.72,8.3,0.95,sz=9.5,color=DARK,wrap=True)

for i,(v,l,c) in enumerate([("10M","Registros por batch",SF_TEAL),
                              ("4","Fontes legadas heterogêneas",SF_ORANGE),
                              ("31/Mar","Deadline regulatório 2027",SF_RED),
                              ("7 anos","Retenção auditável TCU",SF_BLUE)]):
    bx = 9.0+i*1.0 if i<2 else 9.0+(i-2)*1.0
    by = 1.35 if i<2 else 2.1
    rect(s5,bx,by,0.9,0.62,fill=c)
    txt(s5,v,bx,by+0.03,0.9,0.32,sz=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(s5,l,bx+0.04,by+0.35,0.84,0.25,sz=7,color=WHITE,wrap=True,align=PP_ALIGN.CENTER)

rect(s5,0.2,2.86,5.9,4.02,fill=WHITE,line=SF_TEAL,lw=Pt(1))
rect(s5,0.2,2.86,5.9,0.3,fill=SF_TEAL)
txt(s5,"COMO ATENDEREMOS — SOLUÇÃO SALESFORCE",0.3,2.88,5.7,0.24,sz=9,bold=True,color=WHITE)
sol5 = [
    ("MuleSoft CloudHub 2.0","Integração das 4 fontes: SITARWEB (REST), DB_TELECOM (JDBC), SMS/FISTEL (MongoDB), "
     "MOSAICO (JDBC legado) · Batch noturno 2x/mês · ~10M registros por execução"),
    ("Data Cloud — Identity Resolution","Resolução de identidade entre os 3 sistemas via chave composta · "
     "Unified Profile das entidades reguladas · ~10M registros consolidados · Qualidade de dados com regras automáticas"),
    ("Data Cloud — Big Objects","Retenção 7 anos auditável sem degradar performance operacional · "
     "Data Spaces para separação lógica dos períodos · Atende TCU/CGU e Lei de Arquivos"),
    ("Marketing Cloud — Jornada Anti-Inadimplência","Régua D-30/D-15/D-7/D+1/D+15 via WhatsApp+Email · "
     "3 jornadas: Renovação de Outorga, Adimplência TFF, Onboarding novo licenciado"),
    ("Agentforce v2","Reemissão autônoma de boleto TFF via WhatsApp · Consulta de débitos em tempo real · "
     "Handoff para servidor público se necessário"),
]
for i,(cloud,desc) in enumerate(sol5):
    iy = 3.24+i*0.72
    rect(s5,0.28,iy,5.74,0.64,fill=RGBColor(0xE4,0xF7,0xFA),line=SF_TEAL,lw=Pt(0.5))
    rect(s5,0.28,iy,0.08,0.64,fill=SF_TEAL)
    txt(s5,cloud,0.44,iy+0.04,5.4,0.22,sz=9,bold=True,color=SF_TEAL)
    txt(s5,desc,0.44,iy+0.26,5.38,0.35,sz=8,color=DGRAY,wrap=True)

rect(s5,6.3,2.86,3.3,4.02,fill=WHITE,line=SF_TEAL,lw=Pt(1))
rect(s5,6.3,2.86,3.3,0.3,fill=SF_TEAL)
txt(s5,"PRINCIPAIS ENTREGAS",6.4,2.88,3.1,0.24,sz=9,bold=True,color=WHITE)
for i,e in enumerate(["MuleSoft: 4 conectores legados","Identity Resolution 10M registros",
                       "Big Objects retenção 7 anos","3 jornadas MC (WhatsApp+Email)",
                       "Portal do contribuinte TFF","Agentforce v2 reemissão boleto",
                       "Dashboard inadimplência real-time","Go-live início Mar/2027"]):
    iy = 3.24+i*0.41
    rect(s5,6.38,iy+0.07,0.08,0.08,fill=SF_TEAL)
    txt(s5,e,6.52,iy,3.0,0.38,sz=8.5,color=DARK)

rect(s5,9.8,2.86,3.33,4.02,fill=WHITE,line=SF_RED,lw=Pt(1))
rect(s5,9.8,2.86,3.33,0.3,fill=SF_RED)
txt(s5,"FORA DO ESCOPO",9.9,2.88,3.1,0.24,sz=9,bold=True,color=WHITE)
for i,f in enumerate(["Motor de cálculo fiscal (externo)","SIAFI","Contestação/recurso de TFF",
                       "Histórico de cálculos anteriores","Mais de 4 fontes de origem",
                       "Mais de 3 jornadas MC","SMS / Push / outros canais"]):
    iy = 3.24+i*0.48
    rect(s5,9.88,iy+0.08,0.08,0.08,fill=SF_RED)
    txt(s5,f,10.02,iy,3.0,0.42,sz=8.5,color=DARK)

# ══════════════════════════════════════════════════════
# SLIDE 6 — REVENUE CLOUD + CLM
# ══════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
header(s6,"Revenue Cloud + CLM — Arrecadação e Contratos",
       "Fases 2–3 · TFI/TFF · Contratos de Outorga · 5 templates · 7 anos de retenção",
       SF_ORANGE,"RevCloud")
footer(s6,"Slide 6 de 9")
rect(s6,0,1.23,13.33,5.8,fill=LGRAY)

rect(s6,0.2,1.35,8.5,1.42,fill=WHITE,line=SF_ORANGE,lw=Pt(1))
rect(s6,0.2,1.35,8.5,0.3,fill=SF_ORANGE)
txt(s6,"CONTEXTO DE NEGÓCIO",0.3,1.37,8.3,0.24,sz=9,bold=True,color=WHITE)
txt(s6,"A ANATEL gera guias de recolhimento (TFI — Taxa de Fiscalização de Instalação) durante o processo de licenciamento "
       "MMAR e gerencia contratos de outorga com ciclo de vida complexo (solicitação → análise jurídica → aprovação multinível "
       "→ vigência → renovação). Hoje, esses processos são manuais, sem rastreabilidade de versões contratuais e sem "
       "alertas automatizados de vencimento — gerando inadimplência e risco jurídico para o órgão.",
    0.3,1.72,8.3,0.95,sz=9.5,color=DARK,wrap=True)

for i,(v,l,c) in enumerate([("5","Templates de outorga no escopo",SF_ORANGE),
                              ("3","Níveis de aprovação (ger+jur+dir)",SF_BLUE),
                              ("510–760h","Estimativa CLM (ROM)",SF_GREEN),
                              ("7 anos","Retenção contratos (Shield)",SF_DARK)]):
    bx = 9.0+i*1.0 if i<2 else 9.0+(i-2)*1.0
    by = 1.35 if i<2 else 2.1
    rect(s6,bx,by,0.9,0.62,fill=c)
    txt(s6,v,bx,by+0.03,0.9,0.32,sz=11,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(s6,l,bx+0.04,by+0.35,0.84,0.25,sz=7,color=WHITE,wrap=True,align=PP_ALIGN.CENTER)

rect(s6,0.2,2.86,5.9,4.02,fill=WHITE,line=SF_ORANGE,lw=Pt(1))
rect(s6,0.2,2.86,5.9,0.3,fill=SF_ORANGE)
txt(s6,"COMO ATENDEREMOS — SOLUÇÃO SALESFORCE",0.3,2.88,5.7,0.24,sz=9,bold=True,color=WHITE)
sol6 = [
    ("Revenue Cloud — TFI/TFF","Orquestração de geração e conciliação de guias TFI durante o processo MMAR · "
     "Integração com boleto bancário (GRU ou SF Billing) · Conciliação automática"),
    ("CLM — Ciclo de Vida de Contratos","5 etapas: Solicitação → Análise Jurídica → Aprovação → Vigência → Renovação · "
     "5 templates dinâmicos com cláusulas variáveis por tipo de outorga"),
    ("Aprovação Multinível","Workflow gerente + jurídico + diretoria · Sem assinatura digital ICP (fora do escopo) · "
     "Aprovação eletrônica com trilha de auditoria completa · Versionamento nativo SF"),
    ("Alertas de Vencimento e Renovação","D-90/D-60/D-30 automáticos · Fluxo de renovação ou encerramento · "
     "Vinculado ao Mosaico via MuleSoft · Integração com MC para comunicação externa"),
    ("Shield — Retenção e Auditoria","Histórico de contratos arquivado por 7 anos · FLS + Audit Trail · "
     "Exportação para auditorias TCU/CGU · Event Monitoring para rastreabilidade"),
]
for i,(cloud,desc) in enumerate(sol6):
    iy = 3.24+i*0.72
    rect(s6,0.28,iy,5.74,0.64,fill=RGBColor(0xFF,0xF3,0xE0),line=SF_ORANGE,lw=Pt(0.5))
    rect(s6,0.28,iy,0.08,0.64,fill=SF_ORANGE)
    txt(s6,cloud,0.44,iy+0.04,5.4,0.22,sz=9,bold=True,color=SF_ORANGE)
    txt(s6,desc,0.44,iy+0.26,5.38,0.35,sz=8,color=DGRAY,wrap=True)

rect(s6,6.3,2.86,3.3,4.02,fill=WHITE,line=SF_ORANGE,lw=Pt(1))
rect(s6,6.3,2.86,3.3,0.3,fill=SF_ORANGE)
txt(s6,"PRINCIPAIS ENTREGAS",6.4,2.88,3.1,0.24,sz=9,bold=True,color=WHITE)
for i,e in enumerate(["Revenue Cloud — geração de TFI","CLM — 5 templates de outorga",
                       "Fluxo aprovação 3 níveis","Alertas D-90/D-60/D-30 automáticos",
                       "Versionamento nativo de contratos","Integração Mosaico via MuleSoft",
                       "Histórico 7 anos (Shield)","Relatórios OOTB de status e vencimento"]):
    iy = 3.24+i*0.41
    rect(s6,6.38,iy+0.07,0.08,0.08,fill=SF_ORANGE)
    txt(s6,e,6.52,iy,3.0,0.38,sz=8.5,color=DARK)

rect(s6,9.8,2.86,3.33,4.02,fill=WHITE,line=SF_RED,lw=Pt(1))
rect(s6,9.8,2.86,3.33,0.3,fill=SF_RED)
txt(s6,"FORA DO ESCOPO",9.9,2.88,3.1,0.24,sz=9,bold=True,color=WHITE)
for i,f in enumerate(["Assinatura digital ICP-Brasil / DocuSign","Redline com contrapartes externas",
                       "Mais de 2 fluxos de aprovação","Mais de 5 templates CLM",
                       "KPIs customizados (apenas OOTB)","Integração SEI / SIAFI"]):
    iy = 3.24+i*0.48
    rect(s6,9.88,iy+0.08,0.08,0.08,fill=SF_RED)
    txt(s6,f,10.02,iy,3.0,0.42,sz=8.5,color=DARK)

# ══════════════════════════════════════════════════════
# SLIDE 7 — MARKETING CLOUD
# ══════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
header(s7,"Marketing Cloud — Comunicação Regulatória Inteligente",
       "Fase 2 · 3 jornadas · ~5M contatos · WhatsApp + E-mail · Adimplência + Renovação + Onboarding",
       SF_PURPLE,"MC")
footer(s7,"Slide 7 de 9")
rect(s7,0,1.23,13.33,5.8,fill=LGRAY)

rect(s7,0.2,1.35,8.5,1.42,fill=WHITE,line=SF_PURPLE,lw=Pt(1))
rect(s7,0.2,1.35,8.5,0.3,fill=SF_PURPLE)
txt(s7,"CONTEXTO DE NEGÓCIO",0.3,1.37,8.3,0.24,sz=9,bold=True,color=WHITE)
txt(s7,"A ANATEL se comunica com ~5M de entidades reguladas hoje por canais fragmentados, sem jornadas automatizadas e "
       "sem personalização por perfil de licenciado. O resultado é inadimplência de TFF por falta de notificação efetiva, "
       "baixa taxa de renovação de outorgas dentro do prazo e onboarding manual de novos licenciados. "
       "O MC atua como camada de comunicação proativa integrada ao CLM, MuleSoft e Data Cloud.",
    0.3,1.72,8.3,0.95,sz=9.5,color=DARK,wrap=True)

for i,(v,l,c) in enumerate([("5M","Contatos / entidades reguladas",SF_PURPLE),
                              ("3","Jornadas no escopo",SF_BLUE),
                              ("WhatsApp\n+ E-mail","Canais ativos",SF_GREEN),
                              ("Enterprise","Edição MC Engagement",SF_ORANGE)]):
    bx = 9.0+i*1.0 if i<2 else 9.0+(i-2)*1.0
    by = 1.35 if i<2 else 2.1
    rect(s7,bx,by,0.9,0.62,fill=c)
    txt(s7,v,bx,by+0.03,0.9,0.32,sz=11,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(s7,l,bx+0.04,by+0.35,0.84,0.25,sz=7,color=WHITE,wrap=True,align=PP_ALIGN.CENTER)

rect(s7,0.2,2.86,5.9,4.02,fill=WHITE,line=SF_PURPLE,lw=Pt(1))
rect(s7,0.2,2.86,5.9,0.3,fill=SF_PURPLE)
txt(s7,"COMO ATENDEREMOS — 3 JORNADAS",0.3,2.88,5.7,0.24,sz=9,bold=True,color=WHITE)
sol7 = [
    ("Jornada 1 — Renovação de Outorga",
     "Trigger: D-90/D-60/D-30 do vencimento no CLM · Ação D-0: Journey Builder abre case no CLM via API/Flow · "
     "Canal: WhatsApp + E-mail · Dados: CLM → Data Extension"),
    ("Jornada 2 — Adimplência TFF/TFI",
     "Régua: D-30/D-15/D-7/D+1/D+15 · Dados de vencimento via MuleSoft TFF → Data Extension · "
     "WhatsApp com link de reemissão de boleto via Agentforce v2 · E-mail como fallback"),
    ("Jornada 3 — Onboarding Novo Licenciado",
     "Trigger: evento de aprovação de outorga no CLM · Sequência de boas-vindas + orientações regulatórias · "
     "Personalização por tipo de estação / segmento"),
    ("Infraestrutura de Dados MC","Data Cloud como camada de unificação · ~5M contatos em Data Extensions · "
     "Criativos (templates e-mail + HSMs WhatsApp) fornecidos pela ANATEL — PS faz somente configuração técnica"),
    ("Integração MC ↔ CLM ↔ MuleSoft","Journey Builder → CLM via API/Flow (Jornada 1, ação D-0) · "
     "MuleSoft TFF → Data Extension (Jornada 2) · Entrega incremental após CLM e Data Cloud configurados"),
]
for i,(cloud,desc) in enumerate(sol7):
    iy = 3.24+i*0.72
    rect(s7,0.28,iy,5.74,0.64,fill=RGBColor(0xF3,0xE8,0xFF),line=SF_PURPLE,lw=Pt(0.5))
    rect(s7,0.28,iy,0.08,0.64,fill=SF_PURPLE)
    txt(s7,cloud,0.44,iy+0.04,5.4,0.22,sz=9,bold=True,color=SF_PURPLE)
    txt(s7,desc,0.44,iy+0.26,5.38,0.35,sz=8,color=DGRAY,wrap=True)

rect(s7,6.3,2.86,3.3,4.02,fill=WHITE,line=SF_PURPLE,lw=Pt(1))
rect(s7,6.3,2.86,3.3,0.3,fill=SF_PURPLE)
txt(s7,"PRINCIPAIS ENTREGAS",6.4,2.88,3.1,0.24,sz=9,bold=True,color=WHITE)
for i,e in enumerate(["3 jornadas configuradas e testadas","Integration MC → CLM (ação D-0)",
                       "Data Extensions ~5M contatos","HSMs WhatsApp configurados",
                       "Templates e-mail configurados","Régua adimplência 5 pontos de contato",
                       "Documentação das jornadas","Teste de performance segmentação"]):
    iy = 3.24+i*0.41
    rect(s7,6.38,iy+0.07,0.08,0.08,fill=SF_PURPLE)
    txt(s7,e,6.52,iy,3.0,0.38,sz=8.5,color=DARK)

rect(s7,9.8,2.86,3.33,4.02,fill=WHITE,line=SF_RED,lw=Pt(1))
rect(s7,9.8,2.86,3.33,0.3,fill=SF_RED)
txt(s7,"FORA DO ESCOPO",9.9,2.88,3.1,0.24,sz=9,bold=True,color=WHITE)
for i,f in enumerate(["Criação de criativos / copywriting","Aprovação WABA junto à Meta",
                       "Mais de 3 jornadas","SMS / Push / outros canais",
                       "Einstein AI além das 3 jornadas","Segmentação avançada ML customizado"]):
    iy = 3.24+i*0.48
    rect(s7,9.88,iy+0.08,0.08,0.08,fill=SF_RED)
    txt(s7,f,10.02,iy,3.0,0.42,sz=8.5,color=DARK)

# ══════════════════════════════════════════════════════
# SLIDE 8 — PREMISSAS
# ══════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
header(s8,"Premissas Críticas do Programa",
       "Itens que devem ser confirmados antes do kick-off de cada fase — impacto direto em prazo e estimativa",
       SF_DARK)
footer(s8,"Slide 8 de 9")
rect(s8,0,1.23,13.33,5.8,fill=LGRAY)

premissas = [
    (SF_RED,"BLOQUEANTE","MuleSoft licenciado antes do kick-off Fase 0",
     "Sem MuleSoft, a integração das 4 fontes TFF não pode começar → atraso em cascata até prazo regulatório"),
    (SF_RED,"BLOQUEANTE","Data Cloud licenciado antes do início Fase 2",
     "Identity Resolution e Big Objects (retenção 7 anos) são pré-requisitos do go-live TFF"),
    (SF_RED,"BLOQUEANTE","Acesso às 3 bases TFF (SITARWEB, DB_TELECOM, MongoDB) garantido antes Fase 2",
     "Sem acesso às fontes, o data profiling e a configuração dos conectores MuleSoft não podem ocorrer"),
    (SF_ORANGE,"CRÍTICA","OmniStudio licenciado com aquisição iniciada na Fase 1",
     "Lead time de licenciamento Industries/Vlocity — se não iniciado na F1, bloqueia a F3 (MMAR)"),
    (SF_ORANGE,"CRÍTICA","Regras de cálculo TFF documentadas antes do início Fase 2",
     "Motor fiscal permanece externo — PS não pode estimar integração sem conhecer o schema de saída"),
    (SF_ORANGE,"CRÍTICA","API do Mosaico documentada antes do início Fase 3",
     "Mosaico via JDBC direto (sem API) — confirmação de acesso e permissões de leitura é pré-requisito hard"),
    (SF_ORANGE,"CRÍTICA","WhatsApp Business API (Meta) aprovado — processo iniciado na Fase 0",
     "Aprovação Meta leva 4–8 semanas — e-mail como fallback, mas WhatsApp é o canal primário"),
    (SF_YELLOW,"IMPORTANTE","PO dedicado ANATEL com poder de decisão em todas as fases",
     "Sem PO com autoridade, decisões de escopo ficam represadas e comprometem o prazo regulatório TFF"),
    (SF_YELLOW,"IMPORTANTE","Sponsor executivo com autoridade sobre as 6 gerências",
     "Change management de 6 gerências simultâneas exige patrocinador com poder de mobilizar todas elas"),
    (SF_YELLOW,"IMPORTANTE","IdP corporativo ANATEL definido (LDAP, Azure AD ou Gov.br) antes Fase 0",
     "SSO crítico para acesso unificado — sem IdP definido, o modelo de autenticação do projeto fica em aberto"),
]

cols = [[0,1,2,3,4],[5,6,7,8,9]]
col_x = [0.2, 6.78]
col_w = 6.35

for ci,col_items in enumerate(cols):
    cx = col_x[ci]
    for ri,idx in enumerate(col_items):
        cor,nivel,titulo,detalhe = premissas[idx]
        ry = 1.38 + ri*1.04
        rect(s8,cx,ry,col_w,0.94,fill=WHITE,line=cor,lw=Pt(1))
        rect(s8,cx,ry,0.1,0.94,fill=cor)
        badge_tc = WHITE if nivel!="IMPORTANTE" else DARK
        rect(s8,cx+col_w-1.12,ry+0.06,1.02,0.24,fill=cor)
        txt(s8,nivel,cx+col_w-1.12,ry+0.07,1.02,0.2,
            sz=7,bold=True,color=badge_tc,align=PP_ALIGN.CENTER)
        txt(s8,titulo,cx+0.18,ry+0.05,col_w-1.25,0.28,sz=9,bold=True,color=DARK)
        txt(s8,detalhe,cx+0.18,ry+0.38,col_w-0.28,0.48,sz=8,color=DGRAY,italic=True,wrap=True)

# ══════════════════════════════════════════════════════
# SLIDE 9 — RISCOS
# ══════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
header(s9,"Riscos do Programa — Matriz de Exposição",
       "7 riscos mapeados · Probabilidade × Impacto · Mitigações definidas por risco",SF_RED)
footer(s9,"Slide 9 de 9")
rect(s9,0,1.23,13.33,5.8,fill=LGRAY)

riscos = [
    (SF_RED,"CRÍTICO","Alta","Prazo TFF 31/março não respeitado",
     "Qualquer atraso na Fase 2 implica em risco jurídico para a ANATEL — prazo não negociável",
     "Go-live Fase 2 planejado para início de março/2027 — 4 semanas de margem. "
     "Milestone tracker semanal a partir da Fase 1."),
    (SF_RED,"CRÍTICO","Alta","Regras TFF com alta variabilidade não mapeada",
     "Motor fiscal externo recebe dados consolidados — se as regras tiverem exceções não documentadas, "
     "o schema de saída muda e impacta o MuleSoft",
     "Discovery dedicado de regras TFF antes da Fase 2. "
     "Motor fiscal permanece externo e fora do escopo PS."),
    (SF_ORANGE,"ALTO","Alta","Qualidade dos dados nas 3 fontes TFF",
     "Dados fragmentados sem chave de identidade unificada — registros duplicados ou sem chave válida "
     "comprometem a Identity Resolution do Data Cloud",
     "Data profiling obrigatório na Fase 0. "
     "Regras de qualidade no Data Cloud antes da carga dos 10M registros."),
    (SF_ORANGE,"ALTO","Média","API do Mosaico indisponível ou sem documentação",
     "Mosaico legado via JDBC direto — se o acesso ao banco não for liberado, "
     "a integração MMAR (Fase 3) e TFF (Fase 2) ficam bloqueadas",
     "Discovery técnico do Mosaico durante a Fase 2. "
     "Confirmar permissões de leitura no kick-off da Fase 0."),
    (SF_ORANGE,"ALTO","Média","Lead time licenciamento OmniStudio não iniciado a tempo",
     "Industries/Vlocity tem lead time de aquisição — se não iniciado na Fase 1, "
     "bloqueia o início da Fase 3 (MMAR)",
     "Aquisição obrigatória durante a Fase 1. "
     "Validar status de procurement a cada status report."),
    (SF_ORANGE,"ALTO","Média","Aprovação Meta WhatsApp Business atrasada",
     "Processo de aprovação WABA leva 4–8 semanas — se não iniciado na Fase 0, "
     "o canal principal da Fase 1 (Agentforce v1) não está disponível no go-live",
     "Iniciar processo na Fase 0. "
     "E-mail configurado como fallback obrigatório para todos os fluxos de atendimento."),
    (SF_YELLOW,"ALTO","Alta","Resistência ao change management (6 gerências)",
     "Sponsor executivo precisa mobilizar GIDS, GIMR, GIIB, ORLE, ORER e AFO simultaneamente — "
     "resistência de qualquer gerência compromete adoção e SLA",
     "Sponsor executivo com autoridade formal sobre as 6 gerências como premissa hard. "
     "UX/HCC Consultant nas Fases 0, 1 e 3."),
]

rh = 0.88
for i,(_,nivel,prob,titulo,contexto,mitigacao) in enumerate(riscos):
    col = i % 2
    row = i // 2
    cx = 0.2 + col*6.57
    ry = 1.38 + row*rh*1.04
    cor = SF_RED if nivel=="CRÍTICO" else SF_ORANGE if nivel=="ALTO" else SF_YELLOW
    badge_tc = WHITE if nivel!="IMPORTANTE" else DARK

    rect(s9,cx,ry,6.37,rh,fill=WHITE,line=cor,lw=Pt(1))
    rect(s9,cx,ry,0.1,rh,fill=cor)

    # badges
    rect(s9,cx+0.16,ry+0.06,0.88,0.22,fill=cor)
    txt(s9,nivel,cx+0.16,ry+0.07,0.88,0.18,sz=7,bold=True,
        color=badge_tc,align=PP_ALIGN.CENTER)

    prob_c = SF_RED if prob=="Alta" else SF_ORANGE if prob=="Média" else SF_GREEN
    rect(s9,cx+1.1,ry+0.06,1.1,0.22,fill=prob_c)
    txt(s9,f"Prob: {prob}",cx+1.1,ry+0.07,1.1,0.18,sz=7,bold=True,
        color=WHITE,align=PP_ALIGN.CENTER)

    txt(s9,titulo,cx+0.16,ry+0.34,6.05,0.22,sz=9,bold=True,color=DARK)
    txt(s9,f"Contexto: {contexto}",cx+0.16,ry+0.56,4.0,0.18,sz=7.5,color=DGRAY,wrap=True,italic=True)
    rect(s9,cx+4.3,ry+0.28,0.03,rh-0.36,fill=MGRAY)
    txt(s9,"Mitigação",cx+4.4,ry+0.28,1.8,0.2,sz=7.5,bold=True,color=cor)
    txt(s9,mitigacao,cx+4.4,ry+0.46,1.85,0.38,sz=7.5,color=DARK,wrap=True)

# Último item (7°) centralizado
i=6
_,nivel,prob,titulo,contexto,mitigacao = riscos[6]
cor = SF_YELLOW; badge_tc = DARK
cx = 0.2+3.285; ry = 1.38+3*rh*1.04
rect(s9,cx,ry,6.37,rh,fill=WHITE,line=cor,lw=Pt(1))
rect(s9,cx,ry,0.1,rh,fill=cor)
rect(s9,cx+0.16,ry+0.06,0.88,0.22,fill=cor)
txt(s9,nivel,cx+0.16,ry+0.07,0.88,0.18,sz=7,bold=True,color=badge_tc,align=PP_ALIGN.CENTER)
prob_c = SF_ORANGE
rect(s9,cx+1.1,ry+0.06,1.1,0.22,fill=prob_c)
txt(s9,f"Prob: {prob}",cx+1.1,ry+0.07,1.1,0.18,sz=7,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
txt(s9,titulo,cx+0.16,ry+0.34,6.05,0.22,sz=9,bold=True,color=DARK)
txt(s9,f"Contexto: {contexto}",cx+0.16,ry+0.56,4.0,0.18,sz=7.5,color=DGRAY,wrap=True,italic=True)
rect(s9,cx+4.3,ry+0.28,0.03,rh-0.36,fill=MGRAY)
txt(s9,"Mitigação",cx+4.4,ry+0.28,1.8,0.2,sz=7.5,bold=True,color=cor)
txt(s9,mitigacao,cx+4.4,ry+0.46,1.85,0.38,sz=7.5,color=DARK,wrap=True)

# ── Salvar ──────────────────────────────────────────
path = "/Users/nfilho/claude/ANATEL_OTT_Defesa_Tecnica.pptx"
prs.save(path)
print(f"Salvo: {path}")
