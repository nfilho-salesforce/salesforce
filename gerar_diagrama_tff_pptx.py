from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
import copy

# ── Colors ────────────────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
SF_BLUE    = RGBColor(0x00, 0x70, 0xD2)   # Salesforce primary
SF_BLUE2   = RGBColor(0x00, 0x52, 0x9E)   # darker blue
SF_NAVY    = RGBColor(0x03, 0x2D, 0x60)   # Salesforce navy
SF_ORANGE  = RGBColor(0xFF, 0x7A, 0x00)   # MuleSoft orange
GRAY_BG    = RGBColor(0xF4, 0xF6, 0xF9)   # light section bg
GRAY_BOR   = RGBColor(0xD8, 0xDF, 0xE6)   # border
GRAY_TEXT  = RGBColor(0x54, 0x65, 0x7D)   # secondary text
GREEN      = RGBColor(0x2E, 0x84, 0x4E)   # MongoDB green
BLUE_SQL   = RGBColor(0x0F, 0x62, 0xFE)   # SQL Server
ORANGE_LEG = RGBColor(0xFF, 0x5A, 0x00)   # legacy / SITARWEB
AMBER      = RGBColor(0xCC, 0x7A, 0x00)   # warning
AMBER_BG   = RGBColor(0xFF, 0xF3, 0xCD)
PURPLE     = RGBColor(0x5A, 0x2D, 0x82)   # outputs accent
TEAL       = RGBColor(0x02, 0xA3, 0xBF)   # flow/automation

# ── Slide setup ───────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Cm(25.4)
prs.slide_height = Cm(14.29)
layout = prs.slide_layouts[6]  # blank
slide  = prs.slides.add_slide(layout)

def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_width_pt=1.0, radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Cm(x), Cm(y), Cm(w), Cm(h)
    )
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = Pt(line_width_pt)
    else:
        line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=9, bold=False, color=None,
             align=PP_ALIGN.CENTER, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return tb

def add_arrow_h(slide, x1, y, x2, color_rgb, dash=False):
    """Horizontal arrow left→right using a connector shape."""
    from pptx.util import Cm, Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    cx = slide.shapes.add_connector(1, Cm(x1), Cm(y), Cm(x2), Cm(y))
    cx.line.color.rgb = color_rgb
    cx.line.width = Pt(1.8)
    # arrowhead
    spPr = cx._element
    ln = spPr.find('.//' + qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(spPr, qn('a:ln'))
    tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', 'none')
    headEnd = etree.SubElement(ln, qn('a:headEnd'))
    headEnd.set('type', 'arrow')
    headEnd.set('w', 'med')
    headEnd.set('len', 'med')
    if dash:
        prstDash = etree.SubElement(ln, qn('a:prstDash'))
        prstDash.set('val', 'dash')
    return cx

# ── Background ────────────────────────────────────────────────────────────────
add_rect(slide, 0, 0, 25.4, 14.29, fill_rgb=WHITE)

# ── Header bar ────────────────────────────────────────────────────────────────
add_rect(slide, 0, 0, 25.4, 1.35, fill_rgb=SF_NAVY)
# accent stripe
add_rect(slide, 0, 1.35, 25.4, 0.12, fill_rgb=SF_ORANGE)

add_text(slide, 'ANATEL — TFF/TFI  |  Arquitetura de Dados',
         0.4, 0.15, 16, 0.7, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(slide, 'Fase 2  ·  MuleSoft como Orquestrador  ·  Fluxo de Dados',
         0.4, 0.82, 16, 0.45, font_size=8.5, color=RGBColor(0xB0,0xC4,0xDE),
         align=PP_ALIGN.LEFT, italic=True)
add_text(slide, 'Salesforce Professional Services  |  LATAM',
         17.5, 0.35, 7.5, 0.55, font_size=8, color=RGBColor(0xB0,0xC4,0xDE),
         align=PP_ALIGN.RIGHT)

# ── Column headers ────────────────────────────────────────────────────────────
cols = [
    (0.25,  'FONTES DE DADOS',     GRAY_TEXT),
    (6.05,  'ORQUESTRADOR',        SF_ORANGE),
    (11.2,  'SALESFORCE PLATFORM', SF_BLUE),
    (17.6,  'SAÍDAS / USUÁRIO',    PURPLE),
]
for (cx, ct, cc) in cols:
    cw = 5.5 if cx < 6 else (5.9 if cx < 11 else (6.2 if cx < 17 else 7.5))
    add_rect(slide, cx, 1.65, cw, 0.5, fill_rgb=GRAY_BG, line_rgb=GRAY_BOR)
    add_text(slide, ct, cx, 1.67, cw, 0.45,
             font_size=7.5, bold=True, color=cc)

# ── FONTES ────────────────────────────────────────────────────────────────────
sources = [
    (0.3,  2.4,  'SITARWEB',    'Sistema de fiscalização\nde estações de rádio', 'Legacy / API REST', ORANGE_LEG),
    (0.3,  4.55, 'DB_TELECOM',  'Obrigações e contribuintes\nTFF principais',     'SQL Server / JDBC', BLUE_SQL),
    (0.3,  6.7,  'SMS / FISTEL','Registros e histórico\nde arrecadação',          'MongoDB Connector', GREEN),
]
for (sx, sy, st, sd, stag, sc) in sources:
    add_rect(slide, sx, sy,    5.5, 1.85, fill_rgb=GRAY_BG, line_rgb=sc, line_width_pt=1.5)
    add_text(slide, st,        sx+0.15, sy+0.12, 5.2, 0.5, font_size=10, bold=True, color=sc, align=PP_ALIGN.LEFT)
    add_text(slide, sd,        sx+0.15, sy+0.6,  5.1, 0.75, font_size=8, color=GRAY_TEXT, align=PP_ALIGN.LEFT)
    add_rect(slide, sx+0.15, sy+1.38, 2.8, 0.35, fill_rgb=WHITE, line_rgb=sc)
    add_text(slide, stag,      sx+0.15, sy+1.39, 2.8, 0.33, font_size=7.5, bold=True, color=sc)

# Alerta
add_rect(slide, 0.3, 8.85, 5.5, 1.1, fill_rgb=AMBER_BG, line_rgb=AMBER, line_width_pt=1.5)
add_text(slide, '⚠  ~10M registros',
         0.45, 8.92, 5.2, 0.4, font_size=9, bold=True, color=AMBER, align=PP_ALIGN.LEFT)
add_text(slide, 'Identidade fragmentada entre as 3 fontes.\nDeduplicação e resolução tratadas no MuleSoft.',
         0.45, 9.32, 5.2, 0.6, font_size=7.5, color=RGBColor(0x7A,0x4F,0x00), align=PP_ALIGN.LEFT)

# ── MULESOFT ──────────────────────────────────────────────────────────────────
add_rect(slide, 6.1, 2.4, 4.85, 7.55, fill_rgb=RGBColor(0xFF,0xF7,0xF0), line_rgb=SF_ORANGE, line_width_pt=2.5)
# icon area
add_rect(slide, 6.1, 2.4, 4.85, 1.1, fill_rgb=SF_ORANGE)
add_text(slide, 'MuleSoft  Anypoint Platform',
         6.1, 2.52, 4.85, 0.7, font_size=11, bold=True, color=WHITE)

# badge
add_rect(slide, 6.9, 3.7, 3.2, 0.42, fill_rgb=SF_ORANGE, line_rgb=SF_ORANGE)
add_text(slide, 'ORQUESTRADOR CENTRAL',
         6.9, 3.72, 3.2, 0.38, font_size=7.5, bold=True, color=WHITE)

features = [
    'Normalização de esquemas',
    'Transformação de dados',
    'Deduplicação de registros',
    'Roteamento por tipo de dado',
    'Error handling & retry',
    'Auditoria de transações',
    'Rate limiting / throttling',
    'Monitoramento (Anypoint)',
]
for i, feat in enumerate(features):
    fy = 4.35 + i * 0.52
    add_text(slide, f'▸  {feat}', 6.35, fy, 4.5, 0.45,
             font_size=8.5, color=RGBColor(0x7A, 0x3A, 0x00), align=PP_ALIGN.LEFT)

# Protocol footer
add_rect(slide, 6.35, 9.35, 4.35, 0.5, fill_rgb=WHITE, line_rgb=SF_ORANGE)
add_text(slide, 'Salesforce REST API  ·  Platform Events  ·  Bulk API',
         6.35, 9.37, 4.35, 0.45, font_size=7.5, color=SF_ORANGE, bold=True)

# ── SF PLATFORM ───────────────────────────────────────────────────────────────
sf_cards = [
    (11.25, 2.4,  'Flow Builder',  'Jornada anti-inadimplência\nD-30/D-15/D-7/D+1/D+15',    'Automação', TEAL),
    (11.25, 4.75, 'Service Cloud', 'Fila unificada de atendimento\ne gestão de casos TFF',   'Case Mgmt', SF_BLUE),
    (11.25, 7.1,  'Big Objects',   'Retenção imutável de 7 anos\nauditável (TCU/compliance)', 'Compliance', PURPLE),
]
for (fx, fy, ft, fd, ftag, fc) in sf_cards:
    add_rect(slide, fx, fy, 6.1, 1.85, fill_rgb=GRAY_BG, line_rgb=fc, line_width_pt=1.5)
    add_text(slide, ft, fx+0.15, fy+0.12, 5.8, 0.5, font_size=10, bold=True, color=fc, align=PP_ALIGN.LEFT)
    add_text(slide, fd, fx+0.15, fy+0.6,  5.8, 0.75, font_size=8, color=GRAY_TEXT, align=PP_ALIGN.LEFT)
    add_rect(slide, fx+0.15, fy+1.38, 2.5, 0.35, fill_rgb=WHITE, line_rgb=fc)
    add_text(slide, ftag, fx+0.15, fy+1.39, 2.5, 0.33, font_size=7.5, bold=True, color=fc)

# ── SAÍDAS ────────────────────────────────────────────────────────────────────
out_cards = [
    (17.65, 2.4,  'WhatsApp + E-mail',  'Notificações automáticas\nda jornada com link de pagamento', 'Omnichannel', GREEN),
    (17.65, 4.75, 'Atendimento TFF',    'Servidor com contexto completo\ndo contribuinte no SF',       'Service Cloud', SF_BLUE),
    (17.65, 7.1,  'Trilha de Auditoria','7 anos de histórico imutável\npara conformidade e TCU',       'Big Objects', PURPLE),
]
for (ox, oy, ot, od, otag, oc) in out_cards:
    add_rect(slide, ox, oy, 7.45, 1.85, fill_rgb=GRAY_BG, line_rgb=oc, line_width_pt=1.5)
    add_text(slide, ot, ox+0.15, oy+0.12, 7.2, 0.5, font_size=10, bold=True, color=oc, align=PP_ALIGN.LEFT)
    add_text(slide, od, ox+0.15, oy+0.6,  7.2, 0.75, font_size=8, color=GRAY_TEXT, align=PP_ALIGN.LEFT)
    add_rect(slide, ox+0.15, oy+1.38, 2.8, 0.35, fill_rgb=WHITE, line_rgb=oc)
    add_text(slide, otag, ox+0.15, oy+1.39, 2.8, 0.33, font_size=7.5, bold=True, color=oc)

# ── ARROWS ────────────────────────────────────────────────────────────────────
# Fontes → MuleSoft
arrow_ys_src = [3.32, 5.47, 7.62]
for ay in arrow_ys_src:
    add_arrow_h(slide, 5.8, ay, 6.1, SF_ORANGE)

# MuleSoft → SF cards (fan to center of each)
sf_mid_ys = [3.32, 5.67, 8.02]
for ay in sf_mid_ys:
    add_arrow_h(slide, 10.95, ay, 11.25, SF_BLUE)

# SF → Outputs
out_mid_ys = [3.32, 5.67, 8.02]
for ay in out_mid_ys:
    add_arrow_h(slide, 17.35, ay, 17.65, PURPLE)

# ── FOOTER ────────────────────────────────────────────────────────────────────
add_rect(slide, 0, 12.95, 25.4, 1.34, fill_rgb=GRAY_BG, line_rgb=GRAY_BOR, line_width_pt=0.5)

# Warning
add_rect(slide, 0.25, 13.07, 15.5, 1.1, fill_rgb=AMBER_BG, line_rgb=AMBER, line_width_pt=1.5)
add_text(slide, '⚠  Prazo Regulatório Crítico',
         0.4, 13.1, 6, 0.45, font_size=8.5, bold=True, color=AMBER, align=PP_ALIGN.LEFT)
add_text(slide, 'Go-live planejado para início de março/2027 — margem de 4 semanas antes do prazo de 31/03/2027 para geração de boletos TFF.',
         0.4, 13.55, 15.1, 0.5, font_size=7.5, color=RGBColor(0x7A,0x4F,0x00), align=PP_ALIGN.LEFT)

# Legend
legend = [
    (SF_ORANGE, 'Ingestão — MuleSoft Connectors'),
    (SF_BLUE,   'Distribuição — Salesforce APIs'),
    (PURPLE,    'Entrega ao usuário'),
]
lx = 16.2
for (lc, lt) in legend:
    add_rect(slide, lx, 13.25, 0.7, 0.25, fill_rgb=lc)
    add_text(slide, lt, lx + 0.8, 13.2, 4.5, 0.35, font_size=7.5, color=GRAY_TEXT, align=PP_ALIGN.LEFT)
    lx += 4.5 if lx < 20 else 0

# ── Save ─────────────────────────────────────────────────────────────────────
out = '/Users/nfilho/claude/ANATEL_TFF_Arquitetura.pptx'
prs.save(out)
print(f'Salvo: {out}')
