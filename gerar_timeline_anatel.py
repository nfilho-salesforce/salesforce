from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
WHITE    = RGBColor(0xFF,0xFF,0xFF)
NAVY     = RGBColor(0x03,0x2D,0x60)
SF_BLUE  = RGBColor(0x00,0x70,0xD2)
ORANGE   = RGBColor(0xFF,0x7A,0x00)
GRAY_BG  = RGBColor(0xF4,0xF6,0xF9)
GRAY_BOR = RGBColor(0xD8,0xDF,0xE6)
GRAY_TXT = RGBColor(0x54,0x65,0x7D)
GREEN    = RGBColor(0x1B,0x7A,0x43)
GREEN_L  = RGBColor(0xD5,0xF2,0xE2)
PURPLE   = RGBColor(0x5A,0x2D,0x82)
TEAL     = RGBColor(0x02,0xA3,0xBF)
RED      = RGBColor(0xBA,0x23,0x25)
DARK     = RGBColor(0x16,0x1C,0x28)
MULE_G   = RGBColor(0x1E,0x8A,0x44)   # MuleSoft green bars

# Hours heatmap (blue scale for Agentes, green scale for MuleSoft)
def hour_color_blue(h):
    if   h ==  0: return WHITE
    elif h <= 10: return RGBColor(0xCC,0xE5,0xFF)
    elif h <= 20: return RGBColor(0x7A,0xBB,0xF5)
    elif h <= 30: return RGBColor(0x35,0x8A,0xE0)
    else:         return RGBColor(0x00,0x70,0xD2)

def hour_color_green(h):
    if   h ==  0: return WHITE
    elif h <= 10: return RGBColor(0xC6,0xEF,0xD4)
    elif h <= 20: return RGBColor(0x7D,0xD1,0x9D)
    elif h <= 30: return RGBColor(0x3E,0xA8,0x67)
    else:         return RGBColor(0x1E,0x8A,0x44)

def txt_on_h(h, green=False):
    return WHITE if h >= 30 else DARK

# ── Slide setup ───────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Cm(25.4)
prs.slide_height = Cm(14.29)
sl = prs.slides.add_slide(prs.slide_layouts[6])

def R(x,y,w,h, fc=None, lc=None, lw=0.4):
    s = sl.shapes.add_shape(1, Cm(x),Cm(y),Cm(w),Cm(h))
    if fc: s.fill.solid(); s.fill.fore_color.rgb = fc
    else:  s.fill.background()
    if lc: s.line.color.rgb = lc; s.line.width = Pt(lw)
    else:  s.line.fill.background()
    return s

def T(text, x,y,w,h, size=7, bold=False, color=DARK, align=PP_ALIGN.CENTER, italic=False):
    tb = sl.shapes.add_textbox(Cm(x),Cm(y),Cm(w),Cm(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color

# ── Layout ────────────────────────────────────────────────────────────────────
LM  = 0.18   # left margin
LC  = 3.28   # label column width
WC  = 1.357  # week column width (16 cols × 1.357 + 3.28 + 0.36 = 25.4)
RH  = 0.43   # role row height
MH  = 0.41   # milestone row height
SH  = 0.38   # section header height

def cx(week):   # X of week column (0-indexed)
    return LM + LC + week * WC

def bar(row_y, row_h, s, e, label, color, pad=0.055):
    x1  = cx(s)
    bw  = cx(e) + WC - x1 - 0.05
    by  = row_y + pad
    bh  = row_h - pad*2
    R(x1+0.03, by, bw, bh, fc=color)
    # Show label only if bar is wide enough
    if bw > 0.5:
        T(label, x1+0.04, by, bw-0.05, bh, size=5.5, bold=True, color=WHITE)

# ── Background ────────────────────────────────────────────────────────────────
R(0,0,25.4,14.29, fc=WHITE)

# ── Header ────────────────────────────────────────────────────────────────────
R(0,0,25.4,1.32, fc=NAVY)
R(0,1.32,25.4,0.10, fc=ORANGE)
T('ANATEL — Fase 1: Agentforce + MuleSoft  ·  Linha do Tempo com Alocação de Recursos',
  0.3,0.08,20,0.65, size=12,bold=True,color=WHITE,align=PP_ALIGN.LEFT)
T('Semanas S0–S15  ·  Horas/semana por perfil (heatmap)  ·  Sobreposição de marcos visível',
  0.3,0.78,20,0.38, size=7.5,color=RGBColor(0xB0,0xC4,0xDE),italic=True,align=PP_ALIGN.LEFT)

Y = 1.52

# ── Week / Phase header ───────────────────────────────────────────────────────
# Phase color blocks (top half) + S-number (bottom half)
phase_info = [
    ('Week 0',    GRAY_TXT),  # S0
    ('Discovery', SF_BLUE),   # S1
    ('Define',    TEAL),      # S2
    ('Design',    PURPLE),    # S3
    ('Dev',       GREEN),     # S4
    ('Dev',       GREEN),     # S5
    ('Dev',       GREEN),     # S6
    ('Dev',       GREEN),     # S7
    ('Dev',       GREEN),     # S8
    ('Dev',       GREEN),     # S9
    ('Dev',       GREEN),     # S10
    ('SIT/UAT',   ORANGE),    # S11
    ('SIT/UAT',   ORANGE),    # S12
    ('SIT/UAT',   ORANGE),    # S13
    ('Scale',     RED),       # S14
    ('Scale',     RED),       # S15
]
R(LM,Y,LC,0.52, fc=NAVY)
T('Perfil / Semana', LM+0.05,Y+0.05,LC-0.1,0.42, size=7.5,bold=True,color=WHITE,align=PP_ALIGN.LEFT)
for i,(ph,pc) in enumerate(phase_info):
    R(cx(i),  Y,       WC, 0.26, fc=pc)
    T(f'S{i}', cx(i),  Y,       WC, 0.26, size=6.5,bold=True,color=WHITE)
    R(cx(i),  Y+0.26,  WC, 0.26, fc=GRAY_BG, lc=GRAY_BOR)
    T(ph,     cx(i),   Y+0.26,  WC, 0.26, size=5.5,color=GRAY_TXT)
Y += 0.52

# ── Senior Project Manager ────────────────────────────────────────────────────
pm_h = [40]*16
R(LM,Y,LC,RH, fc=RGBColor(0xE4,0xED,0xFF), lc=GRAY_BOR)
T('Senior Project Manager', LM+0.06,Y,LC-0.08,RH, size=7,bold=True,color=NAVY,align=PP_ALIGN.LEFT)
for i,h in enumerate(pm_h):
    R(cx(i),Y,WC,RH, fc=hour_color_blue(h), lc=GRAY_BOR)
    if h: T(h, cx(i),Y,WC,RH, size=6,bold=True,color=WHITE)
Y += RH

# ── FASE 1 — AGENTFORCE ───────────────────────────────────────────────────────
R(LM,Y,25.4-LM*2,SH, fc=SF_BLUE)
T('FASE 1 — AGENTFORCE', LM+0.1,Y,10,SH, size=7.5,bold=True,color=WHITE,align=PP_ALIGN.LEFT)
Y += SH

agent_roles = [
    ('Solution Architect',   [10,40,40,40,40,20,20,20,20,10,10,10,10,10,10,10], NAVY),
    ('Technical Architect',  [10,40,40,40,40,40,40,20,20,20,20,20,20,20,20,20], NAVY),
    ('UX Designer',          [10,40,40,40,40,20,20,20,20,20,20, 0, 0, 0, 0, 0], TEAL),
    ('Business Analyst',     [ 0,40,40,40,40,20,20,20,20,20,20, 0, 0, 0, 0, 0], TEAL),
    ('Developer 1',          [ 0, 0, 0,40,40,40,40,40,40,40,40,40,40,40,40,40], SF_BLUE),
    ('Developer 2',          [ 0, 0, 0, 0, 0,40,40,40,40,40,40,40,40,40, 0, 0], SF_BLUE),
    ('Developer 3',          [ 0, 0, 0, 0, 0, 0,40,40,40,40,40,30,40,40,40, 0], SF_BLUE),
    ('QA Engineer',          [ 0, 0, 0, 0, 0,40,40,40,40,40,40,40,40,40, 0, 0], ORANGE),
]
for i,(name,hours,lc) in enumerate(agent_roles):
    bg = WHITE if i%2==0 else GRAY_BG
    R(LM,Y,LC,RH, fc=bg, lc=GRAY_BOR)
    T(name, LM+0.06,Y,LC-0.08,RH, size=7,bold=(i<4),color=lc,align=PP_ALIGN.LEFT)
    for j,h in enumerate(hours):
        R(cx(j),Y,WC,RH, fc=hour_color_blue(h), lc=GRAY_BOR)
        if h: T(h, cx(j),Y,WC,RH, size=6,color=txt_on_h(h))
    Y += RH

# ── Agentforce Milestone rows ──────────────────────────────────────────────────
ms_labels = ['Knowledge Base','Orquestrador','MMAR','Consumidor','Ouvidoria']
MS_Y = Y
# Background grid
for r in range(5):
    ry = MS_Y + r*MH
    R(LM,ry,LC,MH, fc=RGBColor(0xF8,0xF8,0xFC), lc=GRAY_BOR)
    T(ms_labels[r], LM+0.06,ry,LC-0.08,MH, size=6.5,italic=True,color=PURPLE,align=PP_ALIGN.LEFT)
    for j in range(16):
        R(cx(j),ry,WC,MH, fc=WHITE, lc=RGBColor(0xEE,0xEE,0xF5))

# Milestone bars  (row, start, end_inclusive, label, color)
KB_C = PURPLE
DEV_C = SF_BLUE
HLG_C = ORANGE
GL_C  = GREEN
SETUP_C = GRAY_TXT

agent_bars = [
    # Row 0 — Knowledge Base milestones
    (0,  4,  4, 'KB MMAR',       KB_C),
    (0,  6,  7, 'KB CONSUMIDOR', KB_C),
    (0,  8,  9, 'KB OUVIDORIA',  KB_C),
    # Row 1 — Orquestrador
    (1,  3,  3, 'SETUP',                 SETUP_C),
    (1,  4,  6, 'AGENTE ORQUESTRADOR',   DEV_C),
    (1,  7,  8, 'HLG/E2E',              HLG_C),
    (1, 10, 13, 'GO-LIVE Suporte',        GL_C),
    # Row 2 — MMAR
    (2,  5,  7, 'AGENTE MMAR',   DEV_C),
    (2,  8,  9, 'HLG/E2E',      HLG_C),
    (2, 11, 14, 'Go Live Suporte', GL_C),
    # Row 3 — Consumidor
    (3,  7,  9, 'AGENTE CONSUMIDOR', DEV_C),
    (3, 10, 11, 'HLG/E2E',          HLG_C),
    (3, 12, 13, 'Go Live Suporte',   GL_C),
    # Row 4 — Ouvidoria
    (4,  9, 11, 'AGENTE OUVIDORIA', DEV_C),
    (4, 12, 13, 'HLG/E2E',         HLG_C),
    (4, 14, 15, 'Go Live Suporte',  GL_C),
]
for (row,s,e,label,color) in agent_bars:
    bar(MS_Y + row*MH, MH, s, e, label, color)

Y = MS_Y + 5*MH

# ── MULESOFT ──────────────────────────────────────────────────────────────────
R(LM,Y,25.4-LM*2,SH, fc=MULE_G)
T('MULESOFT — INTEGRAÇÃO', LM+0.1,Y,10,SH, size=7.5,bold=True,color=WHITE,align=PP_ALIGN.LEFT)
Y += SH

mule_roles = [
    ('Mule Technical Architect', [10,40,40,40,20,20,20,20,20,20,20,20,20, 0, 0, 0]),
    ('Mule Developer 1',         [ 0, 0,40,40,40,40,40,40,40,20,20,20,20, 0, 0, 0]),
    ('Mule Developer 2',         [ 0, 0, 0, 0, 0,40,40,40,40, 0, 0, 0, 0, 0, 0, 0]),
]
for i,(name,hours) in enumerate(mule_roles):
    bg = WHITE if i%2==0 else GRAY_BG
    R(LM,Y,LC,RH, fc=bg, lc=GRAY_BOR)
    T(name, LM+0.06,Y,LC-0.08,RH, size=7,bold=True,color=MULE_G,align=PP_ALIGN.LEFT)
    for j,h in enumerate(hours):
        R(cx(j),Y,WC,RH, fc=hour_color_green(h), lc=GRAY_BOR)
        if h: T(h, cx(j),Y,WC,RH, size=6,color=txt_on_h(h,green=True))
    Y += RH

# ── MuleSoft Milestone rows ───────────────────────────────────────────────────
mule_ms_labels = ['GOV.BR','Setup / APIs','Consumidor']
MULE_MS_Y = Y
API_C = TEAL
for r in range(3):
    ry = MULE_MS_Y + r*MH
    R(LM,ry,LC,MH, fc=RGBColor(0xF0,0xFD,0xF5), lc=GRAY_BOR)
    T(mule_ms_labels[r], LM+0.06,ry,LC-0.08,MH, size=6.5,italic=True,color=MULE_G,align=PP_ALIGN.LEFT)
    for j in range(16):
        R(cx(j),ry,WC,MH, fc=WHITE, lc=RGBColor(0xEE,0xF5,0xEE))

mule_bars = [
    (0,  4,  5, 'API GOV.BR',    API_C),
    (1,  2,  3, 'SETUP',         SETUP_C),
    (1,  4,  5, 'API MMAR',      API_C),
    (1,  6,  8, 'API OUVIDORIA', API_C),
    (2,  7,  9, 'API CONSUMIDOR',API_C),
]
for (row,s,e,label,color) in mule_bars:
    bar(MULE_MS_Y + row*MH, MH, s, e, label, color)

Y = MULE_MS_Y + 3*MH

# ── Legend ────────────────────────────────────────────────────────────────────
LY = Y + 0.08
R(LM,LY,25.4-LM*2,0.88, fc=GRAY_BG, lc=GRAY_BOR)
T('Legenda:', LM+0.1,LY+0.08,1.4,0.38, size=7,bold=True,color=DARK,align=PP_ALIGN.LEFT)

legend = [
    (RGBColor(0xCC,0xE5,0xFF),'10h'),
    (RGBColor(0x7A,0xBB,0xF5),'20h'),
    (RGBColor(0x00,0x70,0xD2),'40h  (Agentes)'),
    (RGBColor(0x1E,0x8A,0x44),'40h  (MuleSoft)'),
    (SF_BLUE,  'Desenvolvimento'),
    (KB_C,     'Knowledge Base'),
    (HLG_C,    'HLG / E2E'),
    (GL_C,     'Go-Live / Suporte'),
    (API_C,    'APIs MuleSoft'),
]
lx = LM + 1.6
for (lc2,lt) in legend:
    R(lx,LY+0.14,0.32,0.32, fc=lc2, lc=GRAY_BOR, lw=0.3)
    T(lt, lx+0.38,LY+0.14,2.35,0.32, size=6.5,color=DARK,align=PP_ALIGN.LEFT)
    lx += 2.55

# ── Save ──────────────────────────────────────────────────────────────────────
out = '/Users/nfilho/claude/ANATEL_Timeline_Fase1.pptx'
prs.save(out)
print(f'Salvo: {out}')

# Quick sanity: print total Y
print(f'Content ends at Y ≈ {Y:.2f} cm  (slide height: 14.29 cm)')
