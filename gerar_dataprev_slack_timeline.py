from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Cores
SLACK_PURPLE   = RGBColor(0x4A, 0x15, 0x4B)
SLACK_GREEN    = RGBColor(0x2E, 0xB6, 0x7D)
SLACK_YELLOW   = RGBColor(0xEC, 0xB2, 0x2E)
SLACK_RED      = RGBColor(0xE0, 0x1E, 0x5A)
SLACK_BLUE     = RGBColor(0x36, 0xC5, 0xF0)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
DARK           = RGBColor(0x1A, 0x1D, 0x21)
LGRAY          = RGBColor(0xF8, 0xF8, 0xF8)
MGRAY          = RGBColor(0xE8, 0xE8, 0xE8)
DGRAY          = RGBColor(0x61, 0x61, 0x61)
SF_BLUE        = RGBColor(0x00, 0x70, 0xD2)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w if line_w else Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=11, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def header(slide, title, subtitle=""):
    add_rect(slide, 0, 0, 13.33, 1.15, fill=SLACK_PURPLE)
    add_rect(slide, 0, 1.15, 13.33, 0.04, fill=SLACK_GREEN)
    add_text(slide, "Dataprev · Governança Agentes · Bolsão 3",
             0.3, 0.08, 8, 0.3, size=9, color=RGBColor(0xCC,0xCC,0xCC))
    add_text(slide, title, 0.3, 0.32, 10, 0.55, size=22, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.3, 0.87, 10, 0.28, size=10,
                 color=RGBColor(0xAA,0xBB,0xCC), italic=True)
    # Slack logo placeholder (colored block)
    add_rect(slide, 12.5, 0.22, 0.55, 0.55, fill=SLACK_GREEN)
    add_text(slide, "Slack", 12.5, 0.29, 0.55, 0.3, size=9, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

def footer(slide, note=""):
    add_rect(slide, 0, 7.25, 13.33, 0.25, fill=DARK)
    add_text(slide, "Salesforce Professional Services — LATAM  |  Uso Interno",
             0.3, 7.27, 7, 0.2, size=8, color=RGBColor(0x88,0x88,0x88))
    if note:
        add_text(slide, note, 7, 7.27, 6, 0.2, size=8,
                 color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────
# SLIDE 1 — CAPA
# ─────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
add_rect(s1, 0, 0, 13.33, 7.5, fill=SLACK_PURPLE)
# Decoração geométrica
add_rect(s1, 9.5, -0.5, 5, 5, fill=RGBColor(0x3A,0x05,0x3B))
add_rect(s1, 10.5, 4, 4, 4, fill=RGBColor(0x2A,0x03,0x2B))

# Badges de cor Slack
colors = [SLACK_GREEN, SLACK_YELLOW, SLACK_RED, SLACK_BLUE]
labels = ["#", "#", "#", "#"]
for i, (c, lb) in enumerate(zip(colors, labels)):
    add_rect(s1, 0.4 + i*0.55, 5.6, 0.42, 0.42, fill=c)

add_text(s1, "DATAPREV", 0.4, 1.6, 12, 0.6, size=14, bold=False,
         color=RGBColor(0x8B,0x6B,0x8C))
add_text(s1, "Implantação do Slack", 0.4, 2.1, 12, 1.0, size=42, bold=True, color=WHITE)
add_text(s1, "Governança Agentes Serviço na Ponta · Bolsão 3", 0.4, 3.05, 10, 0.45,
         size=16, color=RGBColor(0xCC,0xDD,0xEE))
add_text(s1, "Timeline, Atividades e Premissas", 0.4, 3.52, 10, 0.35,
         size=13, color=RGBColor(0x8B,0xAA,0xCC), italic=True)

add_rect(s1, 0.4, 4.15, 4.5, 0.04, fill=SLACK_GREEN)

add_text(s1, "Salesforce Professional Services — LATAM", 0.4, 4.35, 8, 0.3,
         size=10, color=RGBColor(0x99,0xAA,0xBB))
add_text(s1, "Junho 2026", 0.4, 4.65, 4, 0.3, size=10, color=RGBColor(0x77,0x88,0x99))

# ─────────────────────────────────────────────
# SLIDE 2 — VISÃO GERAL DA TIMELINE (swimlane)
# ─────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
header(s2, "Timeline de Implantação — Visão Geral",
       "8–10 dias de trabalho técnico · 4 fases sequenciais · 2 workshops de adoção")
footer(s2, "Slide 2 de 6")

# Background
add_rect(s2, 0, 1.22, 13.33, 5.8, fill=LGRAY)

# Semanas
semanas = ["Sem 1", "Sem 2", "Sem 3"]
sem_x   = [1.6, 5.15, 8.7]
sem_w   = 3.5
for sx, sl in zip(sem_x, semanas):
    add_rect(s2, sx, 1.28, sem_w, 0.38, fill=MGRAY)
    add_text(s2, sl, sx, 1.3, sem_w, 0.34, size=11, bold=True,
             color=DGRAY, align=PP_ALIGN.CENTER)

# Swimlanes
lanes = [
    ("Fase 1 · Setup", SLACK_PURPLE),
    ("Fase 2 · Integração SF→Slack", SF_BLUE),
    ("Fase 3 · Automações & RAID", SLACK_GREEN),
    ("Fase 4 · Treinamento & Handoff", SLACK_YELLOW),
]
lane_y = [1.78, 2.78, 3.78, 4.78]
lane_h = 0.82

for (ln, lc), ly in zip(lanes, lane_y):
    add_rect(s2, 0.18, ly, 1.35, lane_h, fill=lc)
    add_text(s2, ln, 0.19, ly+0.05, 1.3, lane_h-0.1, size=8.5, bold=True,
             color=WHITE, wrap=True)

# Barras de atividade
# Formato: (lane_idx, sem_start_x_offset, width_in_semanas, label, cor, linha)
bars = [
    # Fase 1 — Setup (Sem 1 inteiro)
    (0, 1.62, 3.45, "Workspace · canais · permissões · estrutura por agente", SLACK_PURPLE, 0),
    # Fase 2 — Integração (Sem 1 parcial + Sem 2 parcial)
    (1, 1.62, 5.5, "Flows SF → Slack · thresholds DAF · alertas por canal · testes", SF_BLUE, 1),
    # Fase 3 — Automações (Sem 2 + início Sem 3)
    (2, 5.17, 5.0, "RAID Log (Listas Slack) · Status Report automático · validação", SLACK_GREEN, 2),
    # Fase 4 — Treinamento (Sem 3)
    (3, 8.72, 3.45, "Workshop 1 (técnico 4h) · Workshop 2 (gestores 4h) · Handoff", SLACK_YELLOW, 3),
]
for (li, bx, bw, bl, bc, _) in bars:
    ly = lane_y[li]
    add_rect(s2, bx, ly+0.1, bw, lane_h-0.2, fill=bc)
    add_text(s2, bl, bx+0.08, ly+0.2, bw-0.16, lane_h-0.4,
             size=8, color=WHITE, wrap=True)

# Linha de hoje
add_rect(s2, 1.6, 1.28, 0.03, 4.4, fill=SLACK_RED)
add_text(s2, "▼ Kick-off", 1.35, 1.22, 1.0, 0.25, size=8, bold=True, color=SLACK_RED)

# Legenda duração
add_text(s2, "Duração total estimada: 8–10 dias de trabalho técnico  |  Workshops: 2 × 4h",
         1.6, 5.72, 11, 0.3, size=9, color=DGRAY, italic=True)

# ─────────────────────────────────────────────
# SLIDE 3 — FASES 1 E 2 (detalhe)
# ─────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
header(s3, "Fases 1 e 2 — Setup e Integração Salesforce → Slack",
       "Semana 1 · Arquiteto Técnico · 4–6 dias")
footer(s3, "Slide 3 de 6")

add_rect(s3, 0, 1.22, 13.33, 5.8, fill=LGRAY)

# Fase 1
add_rect(s3, 0.2, 1.35, 6.1, 5.5, fill=WHITE,
         line=SLACK_PURPLE, line_w=Pt(1.5))
add_rect(s3, 0.2, 1.35, 6.1, 0.42, fill=SLACK_PURPLE)
add_text(s3, "FASE 1 — Setup do Workspace  |  2–3 dias  |  Arquiteto Técnico",
         0.3, 1.38, 5.9, 0.35, size=10, bold=True, color=WHITE)

ativs1 = [
    ("Criação/configuração do workspace Slack da Dataprev",
     "Greenfield ou aproveitamento de workspace existente — confirmar na Semana 0"),
    ("Estrutura de canais por agente, pilar e papel",
     "7 canais estruturais + 1 canal por agente em produção (mín. 5 canais de agente)"),
    ("Configuração de permissões e perfis de acesso",
     "Quem vê o quê — ministérios não veem canais de outros ministérios"),
    ("Canais privados para war room e comitê de governança",
     "#war-room-incidentes e #governanca-dataprev com acesso restrito"),
    ("Nomenclatura padrão aprovada com o cliente",
     "Definir no kick-off — renomeações posteriores quebram automações"),
]
for i, (titulo, detalhe) in enumerate(ativs1):
    ry = 1.9 + i * 0.88
    add_rect(s3, 0.28, ry, 5.95, 0.78, fill=RGBColor(0xF5,0xEE,0xF5),
             line=SLACK_PURPLE, line_w=Pt(0.5))
    add_rect(s3, 0.28, ry, 0.08, 0.78, fill=SLACK_PURPLE)
    add_text(s3, titulo, 0.45, ry+0.04, 5.6, 0.28, size=9.5, bold=True, color=DARK)
    add_text(s3, detalhe, 0.45, ry+0.33, 5.6, 0.35, size=8.5, color=DGRAY, italic=True)

# Fase 2
add_rect(s3, 6.8, 1.35, 6.3, 5.5, fill=WHITE,
         line=SF_BLUE, line_w=Pt(1.5))
add_rect(s3, 6.8, 1.35, 6.3, 0.42, fill=SF_BLUE)
add_text(s3, "FASE 2 — Integração SF → Slack  |  2–3 dias  |  Arquiteto Técnico",
         6.9, 1.38, 6.1, 0.35, size=10, bold=True, color=WHITE)

ativs2 = [
    ("Instalar e configurar app Slack for Salesforce",
     "Via AppExchange — requer permissão de System Admin no org SF"),
    ("Mapear KPIs e thresholds do DAF Health Score",
     "Definir com Dataprev quais valores disparam alerta leve, crítico e war room"),
    ("Configurar Flows no Salesforce para disparar notificações",
     "Record-triggered flows com action 'Send Slack Message' por canal"),
    ("Validar alertas por canal (agente · executivo · war room)",
     "Teste com dados reais de 1 agente antes de expandir para todos"),
    ("Documentar thresholds acordados e lógica dos alertas",
     "Runbook de alertas — quem recebe, o que fazer, em quanto tempo responder"),
]
for i, (titulo, detalhe) in enumerate(ativs2):
    ry = 1.9 + i * 0.88
    add_rect(s3, 6.88, ry, 6.15, 0.78, fill=RGBColor(0xEE,0xF5,0xFF),
             line=SF_BLUE, line_w=Pt(0.5))
    add_rect(s3, 6.88, ry, 0.08, 0.78, fill=SF_BLUE)
    add_text(s3, titulo, 7.05, ry+0.04, 5.85, 0.28, size=9.5, bold=True, color=DARK)
    add_text(s3, detalhe, 7.05, ry+0.33, 5.85, 0.35, size=8.5, color=DGRAY, italic=True)

# ─────────────────────────────────────────────
# SLIDE 4 — FASES 3 E 4 (detalhe)
# ─────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
header(s4, "Fases 3 e 4 — Automações, Treinamento e Handoff",
       "Semanas 2–3 · GP + Consultor CM + Arquiteto Técnico")
footer(s4, "Slide 4 de 6")

add_rect(s4, 0, 1.22, 13.33, 5.8, fill=LGRAY)

# Fase 3
add_rect(s4, 0.2, 1.35, 6.1, 5.5, fill=WHITE,
         line=SLACK_GREEN, line_w=Pt(1.5))
add_rect(s4, 0.2, 1.35, 6.1, 0.42, fill=SLACK_GREEN)
add_text(s4, "FASE 3 — Automações e RAID Log  |  2 dias  |  GP + Arquiteto",
         0.3, 1.38, 5.9, 0.35, size=10, bold=True, color=WHITE)

ativs3 = [
    ("Configurar Listas do Slack para RAID Log dinâmico",
     "Colunas: Tipo · Descrição · Responsável · Prazo · Status · Agente afetado"),
    ("Definir fluxo de criação e atualização de itens do RAID",
     "Quem cria · quem atribui · quem fecha — SLA de atualização por tipo"),
    ("Configurar automação do Status Report semanal",
     "Template narrativo postado automaticamente toda segunda-feira no canal executivo"),
    ("Validar primeiro ciclo de Status Report com a Dataprev",
     "Aprovação do formato antes de automatizar — ajuste de tom e métricas exibidas"),
    ("Configurar Slackbot para resumo de threads longas",
     "Canal #war-room-incidentes e #comite-governanca — resumo automático para novos membros"),
]
for i, (titulo, detalhe) in enumerate(ativs3):
    ry = 1.9 + i * 0.88
    add_rect(s4, 0.28, ry, 5.95, 0.78, fill=RGBColor(0xEE,0xF8,0xF2),
             line=SLACK_GREEN, line_w=Pt(0.5))
    add_rect(s4, 0.28, ry, 0.08, 0.78, fill=SLACK_GREEN)
    add_text(s4, titulo, 0.45, ry+0.04, 5.6, 0.28, size=9.5, bold=True, color=DARK)
    add_text(s4, detalhe, 0.45, ry+0.33, 5.6, 0.35, size=8.5, color=DGRAY, italic=True)

# Fase 4
add_rect(s4, 6.8, 1.35, 6.3, 5.5, fill=WHITE,
         line=SLACK_YELLOW, line_w=Pt(1.5))
add_rect(s4, 6.8, 1.35, 6.3, 0.42, fill=SLACK_YELLOW)
add_text(s4, "FASE 4 — Treinamento e Handoff  |  3 dias  |  Consultor CM + GP",
         6.9, 1.38, 6.1, 0.35, size=10, bold=True, color=DARK)

ativs4 = [
    ("Workshop 1 — Equipe Técnica Dataprev (4h)",
     "Estrutura de canais · alertas · RAID Log · resposta a incidentes via Slack"),
    ("Workshop 2 — Gestores e Stakeholders Ministérios (4h)",
     "Leitura de canais · protocolo de crise · uso do war room · Status Reports"),
    ("Entrega dos guias de uso rápido por papel",
     "1 página por perfil: Gestor · Técnico · Representante de Ministério"),
    ("Documentação da arquitetura do workspace",
     "Canais · permissões · integrações · thresholds — versionado no canal #aprendizado-agentes"),
    ("Handoff formal ProServ → CSM",
     "Contexto completo preservado nos canais Slack — não depende de memória pessoal"),
]
for i, (titulo, detalhe) in enumerate(ativs4):
    ry = 1.9 + i * 0.88
    add_rect(s4, 6.88, ry, 6.15, 0.78, fill=RGBColor(0xFF,0xFB,0xEE),
             line=SLACK_YELLOW, line_w=Pt(0.5))
    add_rect(s4, 6.88, ry, 0.08, 0.78, fill=SLACK_YELLOW)
    add_text(s4, titulo, 7.05, ry+0.04, 5.85, 0.28, size=9.5, bold=True, color=DARK)
    add_text(s4, detalhe, 7.05, ry+0.33, 5.85, 0.35, size=8.5, color=DGRAY, italic=True)

# ─────────────────────────────────────────────
# SLIDE 5 — PREMISSAS
# ─────────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
header(s5, "Premissas Críticas da Implantação",
       "Itens que devem ser confirmados antes ou no kick-off — impacto direto no prazo")
footer(s5, "Slide 5 de 6")

add_rect(s5, 0, 1.22, 13.33, 5.8, fill=LGRAY)

premissas = [
    (SLACK_RED,    "BLOQUEANTE",
     "Licença Slack adquirida antes do kick-off",
     "Sem licença, nenhuma atividade do Pilar 6 pode começar. Plano mínimo: Slack Pro."),
    (SLACK_RED,    "BLOQUEANTE",
     "Licença inclui Listas, Workflows e integração via API",
     "Slack Free não suporta automações. Necessário Slack Pro ou superior."),
    (SLACK_RED,    "BLOQUEANTE",
     "Administrador Slack designado na Dataprev com acesso de Owner",
     "Sem Owner, PS não consegue configurar canais privados e permissões."),
    (SLACK_YELLOW, "CRÍTICA",
     "App 'Slack for Salesforce' instalável no org SF (AppExchange)",
     "Ambientes com restrição de AppExchange exigem aprovação prévia de TI — iniciar na Semana 0."),
    (SLACK_YELLOW, "CRÍTICA",
     "Thresholds do DAF Health Score definidos até o Dia 5",
     "Sem thresholds acordados, os alertas automáticos não podem ser configurados."),
    (SLACK_YELLOW, "CRÍTICA",
     "Sponsor executivo aprova Slack como canal oficial antes do kick-off",
     "Sem adesão dos ministérios, o Pilar 3 opera em modo degradado."),
    (SLACK_GREEN,  "IMPORTANTE",
     "Mapeamento de usuários (perfis + ministérios) entregue na Semana 1",
     "Atraso impacta configuração de canais privados e permissões."),
    (SLACK_GREEN,  "IMPORTANTE",
     "Nomenclatura de canais definida no kick-off e não alterada",
     "Renomeações posteriores quebram automações e links nos runbooks."),
    (SLACK_GREEN,  "IMPORTANTE",
     "Plano com retenção de mensagens ativada para auditoria TCU/CGU",
     "Slack Business+ ou Enterprise Grid — retenção ilimitada para eDiscovery."),
    (SLACK_BLUE,   "ASSUMIDA",
     "Ambiente SF estável e APIs disponíveis durante a implantação",
     "Restrições de firewall ou whitelist podem adicionar 3–5 dias ao esforço técnico."),
]

cols = [0, 1]
rows_per_col = 5
col_x = [0.22, 6.78]
col_w = 6.35

for i, (cor, nivel, titulo, detalhe) in enumerate(premissas):
    col = i // rows_per_col
    row = i % rows_per_col
    cx = col_x[col]
    ry = 1.42 + row * 1.02

    add_rect(s5, cx, ry, col_w, 0.92, fill=WHITE,
             line=cor, line_w=Pt(1))
    add_rect(s5, cx, ry, 0.1, 0.92, fill=cor)

    # Badge nível
    badge_colors = {
        "BLOQUEANTE": SLACK_RED,
        "CRÍTICA": SLACK_YELLOW,
        "IMPORTANTE": SLACK_GREEN,
        "ASSUMIDA": SLACK_BLUE,
    }
    badge_text_colors = {
        "BLOQUEANTE": WHITE,
        "CRÍTICA": DARK,
        "IMPORTANTE": WHITE,
        "ASSUMIDA": WHITE,
    }
    bc = badge_colors.get(nivel, MGRAY)
    btc = badge_text_colors.get(nivel, DARK)
    add_rect(s5, cx + col_w - 1.1, ry + 0.07, 1.0, 0.22, fill=bc)
    add_text(s5, nivel, cx + col_w - 1.1, ry + 0.08, 1.0, 0.2,
             size=7, bold=True, color=btc, align=PP_ALIGN.CENTER)

    add_text(s5, titulo, cx + 0.18, ry + 0.06, col_w - 1.3, 0.28,
             size=9.5, bold=True, color=DARK)
    add_text(s5, detalhe, cx + 0.18, ry + 0.38, col_w - 0.28, 0.44,
             size=8.5, color=DGRAY, italic=True)

# ─────────────────────────────────────────────
# SLIDE 6 — ESTRUTURA DE CANAIS
# ─────────────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
header(s6, "Estrutura de Canais Sugerida",
       "7 canais estruturais + 1 canal por agente · nomenclatura padrão a aprovar no kick-off")
footer(s6, "Slide 6 de 6")

add_rect(s6, 0, 1.22, 13.33, 5.8, fill=LGRAY)

canais = [
    (SLACK_PURPLE, "#governanca-dataprev",
     "Pilar 4 — Governança",
     "Comitê periódico · atas · decisões · follow-ups · acesso restrito a liderança"),
    (SLACK_RED,    "#war-room-incidentes",
     "Pilar 2 — Crise",
     "Canal de acionamento imediato · protocolo pré-definido · todos os perfis técnicos"),
    (SF_BLUE,      "#monitoramento-operacional",
     "Pilar 1 — Monitoramento",
     "Dashboards diários · alertas automáticos DAF · uptime · fallbacks · override patterns"),
    (SLACK_YELLOW, "#raid-log-riscos",
     "Pilar 2 — RAID",
     "Listas Slack · riscos e issues em tempo real · responsável + prazo por item"),
    (SLACK_GREEN,  "#status-reports",
     "Pilar 3 — Comunicação",
     "Reports semanais automáticos · narrativa executiva · toda segunda-feira"),
    (RGBColor(0x62,0x64,0x6A), "#aprendizado-agentes",
     "Pilar 5 — Transferência",
     "Runbooks · guias de dashboards · protocolos de crise · memória institucional"),
    (SLACK_BLUE,   "#agente-[nome] (1 por agente)",
     "Pilares 1–2 — Por Agente",
     "Ex: #agente-bolsa-familia · #agente-inss-pericias · alertas e decisões por agente"),
]

card_w = 3.8
card_h = 1.4
positions = [
    (0.22, 1.38), (4.22, 1.38), (8.22, 1.38),
    (0.22, 2.95), (4.22, 2.95), (8.22, 2.95),
    (2.22, 4.52),
]

for (cx, cy), (cor, nome, pilar, desc) in zip(positions, canais):
    add_rect(s6, cx, cy, card_w, card_h, fill=WHITE,
             line=cor, line_w=Pt(1.5))
    add_rect(s6, cx, cy, card_w, 0.38, fill=cor)
    add_text(s6, nome, cx + 0.1, cy + 0.05, card_w - 0.15, 0.28,
             size=10, bold=True, color=WHITE)
    add_rect(s6, cx + card_w - 1.3, cy, 1.3, 0.38,
             fill=RGBColor(0,0,0) if False else cor)
    add_text(s6, pilar, cx + 0.1, cy + 0.44, card_w - 0.2, 0.22,
             size=8, bold=True, color=cor)
    add_text(s6, desc, cx + 0.1, cy + 0.66, card_w - 0.2, 0.65,
             size=8.5, color=DGRAY, wrap=True)

add_text(s6,
         "Nota: canais de agente são criados conforme agentes entram em produção — não todos de uma vez.",
         0.22, 6.05, 12.9, 0.28, size=8.5, color=DGRAY, italic=True)

# ─────────────────────────────────────────────
# SALVAR
# ─────────────────────────────────────────────
path = "/Users/nfilho/claude/DATAPREV_Slack_Timeline.pptx"
prs.save(path)
print(f"Salvo em: {path}")
