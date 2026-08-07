#!/usr/bin/env python3
"""ROM Generator — CLARO JSON V2 Migração v2 Attributes"""

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import copy, os

TEMPLATE = '/Users/nfilho/Downloads/URT_UNIFIED_ROM_TEMPLATE_PT (4).pptx'
OUTPUT   = '/Users/nfilho/claude/CLARO_JSON_V2_ROM.pptx'

prs = Presentation(TEMPLATE)

def set_tf(shape, lines, size=None, bold=None, color=None):
    """Replace all text in a text-frame shape with lines (list of str or (str,bold,size))."""
    tf = shape.text_frame
    tf.word_wrap = True
    # clear existing paragraphs beyond first
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)
    first = True
    for line in lines:
        if isinstance(line, tuple):
            txt, b, s = line[0], (line[1] if len(line)>1 else bold), (line[2] if len(line)>2 else size)
        else:
            txt, b, s = line, bold, size
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            from pptx.oxml.ns import qn
            from lxml import etree
            p_elem = copy.deepcopy(tf.paragraphs[-1]._p)
            # clear runs
            for r in p_elem.findall(qn('a:r')):
                p_elem.remove(r)
            tf._txBody.append(p_elem)
            para = tf.paragraphs[-1]
        if para.runs:
            run = para.runs[0]
        else:
            from pptx.oxml.ns import qn
            from lxml import etree
            r_elem = etree.SubElement(para._p, qn('a:r'))
            etree.SubElement(r_elem, qn('a:t'))
            run = para.runs[0]
        run.text = txt
        if s:  run.font.size = Pt(s)
        if b is not None: run.font.bold = b
        if color: run.font.color.rgb = RGBColor(*color)

def replace_text(shape, old, new):
    """Simple find-replace inside a shape's text frame."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)

def fill_table(shape, data):
    """Fill table with data (list of rows, each row is list of str)."""
    table = shape.table
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            if r < len(table.rows) and c < len(table.columns):
                cell = table.cell(r, c)
                if cell.text_frame.paragraphs:
                    para = cell.text_frame.paragraphs[0]
                    if para.runs:
                        para.runs[0].text = val
                    else:
                        from pptx.oxml.ns import qn
                        from lxml import etree
                        r_elem = etree.SubElement(para._p, qn('a:r'))
                        t_elem = etree.SubElement(r_elem, qn('a:t'))
                        t_elem.text = val

# ── Helper: find shape by name fragment ──────────────────────────────────────
def shape_by_name(slide, fragment):
    for s in slide.shapes:
        if fragment in s.name:
            return s
    return None

def shapes_with_text(slide, fragment):
    results = []
    for s in slide.shapes:
        if s.has_text_frame and fragment in s.text_frame.text:
            results.append(s)
    return results

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Capa
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[0]
for shape in s.shapes:
    if not shape.has_text_frame: continue
    t = shape.text_frame.text
    if 'Inserir Nome do Cliente' in t:
        set_tf(shape, ['Claro Brasil'], size=21, bold=True)
    elif 'XXX Cloud' in t:
        set_tf(shape, [
            ('Communications Cloud', True, 33),
            ('Migração JSONAttribute v1 → v2 | Jornada Única', False, 22)
        ])
    elif 'Inserir_Papel_Membro_Equipe' in t or 'INSERIR MÊS E ANO' in t:
        set_tf(shape, [
            ('Nelson Stebulaitis Filho — Engagement Manager', False, 10),
            ('Odair Civelli Junior — Solution Architect', False, 10),
            ('Fabricio Maia — Technical Lead', False, 10),
            ('Luciano Ricardo Emidio e Silva — Sponsor', False, 10),
            ('Junho 2026', True, 10),
        ])

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Agenda (manter estrutura, só trocar logo placeholder)
# ═════════════════════════════════════════════════════════════════════════════
# Logo placeholder — leave as-is (no logo file available)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Missão do Cliente
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[6]
for shape in s.shapes:
    if shape.has_text_frame and 'Inserir Missão do Cliente' in shape.text_frame.text:
        set_tf(shape, [
            'Entendemos que sua missão principal é modernizar a infraestrutura de BSS da Claro Brasil, '
            'estabelecendo uma plataforma de catálogo ágil que suporte a Jornada Única — '
            'eliminando a dívida técnica acumulada no modelo de atributos legado (JSONAttribute v1) '
            'e habilitando inovação rápida em produtos convergentes (Móvel, TV, Banda Larga, Fixo e Aparelhos).'
        ], size=17)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Motores Principais
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[7]
for shape in s.shapes:
    if shape.has_text_frame and 'Motor Principal' in shape.text_frame.text:
        set_tf(shape, [
            ('Esta iniciativa é impulsionada por:', False, 17),
            ('Performance crítica: áreas de negócio da Claro exigem melhoria de performance na Jornada Única — '
             'a otimização do código nativo depende da migração para o modelo v2.', False, 15),
            ('Prazo competitivo: IBM estimou 9 meses para o mesmo escopo; '
             'a Salesforce PS propõe 2 meses com aceleração via IA — diferencial estratégico decisivo.', False, 15),
            ('Dívida técnica acumulada: ~200 componentes custom manipulam o campo JSONAttribute__c de forma legada, '
             'bloqueando evoluções nativas do produto Salesforce Communications Cloud.', False, 15),
        ], size=15)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Objetivos de Negócios (tabela 5x2)
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[8]
for shape in s.shapes:
    if hasattr(shape, 'table'):
        fill_table(shape, [
            ['Objetivo de Negócio', 'KPI / Resultado Esperado'],
            ['Migrar 100% dos componentes custom para JSONAttribute v2',
             'Zero componentes em modo legado após go-live'],
            ['Melhorar performance da Jornada Única',
             'Eliminação de gargalos de código nativo dependentes do backlog de produto'],
            ['Reduzir dívida técnica e risco de upgrade futuro',
             'Conformidade total com framework Salesforce Communications Cloud'],
            ['Entregar em 2 meses (vs. 9 meses IBM)',
             'Aceleração de 78% via uso intensivo de IA no ciclo de desenvolvimento'],
        ])

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Desafios
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[9]
for shape in s.shapes:
    if shape.has_text_frame and 'Desafios que superaremos' in shape.text_frame.text:
        set_tf(shape, [
            ('Os Desafios que Superaremos Juntos…', False, 17),
            ('Volume e crescimento de componentes: estimativa cresceu de 127 (2023) para ~200 componentes impactados — '
             'requer confirmação formal no refinamento técnico (M-01).', False, 14),
            ('Complexidade de exceções: componentes com lógica avançada em APEX e LWC exigem '
             'tratamento manual por arquitetos — não automatizável por IA.', False, 14),
            ('Dependências de integrações externas: consumidores externos do campo JSONAttribute__c '
             'precisam ser mapeados e validados (lacuna crítica identificada no USD).', False, 14),
            ('Governança de backlog das squads: risco de interferência da migração '
             'nas entregas correntes das squads da Claro durante o período de convivência v1/v2.', False, 14),
            ('Ambiente de testes: sandbox Full Copy obrigatória para testes de volumetria — '
             'critérios específicos de volumetria ainda pendentes de confirmação pela Claro.', False, 14),
        ], size=14)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Visão Parceria
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[10]
for shape in s.shapes:
    if shape.has_text_frame and 'transformar essa visão em realidade' in shape.text_frame.text:
        set_tf(shape, [
            'Nosso objetivo é fazer parceria com a Claro Brasil para transformar essa visão em realidade — '
            'entregando a migração completa do modelo de atributos v1 → v2 em 2 meses, '
            'com aceleração via Inteligência Artificial, equipe interna Salesforce e framework ADP, '
            'garantindo conformidade técnica, performance e zero impacto funcional nas operações em curso.'
        ], size=18)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Visão do Futuro
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[12]
for shape in s.shapes:
    if not shape.has_text_frame: continue
    t = shape.text_frame.text
    if 'cada equipe é empoderada' in t:
        set_tf(shape, [
            'Imagine um futuro onde cada squad de desenvolvimento da Claro',
            'opera sobre um catálogo unificado, ágil e sem dívida técnica…'
        ], size=17)
    elif 'cada interação com o cliente' in t:
        set_tf(shape, [
            'Onde lançamentos de novos produtos convergentes são acelerados,',
            'sem dependência de migrações manuais ou barreiras de upgrade…'
        ], size=17)
    elif 'tecnologia não é mais uma barreira' in t:
        set_tf(shape, [
            'Onde a plataforma Salesforce evolui de forma nativa,',
            'sustentando a Jornada Única com performance e confiabilidade.'
        ], size=17)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Tabela de Valor (5x3)
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[13]
for shape in s.shapes:
    if hasattr(shape, 'table'):
        fill_table(shape, [
            ['Dimensão de Valor', 'Estado Atual (AS-IS)', 'Estado Futuro (TO-BE)'],
            ['Performance da Jornada Única',
             'Gargalos de código nativo — backlog bloqueado por atributos legados',
             'Código nativo otimizado; backlog de produto desbloqueado'],
            ['Velocidade de Entrega',
             'Migração estimada em 9 meses (IBM) com abordagem componente a componente',
             '2 meses com refatoração automatizada via IA (framework ADP)'],
            ['Dívida Técnica',
             '~200 componentes custom em modo legado (JSONAttribute v1)',
             'Zero componentes legados; conformidade plena com Communications Cloud v2'],
            ['Risco de Upgrade',
             'Incompatibilidade crescente a cada release do produto Salesforce',
             'Alinhamento nativo com roadmap Salesforce — upgrades sem retrabalho'],
        ])

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Arquitetura
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[16]
for shape in s.shapes:
    if shape.has_text_frame and 'Nosso Projeto de Solução Recomendado' in shape.text_frame.text:
        set_tf(shape, [
            ('Nosso Projeto de Solução Recomendado…', False, 17),
            ('A arquitetura concentra-se na migração de ~200 componentes customizados '
             '(Data Raptors, Integration Procedures, FlexCards, OmniScripts, Classes Apex e LWC) '
             'do modelo JSONAttribute v1 para o novo modelo v2 '
             '(vlocity_cmt__AttributeSelectedValues__c), controlado via feature flag '
             'EnableV2AttributeModel. O período de convivência (co-living) entre v1 e v2 '
             'garante zero impacto nos sistemas em produção durante a transição.', False, 11),
        ])
    elif shape.has_text_frame and 'Inserir MAPA da Arquitetura' in shape.text_frame.text:
        set_tf(shape, [
            'Fluxo: Identificação via IA → Refatoração Automatizada (ADP) → '
            'Testes Unitários (IA) → Testes de Volumetria (Sandbox Full Copy) → '
            'Validação OK? → Migração em Produção (lotes/finais de semana) → '
            'Suporte Pós-Lançamento (4 semanas)\n'
            '[Período de Convivência v1/v2 durante toda a fase de desenvolvimento]'
        ], size=12)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Capacidades Salesforce
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[17]
for shape in s.shapes:
    if shape.has_text_frame and 'capacidades do Salesforce' in shape.text_frame.text:
        set_tf(shape, [
            ('As Capacidades que Iremos Aproveitar', False, 17),
            ('Communications Cloud / OmniStudio: plataforma nativa de catálogo de produtos '
             'e componentes low-code (Data Raptors, FlexCards, OmniScripts, IPs) — '
             'base dos ~200 componentes a migrar.', True, 11),
            ('Framework ADP (Agile Delivery Platform): framework interno Salesforce PS '
             'para automação do ciclo de desenvolvimento — obrigatório para uso das ferramentas de IA.', True, 11),
            ('Inteligência Artificial (Claude / Copilot): usada em todo o ciclo da Atividade 1 '
             '(checkout → refatoração → testes unitários → checkin), reduzindo o esforço de 4 para 2 meses.', True, 11),
            ('JSONAttributeSupport API: método nativo Salesforce que, quando utilizado, '
             'isenta o componente de alteração — componentes que já o utilizam estão fora do escopo.', True, 11),
        ])

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Decisões Arquitetônicas
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[18]
for shape in s.shapes:
    if shape.has_text_frame and 'Decisões Arquitetônicas' in shape.text_frame.text:
        set_tf(shape, [
            ('As Decisões Arquitetônicas Chave e Justificativa', False, 17),
            ('Decisão: Estimativa por modelo de IA (por tipo/cenário), não por componente individual. '
             'Justificativa: com ~200 componentes, a contagem individual não escala; '
             'modelos por cenário (Data Raptor, APEX, LWC etc.) são mais precisos e automatizáveis.', True, 11),
            ('Decisão: Período de convivência v1/v2 durante todo o desenvolvimento. '
             'Justificativa: garante que os ambientes em produção continuem operacionais '
             'enquanto os componentes são migrados em lotes, eliminando risco de downtime.', True, 11),
            ('Decisão: Migração de dados em lotes por objeto, nos finais de semana. '
             'Justificativa: objetos grandes (Assets, Inventory Items, Order Items) '
             'não podem ser migrados de uma só vez sem risco de impacto operacional.', True, 11),
            ('Decisão: Time 100% interno Salesforce PS. '
             'Justificativa: uso do framework ADP exige profissionais com acesso às ferramentas internas; '
             'parceiros externos não têm esse acesso.', True, 11),
        ])

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Princípios Orientadores
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[19]
for shape in s.shapes:
    if shape.has_text_frame and 'Princípios Orientadores' in shape.text_frame.text:
        set_tf(shape, [
            ('Nossos Princípios Orientadores para o Sucesso da Claro…', False, 17),
            ('IA como Acelerador, Humanos como Árbitros: a IA automatiza as atividades repetitivas; '
             'arquitetos técnicos tratam os cenários de exceção e realizam a revisão final — '
             'garantindo qualidade sem abrir mão de velocidade.', True, 11),
            ('Migração Puramente Técnica: o escopo é estritamente técnico — '
             'sem revisão funcional, sem alteração de regras de negócio. '
             'Premissa crítica para manter o cronograma de 2 meses.', True, 11),
            ('Co-Living como Estratégia de Risco Zero: v1 e v2 coexistem durante todo o projeto; '
             'o cutover final ocorre apenas após validação técnica, funcional e de negócios (UAT).', True, 11),
            ('Processamento Local — Sem Exposição da Org: a IA opera sobre repositórios locais; '
             'não há conexão direta com a org Salesforce da Claro, eliminando riscos de segurança e '
             'questões legais de autorização.', True, 11),
        ])

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Resumo do Escopo (tabela 5x3)
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[21]
for shape in s.shapes:
    if hasattr(shape, 'table'):
        fill_table(shape, [
            ['Atividade', 'Responsável', 'Duração'],
            ['0. Identificação de Componentes e Governança',
             'Salesforce PS', '1 semana (M-01)'],
            ['1. Adaptar e Testar Componentes v1/v2 (caminho crítico)',
             'Salesforce PS', '2 meses (M-01 a M-02)'],
            ['2. Migração Contínua de Dados em Produção (XLi)',
             'Salesforce PS entrega plano/scripts; Claro executa', '4 semanas (paralelo a 1)'],
            ['3-5. Habilitar v2 (não-prod) → Teste Regressivo → Habilitar em Produção',
             'Claro executa; Salesforce suporta', '2 semanas + 1 dia'],
        ])

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Atividades e Entregáveis
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[22]
for shape in s.shapes:
    if not shape.has_text_frame: continue
    t = shape.text_frame.text
    if 'Atividades' in t and 'Entregáveis' not in t:
        set_tf(shape, [
            ('Atividades', True, 15),
            ('• Refinamento técnico e listagem dos ~200 componentes impactados', False, 12),
            ('• Categorização por tipo (Data Raptor, IP, FlexCard, OmniScript, APEX, LWC)', False, 12),
            ('• Treinamento da IA por arquitetos por cenário de componente', False, 12),
            ('• Refatoração automatizada + testes unitários (IA)', False, 12),
            ('• Tratamento manual de exceções por arquitetos técnicos', False, 12),
            ('• Testes integrados e regressivos (Claro executa)', False, 12),
            ('• Preparação e execução da migração de dados em lotes', False, 12),
            ('• Suporte pós-lançamento por 4 semanas', False, 12),
        ])
    elif 'Entregáveis' in t:
        set_tf(shape, [
            ('Entregáveis', True, 15),
            ('• Lista consolidada de componentes impactados com criticidade', False, 12),
            ('• Código migrado para v2 (repositório)', False, 12),
            ('• Testes unitários gerados por IA para cada componente', False, 12),
            ('• Plano e scripts de migração de dados em produção', False, 12),
            ('• Relatório de testes integrados e regressivos', False, 12),
            ('• Documentação de exceções e cenários tratados manualmente', False, 12),
            ('• Release aprovado em produção', False, 12),
        ])

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Dentro e Fora do Escopo
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[23]
for shape in s.shapes:
    if not shape.has_text_frame: continue
    t = shape.text_frame.text
    if 'Dentro do Escopo' in t:
        set_tf(shape, [
            ('DENTRO DO ESCOPO', True, 14),
            ('✓ Todos os componentes custom que lêem/escrevem vlocity_cmt__JSONAttribute__c', False, 11),
            ('✓ OmniStudio: Data Raptors, Integration Procedures, FlexCards, OmniScripts', False, 11),
            ('✓ Classes APEX e LWC (custom) que manipulam o campo JSON', False, 11),
            ('✓ Testes unitários por componente (gerados via IA)', False, 11),
            ('✓ Plano e scripts de migração de dados em produção', False, 11),
            ('✓ Suporte durante testes integrados e regressivos (Claro executa)', False, 11),
            ('✓ Suporte pós-lançamento: 4 semanas', False, 11),
        ])
    elif 'Fora do Escopo' in t:
        set_tf(shape, [
            ('FORA DO ESCOPO', True, 14),
            ('✗ Classes Apex que utilizam o método nativo JSONAttributeSupport (já compatíveis)', False, 11),
            ('✗ Data Raptors que utilizam @Attribute (não extraem o campo JSON)', False, 11),
            ('✗ Componentes nativos Salesforce (já suportam v1 e v2)', False, 11),
            ('✗ Revisão funcional ou alteração de regras de negócio', False, 11),
            ('✗ Execução dos testes integrados e regressivos (responsabilidade da Claro)', False, 11),
            ('✗ Estabilização do ambiente DEV (responsabilidade da Claro)', False, 11),
            ('✗ Execução da migração de dados em produção (Salesforce entrega plano; Claro executa)', False, 11),
            ('✗ AMS/suporte pós período de 4 semanas', False, 11),
        ])

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 26 — Requisitos e Premissas
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[25]
for shape in s.shapes:
    if not shape.has_text_frame: continue
    t = shape.text_frame.text
    if 'Requisitos' in t and 'Pressupostos' not in t:
        set_tf(shape, [
            ('REQUISITOS', True, 14),
            ('• Sandbox Full Copy disponibilizada pela Claro para testes de volumetria', False, 11),
            ('• Lista consolidada de componentes impactados com criticidade (pendente da Claro)', False, 11),
            ('• Confirmação dos critérios de volumetria para testes em sandbox', False, 11),
            ('• Mapeamento de dependências externas que consomem JSONAttribute__c', False, 11),
            ('• Acesso ao repositório de código (local — sem conexão com org produtiva)', False, 11),
            ('⚠ LACUNA: Lista oficial dos 127+ componentes com nível de criticidade — pendente', False, 11),
            ('⚠ LACUNA: Estratégia de gestão do backlog das squads durante o período de migração', False, 11),
        ])
    elif 'Pressupostos' in t:
        set_tf(shape, [
            ('PREMISSAS', True, 14),
            ('• Migração puramente técnica — sem revisão funcional', False, 11),
            ('• Time 100% interno Salesforce PS (framework ADP)', False, 11),
            ('• IA opera em repositórios locais — sem acesso direto à org da Claro', False, 11),
            ('• Claro é responsável pela estabilização do ambiente DEV', False, 11),
            ('• Claro executa os testes integrados e regressivos; Salesforce suporta', False, 11),
            ('• Claro executa a migração de dados em produção com base nos scripts Salesforce', False, 11),
            ('• Período de convivência v1/v2 durante todo o desenvolvimento', False, 11),
            ('• Risco de escopo: se >200 cenários de exceção, +2 pessoas por 1 mês adicional', False, 11),
        ])

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 28 — Roteiro por Fases
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[27]
for shape in s.shapes:
    if shape.has_text_frame and 'Roteiro para Transformação' in shape.text_frame.text:
        set_tf(shape, [
            ('Seu Roteiro para a Migração JSON V2…', False, 17),
            ('Fase 1 — Refinamento e Identificação (M-01, 2 pessoas): '
             'confirmação da lista de componentes, categorização por tipo, '
             'aprovação TDB no COE. Entregável: lista consolidada com criticidade.', False, 14),
            ('Fase 2 — Desenvolvimento e Testes Unitários (M-01 a M-02, 3 pessoas com IA): '
             'refatoração automatizada por cenário, tratamento manual de exceções, '
             'testes unitários gerados por IA. Entregável: código migrado + testes.', False, 14),
            ('Fase 3 — Testes Integrados e Regressivos (M-05 a M-06, Claro executa, SF suporta 2p): '
             'validação funcional em paralelo v1/v2, testes de volumetria em Sandbox Full Copy.', False, 14),
            ('Fase 4 — Release e Migração de Dados em Produção (M-07 + rolling): '
             'deploy dos metadados, migração de dados em lotes por objeto (finais de semana).', False, 14),
            ('Fase 5 — Suporte Pós-Lançamento (4 semanas, M-07 a M-08, 2 pessoas): '
             'suporte técnico contínuo, correção de defeitos dentro da garantia.', False, 14),
        ])

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 29 — Timeline (tabela 11x12) — preencher linha de fases
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[28]
for shape in s.shapes:
    if hasattr(shape, 'table'):
        tbl = shape.table
        # Row 0: header months
        headers = ['Atividade', 'M-01', '', '', '', 'M-02', '', '', '', 'M-03+', '', '']
        for c, h in enumerate(headers):
            if c < len(tbl.columns):
                cell = tbl.cell(0, c)
                if cell.text_frame.paragraphs[0].runs:
                    cell.text_frame.paragraphs[0].runs[0].text = h
        # Key rows
        rows_data = [
            ['Refinamento Técnico', 'X', 'X', '', '', '', '', '', '', '', '', ''],
            ['Desenvolvimento & Testes Unitários (IA)', '', 'X', 'X', 'X', 'X', 'X', 'X', 'X', '', '', ''],
            ['Testes Integrados (Claro)', '', '', '', '', '', 'X', 'X', '', '', '', ''],
            ['Testes Regressivos (Claro)', '', '', '', '', '', '', 'X', 'X', '', '', ''],
            ['Release / Deploy Metadados', '', '', '', '', '', '', '', 'X', '', '', ''],
            ['Suporte Pós-Lançamento (4 sem)', '', '', '', '', '', '', '', '', 'X', 'X', 'X'],
            ['Migração de Dados em Produção (rolling)', '', 'X', 'X', 'X', 'X', 'X', 'X', 'X', 'X', 'X', ''],
            ['Co-Living v1/v2', '', 'X', 'X', 'X', 'X', 'X', 'X', '', '', '', ''],
            ['Marco: Aprovação TDB no COE', '', 'X', '', '', '', '', '', '', '', '', ''],
            ['Marco: Deploy Metadados em Produção', '', '', '', '', '', '', '', 'X', '', '', ''],
        ]
        for r, row in enumerate(rows_data):
            for c, val in enumerate(row):
                if r+1 < len(tbl.rows) and c < len(tbl.columns):
                    cell = tbl.cell(r+1, c)
                    if cell.text_frame.paragraphs[0].runs:
                        cell.text_frame.paragraphs[0].runs[0].text = val

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 41 — Equipe de Projeto
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[40]
for shape in s.shapes:
    if shape.has_text_frame and 'Equipo Projeto' in shape.text_frame.text:
        set_tf(shape, [
            ('EQUIPE SALESFORCE PS', True, 13),
            ('Nelson Stebulaitis Filho — Engagement Manager', False, 11),
            ('Odair Civelli Junior — Solution Architect / Tech Lead', False, 11),
            ('Fabricio Maia — Technical Lead (framework ADP)', False, 11),
            ('Tar [sobrenome pendente] — Technical Architect', False, 11),
            ('Dev 1 — Senior Developer (APEX/LWC)', False, 11),
            ('Dev 2 — Senior Developer (OmniStudio/IA)', False, 11),
            ('', False, 11),
            ('EQUIPE CLARO BRASIL', True, 13),
            ('Luciano Ricardo Emidio e Silva — Sponsor / Manager', False, 11),
            ('Juliana Brites — Stakeholder', False, 11),
            ('[Product Owner por LOB] — a confirmar', False, 11),
            ('[Equipe QA Claro] — responsável pelos testes integrados/regressivos', False, 11),
        ])

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 42 — Papéis e Responsabilidades (tabela 6x3)
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[41]
for shape in s.shapes:
    if hasattr(shape, 'table'):
        fill_table(shape, [
            ['Atividade', 'Salesforce PS', 'Claro Brasil'],
            ['Refinamento Técnico e Listagem de Componentes', 'Lidera', 'Participa'],
            ['Desenvolvimento e Testes Unitários (com IA)', 'Executa', 'Aprova'],
            ['Testes Integrados e Regressivos', 'Suporta (2p)', 'Executa'],
            ['Migração de Dados em Produção', 'Entrega plano e scripts', 'Executa'],
            ['Suporte Pós-Lançamento (4 semanas)', 'Executa', 'Monitora e valida'],
        ])

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 50 — Resumo do Investimento (tabela 4x2)
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[49]
for shape in s.shapes:
    if hasattr(shape, 'table'):
        fill_table(shape, [
            ['Componente', 'Estimativa'],
            ['Desenvolvimento (3 pessoas × 2 meses)', 'A confirmar após refinamento técnico (M-01)'],
            ['Suporte Pós-Lançamento (2 pessoas × 4 semanas)', 'Incluído no investimento total'],
            ['Risco (contingência): +2 pessoas × 1 mês', 'Opcional — ativado se exceções > estimado'],
        ])
for shape in s.shapes:
    if shape.has_text_frame and 'Detalhes & Condições do ROM' in shape.text_frame.text:
        set_tf(shape, [
            ('Detalhes & Condições do ROM', True, 12),
            ('Estimativa Não Vinculativa: os valores são para fins de planejamento. '
             'A proposta formal será elaborada após recebimento da lista consolidada de componentes '
             'e confirmação dos critérios de volumetria pela Claro.', False, 10),
            ('Premissa de escopo: ~200 componentes custom identificados. '
             'Se o volume real for significativamente superior, o esforço será revisado.', False, 10),
            ('Contingência de risco: +2 pessoas por 1 mês adicional pode ser ativada '
             'caso surjam cenários de exceção imprevistos no refinamento técnico.', False, 10),
            ('Próximo passo para fechar valor: Odair Civelli Junior compartilhará documento '
             'macro com pontos de estimativa detalhados (pendente).', False, 10),
        ])

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 52 — Próximos Passos
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides[51]
for shape in s.shapes:
    if shape.has_text_frame and 'Próximos Passos' in shape.text_frame.text:
        set_tf(shape, [
            ('Próximos Passos…', False, 17),
            ('1. Alinhamento do ROM (esta semana): '
             'revisão desta apresentação com Luciano e Odair para alinhar escopo, '
             'premissas e estrutura da proposta formal.', False, 14),
            ('2. Recebimento do documento macro de estimativas (Odair → Nelson): '
             'confirmação dos perfis, horas por fase e valor unitário.', False, 14),
            ('3. Confirmação da lista de componentes pela Claro: '
             'lista consolidada dos 127-200 componentes com criticidade — '
             'item crítico para fechar a proposta formal.', False, 14),
            ('4. Elaboração da proposta comercial formal (Nelson): '
             'SOW com escopo, premissas, valor e condições de pagamento.', False, 14),
            ('5. Kick-off do Refinamento Técnico (M-01): '
             'confirmação do ambiente DEV, sandbox Full Copy e início do trabalho.', False, 14),
        ])

# ═════════════════════════════════════════════════════════════════════════════
# SAVE
# ═════════════════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f'Saved: {OUTPUT}')
