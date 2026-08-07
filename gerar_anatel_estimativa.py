from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
WHITE    = RGBColor(0xFF,0xFF,0xFF)
NAVY     = RGBColor(0x03,0x2D,0x60)
SF_BLUE  = RGBColor(0x00,0x70,0xD2)
ORANGE   = RGBColor(0xFF,0x7A,0x00)
GRAY_BG  = RGBColor(0xF4,0xF6,0xF9)
GRAY_BOR = RGBColor(0xD8,0xDF,0xE6)
GRAY_TXT = RGBColor(0x54,0x65,0x7D)
RED      = RGBColor(0xC2,0x32,0x34)
RED_BG   = RGBColor(0xFF,0xED,0xED)
AMBER    = RGBColor(0xCC,0x7A,0x00)
AMB_BG   = RGBColor(0xFF,0xF3,0xCD)
GREEN    = RGBColor(0x2E,0x84,0x4E)
GRN_BG   = RGBColor(0xE6,0xF5,0xEC)
PURPLE   = RGBColor(0x5A,0x2D,0x82)
TEAL     = RGBColor(0x02,0xA3,0xBF)
DARK     = RGBColor(0x16,0x1C,0x28)

prs = Presentation()
prs.slide_width  = Cm(25.4)
prs.slide_height = Cm(14.29)

# ── Helpers ───────────────────────────────────────────────────────────────────
def blank():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    # white bg
    r = sl.shapes.add_shape(1, Cm(0),Cm(0),Cm(25.4),Cm(14.29))
    r.fill.solid(); r.fill.fore_color.rgb = WHITE; r.line.fill.background()
    return sl

def rect(sl, x,y,w,h, fc=None, lc=None, lw=1.0):
    s = sl.shapes.add_shape(1, Cm(x),Cm(y),Cm(w),Cm(h))
    if fc: s.fill.solid(); s.fill.fore_color.rgb = fc
    else:  s.fill.background()
    if lc: s.line.color.rgb = lc; s.line.width = Pt(lw)
    else:  s.line.fill.background()
    return s

def txt(sl, text, x,y,w,h, size=9, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = sl.shapes.add_textbox(Cm(x),Cm(y),Cm(w),Cm(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def header(sl, title, subtitle=''):
    rect(sl, 0,0,25.4,1.35, fc=NAVY)
    rect(sl, 0,1.35,25.4,0.1, fc=ORANGE)
    txt(sl, title, 0.4,0.1,20,0.75, size=13, bold=True, color=WHITE)
    if subtitle:
        txt(sl, subtitle, 0.4,0.85,20,0.42, size=8, color=RGBColor(0xB0,0xC4,0xDE), italic=True)
    txt(sl, 'ANATEL  |  ROM Ballpark', 18,0.4,7,0.5, size=8,
        color=RGBColor(0xB0,0xC4,0xDE), align=PP_ALIGN.RIGHT)

def chip(sl, text, x,y,w, fc, tc):
    rect(sl, x,y,w,0.38, fc=fc, lc=tc)
    txt(sl, text, x,y+0.01,w,0.36, size=7.5, bold=True, color=tc, align=PP_ALIGN.CENTER)

def bullet_block(sl, title, items, x,y,w,h, tc, bg=GRAY_BG):
    rect(sl, x,y,w,h, fc=bg, lc=tc, lw=1.5)
    txt(sl, title, x+0.12,y+0.1,w-0.2,0.4, size=8.5, bold=True, color=tc)
    iy = y+0.52
    for item in items:
        txt(sl, f'▸  {item}', x+0.12,iy,w-0.2,0.42, size=8, color=GRAY_TXT)
        iy += 0.42
    return iy

def risk_chip(sl, text, x, y, level='A'):
    colors = {'A':(RED,RED_BG), 'M':(AMBER,AMB_BG), 'B':(GREEN,GRN_BG)}
    tc,fc = colors.get(level,(GRAY_TXT,GRAY_BG))
    rect(sl, x,y,0.28,0.28, fc=tc)
    txt(sl, level, x,y,0.28,0.28, size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, text, x+0.35,y,7.0,0.3, size=7.5, color=GRAY_TXT)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Capa
# ═════════════════════════════════════════════════════════════════════════════
sl = blank()
rect(sl, 0,0,25.4,14.29, fc=NAVY)
rect(sl, 0,9.5,25.4,0.18, fc=ORANGE)

txt(sl,'ANATEL',           0.8,2.2,24,1.8, size=52,bold=True,color=WHITE,align=PP_ALIGN.LEFT)
txt(sl,'Inovação Digital', 0.8,4.0,24,1.2, size=28,bold=False,color=RGBColor(0xB0,0xC4,0xDE),align=PP_ALIGN.LEFT)
txt(sl,'ROM Ballpark  ·  Mapa Mental de Escopo  ·  Estimativa de Horas e Recursos',
    0.8,5.4,23,0.7, size=11,color=RGBColor(0x7A,0xA3,0xCC),align=PP_ALIGN.LEFT)

# tags clouds
clouds = [
    (0.8,  7.2, 'Agentforce',    SF_BLUE),
    (4.0,  7.2, 'CLM / Contracts', PURPLE),
    (8.5,  7.2, 'Marketing Cloud', ORANGE),
    (13.2, 7.2, 'Tableau',        TEAL),
    (16.8, 7.2, 'MuleSoft Atend.', GREEN),
    (21.0, 7.2, 'MuleSoft TFF',   RED),
]
for (cx,cy,ct,cc) in clouds:
    rect(sl, cx,cy,3.5,0.5, fc=RGBColor(0x05,0x1A,0x35), lc=cc)
    txt(sl, ct, cx,cy,3.5,0.5, size=9,bold=True,color=cc,align=PP_ALIGN.CENTER)

txt(sl,'Nelson Stebulaitis Filho  |  Salesforce PS LATAM  |  2026',
    0.8,13.3,20,0.6, size=9, color=RGBColor(0x64,0x74,0x8B))

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Mapa Mental (Linha de Raciocínio)
# ═════════════════════════════════════════════════════════════════════════════
sl = blank()
header(sl,'Mapa Mental — Linha de Raciocínio para Estimativa',
       'Lógica de construção: Escopo → Premissas → Riscos → Perfis → Horas')

# Centro
cx,cy,cw,ch = 10.5,5.5,4.4,1.5
rect(sl, cx,cy,cw,ch, fc=NAVY, lc=SF_BLUE, lw=2)
txt(sl,'ANATEL\nROM Ballpark', cx,cy,cw,ch, size=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER)

# Nodes radiais
nodes = [
    (0.3,  1.8,  4.8,1.2, '① Agentforce\n3 Agentes + Orquestrador', SF_BLUE,   SF_BLUE),
    (0.3,  3.5,  4.8,1.2, '② CLM\nContratos OOTB 5 etapas',         PURPLE,    PURPLE),
    (0.3,  5.2,  4.8,1.2, '③ Marketing Cloud\n3 Jornadas',          ORANGE,    ORANGE),
    (0.3,  6.9,  4.8,1.2, '④ Tableau\n5 Painéis / 10 KPIs',         TEAL,      TEAL),
    (16.5, 1.8,  8.3,1.2, '⑤ MuleSoft Atendimento\n7 endpoints / 15–20 APIs', GREEN, GREEN),
    (16.5, 3.5,  8.3,1.2, '⑥ MuleSoft TFF\n4 fontes / batch / ~10M registros', RED, RED),
    (16.5, 5.2,  8.3,1.2, '⑦ Premissas Críticas\nGOV.BR · MOSAICO · sem SWAG', AMBER, AMBER),
    (16.5, 6.9,  8.3,1.2, '⑧ Estimativa Consolidada\nPerfis · Horas · Riscos',  NAVY,  SF_BLUE),
]
for (nx,ny,nw,nh,nt,nc,lc) in nodes:
    rect(sl, nx,ny,nw,nh, fc=GRAY_BG, lc=lc, lw=1.5)
    txt(sl, nt, nx+0.12,ny+0.12,nw-0.2,nh-0.2, size=9,bold=True,color=nc)

# Linha de raciocínio (footer)
rect(sl, 0.3,9.1,24.8,2.85, fc=GRAY_BG, lc=GRAY_BOR)
txt(sl,'Linha de Raciocínio para Construção da Estimativa', 0.5,9.2,24,0.45, size=9,bold=True,color=NAVY)

steps = [
    ('1. Síntese de Escopo',   'Consolidar os 6 módulos e seus entregáveis confirmados'),
    ('2. Premissas',           'Documentar o que precisa ser verdade para a estimativa ser válida'),
    ('3. Riscos e Incertezas', 'Identificar o que NÃO sabemos e como isso impacta horas'),
    ('4. Perfis e Atividades', 'Mapear quais skills são necessários em cada módulo'),
    ('5. Estimativa em Faixas','Definir min/mid/max por módulo dado o nível de risco'),
    ('6. Ballpark ROM',        'Consolidar em uma proposta premissada e comunicável ao cliente'),
]
sx = 0.5
for (st,sd) in steps:
    rect(sl, sx,9.7,3.9,2.0, fc=WHITE, lc=SF_BLUE, lw=1)
    txt(sl, st, sx+0.1,9.8,3.7,0.45, size=8.5,bold=True,color=SF_BLUE)
    txt(sl, sd, sx+0.1,10.28,3.7,1.2, size=8,color=GRAY_TXT)
    sx += 4.1

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Agentforce
# ═════════════════════════════════════════════════════════════════════════════
sl = blank()
header(sl,'① Agentforce — 3 Agentes Especializados + Orquestrador',
       'Atendimento autônomo com autenticação GOV.BR e integração MOSAICO via MuleSoft')

agents = [
    ('Agente 1\nFAQ Geral',       'Atendimento padrão.\nBase de conhecimento\ndo legado.',  SF_BLUE),
    ('Agente 2\nRequisições',     'Requisições padrão.\nConsultas e status\nno legado.',    TEAL),
    ('Agente 3\nOutorga/MOSAICO','Especializado.\nRegistra pedido de\nconcessão e outorga\nno MOSAICO.\nAutentica GOV.BR.',  PURPLE),
    ('Agente 4\n(Orquestrador)',  'CONDICIONAL:\napenas se canal\nfor único.\nRoteamento entre\nos 3 agentes.',   AMBER),
]
ax = 0.3
for (at,ad,ac) in agents:
    rect(sl, ax,1.85,5.85,5.6, fc=GRAY_BG, lc=ac, lw=2)
    rect(sl, ax,1.85,5.85,0.8, fc=ac)
    txt(sl, at, ax,1.87,5.85,0.76, size=9.5,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(sl, ad, ax+0.15,2.75,5.55,4.5, size=9,color=GRAY_TXT)
    ax += 6.2

# MuleSoft bar
rect(sl, 0.3,7.7,24.8,0.7, fc=RGBColor(0xFF,0xF0,0xE0), lc=ORANGE, lw=1.5)
txt(sl,'MuleSoft  ·  4 plataformas legado integradas  ·  APIs/Métodos não documentados → RISCO ALTO',
    0.5,7.75,24,0.6, size=9,bold=True,color=ORANGE)

# Premissas & Riscos
bullet_block(sl,'Premissas',[
    'Catálogo de serviços e regras disponível (ausência eleva risco crítico)',
    'GOV.BR autenticação disponível (padrão Prodesp/Dataprev)',
    'Legados acessíveis via MuleSoft com APIs documentadas',
    'Reusable asset para 3 agentes base (reduz esforço de dev)',
], 0.3,8.55,12.0,4.5, tc=SF_BLUE)

bullet_block(sl,'Riscos',[
    '🔴 ALTO: Sem catálogo → toda regra de negócio no agente (erro arquitetural)',
    '🔴 ALTO: 4 plataformas sem visibilidade de APIs/métodos → estimativa frágil',
    '🟡 MED: Canal único não confirmado → Agente 4 pode ser necessário',
    '🟡 MED: Processo de outorga MOSAICO pouco documentado',
], 12.5,8.55,12.6,4.5, tc=RED, bg=RED_BG)

# Horas
rect(sl, 0.3,12.25,24.8,1.8, fc=GRAY_BG, lc=NAVY, lw=1.5)
txt(sl,'Perfis Necessários  ·  Agentforce', 0.5,12.3,24,0.4, size=8.5,bold=True,color=NAVY)
profiles = [
    ('Agentforce Specialist','400–600h'),
    ('MuleSoft Integration Dev','200–350h'),
    ('Solution Architect','120–180h'),
    ('QA Engineer','80–120h'),
    ('Project Manager','60–80h'),
]
px = 0.5
for (pr,ph) in profiles:
    txt(sl, pr, px,12.75,4.5,0.4, size=8,bold=True,color=NAVY)
    txt(sl, ph, px,13.15,4.5,0.4, size=9,bold=False,color=SF_BLUE)
    px += 4.9

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — CLM
# ═════════════════════════════════════════════════════════════════════════════
sl = blank()
header(sl,'② CLM — Gestão de Contratos OOTB',
       'Fluxo em 5 etapas · Benchmark setor público · Integração MOSAICO + MuleSoft')

# Fluxo 5 etapas
etapas = ['Fase 1\nSolicitação\ne Triagem','Fase 2\nAnálise\nJurídica','Fase 3\nNegociação\ne Aprovação','Fase 4\nExecução\ne Vigência','Fase 5\nRenovação e\nEncerramento']
colors_etapa = [SF_BLUE, TEAL, PURPLE, ORANGE, GREEN]
ex = 0.4
for i,(et,ec) in enumerate(zip(etapas,colors_etapa)):
    rect(sl, ex,1.85,4.4,2.2, fc=GRAY_BG, lc=ec, lw=2)
    rect(sl, ex,1.85,4.4,0.65, fc=ec)
    txt(sl, f'Etapa {i+1}', ex,1.87,4.4,0.62, size=8,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(sl, et, ex+0.1,2.58,4.2,1.3, size=9,color=GRAY_TXT,align=PP_ALIGN.CENTER)
    if i < 4:
        txt(sl, '→', ex+4.4,2.7,0.5,0.5, size=14,bold=True,color=GRAY_BOR,align=PP_ALIGN.CENTER)
    ex += 4.88

# Capacidades habilitadas
caps = [
    '5 templates dinâmicos com cláusulas variáveis',
    'Versionamento nativo Salesforce (sem redline/track changes)',
    '2 fluxos de aprovação multi-nível (gerente + jurídico + diretoria)',
    'Alertas automáticos D-90/D-60/D-30 para renovação',
    'Fluxo automático de renovação ou encerramento (Fase 5)',
    'Integração MOSAICO via MuleSoft (2 APIs)',
    'Retenção 7 anos via Salesforce Shield',
    'Relatórios básicos OOTB (status, vencimento, volume)',
]
bullet_block(sl,'Capacidades Confirmadas — Premissas ROM', caps, 0.4,4.3,12.2,5.5, tc=PURPLE)

bullet_block(sl,'Limite de Alcance (FORA do escopo)',[
    '✖ Assinatura digital ICP-Brasil ou DocuSign',
    '✖ Track changes / redline com contraparte',
    '✖ Negociação de cláusulas com participação externa',
    '✖ Mais de 2 fluxos ou 5 templates',
    '✖ Integração com SEI, SIAFI ou outros além de MOSAICO',
    '🟡 Risco residual: tipos de contrato ainda não confirmados',
], 12.8,4.3,12.2,5.5, tc=RED, bg=RED_BG)

rect(sl, 0.4,9.95,24.8,1.8, fc=GRAY_BG, lc=NAVY, lw=1.5)
txt(sl,'Perfis Necessários  ·  CLM  (revisado)', 0.6,10.0,24,0.4, size=8.5,bold=True,color=NAVY)
profiles = [
    ('CLM Consultant','320–440h'),
    ('MuleSoft Dev','80–120h'),
    ('Solution Architect','90–110h'),
    ('QA Engineer','85–110h'),
    ('PM','55–70h'),
]
px = 0.5
for (pr,ph) in profiles:
    txt(sl, pr, px,10.45,4.5,0.4, size=8,bold=True,color=NAVY)
    txt(sl, ph, px,10.85,4.5,0.4, size=9,color=PURPLE)
    px += 4.9

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Marketing Cloud
# ═════════════════════════════════════════════════════════════════════════════
sl = blank()
header(sl,'③ Marketing Cloud — 3 Jornadas Automatizadas',
       'Setup · Segmentação · Jornadas de Renovação, Adimplência e Onboarding')

journeys = [
    ('Jornada 1\nRenovação de Outorga',
     ['Segmento: Contratos por prazo\nde vencimento',
      'Trigger: D-90/D-60/D-30',
      'Canal: E-mail + WhatsApp',
      'Ação: Link de renovação'],
     SF_BLUE),
    ('Jornada 2\nAdimplência TFF/TFI',
     ['Segmento: TFF por vencimento',
      'Trigger: D-30/D-15/D-7/D+1/D+15',
      'Canal: E-mail + WhatsApp',
      'Ação: 2ª via de boleto'],
     ORANGE),
    ('Jornada 3\nOnboarding Novo Licenciado',
     ['Segmento: Novos contratos\npor emissão de outorga',
      'Trigger: Aprovação outorga',
      'Canal: E-mail',
      'Ação: Welcome + tutorial'],
     GREEN),
]
jx = 0.4
for (jt,jd,jc) in journeys:
    rect(sl, jx,1.85,7.8,7.5, fc=GRAY_BG, lc=jc, lw=2)
    rect(sl, jx,1.85,7.8,0.75, fc=jc)
    txt(sl, jt, jx,1.87,7.8,0.72, size=9.5,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    iy = 2.75
    for item in jd:
        txt(sl, f'▸  {item}', jx+0.2,iy,7.4,0.7, size=9,color=GRAY_TXT)
        iy += 0.72
    jx += 8.3

# Setup bar
rect(sl, 0.4,9.55,24.8,0.7, fc=RGBColor(0xE8,0xF4,0xFF), lc=SF_BLUE, lw=1.5)
txt(sl,'Setup incluído:  Conexão de canais (E-mail + WhatsApp)  ·  Configuração da conta MC  ·  Data Extension básica  ·  Integração com Core SF',
    0.6,9.6,24,0.6, size=8.5,color=SF_BLUE)

rect(sl, 0.4,10.4,24.8,1.8, fc=GRAY_BG, lc=NAVY, lw=1.5)
txt(sl,'Perfis Necessários  ·  Marketing Cloud', 0.6,10.45,24,0.4, size=8.5,bold=True,color=NAVY)
profiles = [
    ('MC Consultant','200–300h'),
    ('MC Developer','120–180h'),
    ('Solution Architect','60–80h'),
    ('QA','40–60h'),
    ('PM','30–50h'),
]
px = 0.5
for (pr,ph) in profiles:
    txt(sl, pr, px,10.9,4.5,0.4, size=8,bold=True,color=NAVY)
    txt(sl, ph, px,11.3,4.5,0.4, size=9,color=ORANGE)
    px += 4.9

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Tableau
# ═════════════════════════════════════════════════════════════════════════════
sl = blank()
header(sl,'④ Tableau — 5 Painéis de KPIs',
       'Ambiente Greenfield Cloud · Dados do Core SF ou Data Cloud · Fórmulas conhecidas')

dashboards = [
    ('Painel 1\nLicenciamento',   ['Status de outorgas','Tempo de aprovação','Volume por tipo de estação','Pendências MMAR']),
    ('Painel 2\nArrecadação TFF', ['Taxa de inadimplência','Arrecadação x meta','Volume por categoria','Aging de débitos']),
    ('Painel 3\nAtendimento',     ['Volume de chamados','SLA de resolução','Satisfação do cidadão','Top motivos']),
    ('Painel 4\nContratos',       ['Contratos vencendo','Status por etapa','Tempo médio de ciclo','Renovações x novos']),
    ('Painel 5\nRegulação',       ['Cobertura regulatória','Fiscalizações','Infrações por região','Conformidade %']),
]
dx = 0.3
for (dt,di) in dashboards:
    rect(sl, dx,1.85,4.7,7.8, fc=GRAY_BG, lc=TEAL, lw=1.5)
    rect(sl, dx,1.85,4.7,0.7, fc=TEAL)
    txt(sl, dt, dx,1.87,4.7,0.67, size=9,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(sl,'até 10 métricas', dx+0.1,2.65,4.5,0.35, size=7.5,italic=True,color=TEAL)
    iy = 3.1
    for item in di:
        txt(sl, f'▸  {item}', dx+0.15,iy,4.4,0.42, size=8.5,color=GRAY_TXT)
        iy += 0.47
    dx += 4.9

rect(sl, 0.3,9.85,24.8,0.65, fc=GRN_BG, lc=GREEN, lw=1.5)
txt(sl,'Premissas: Greenfield cloud (sem on-premise) · Fórmulas listadas e conhecidas · P&D apenas para refinamento · Dados disponíveis no Core SF / Data Cloud',
    0.5,9.9,24,0.56, size=8.5,color=GREEN)

rect(sl, 0.3,10.65,24.8,1.8, fc=GRAY_BG, lc=NAVY, lw=1.5)
txt(sl,'Perfis Necessários  ·  Tableau', 0.5,10.7,24,0.4, size=8.5,bold=True,color=NAVY)
profiles = [
    ('Tableau Developer','150–220h'),
    ('Data Analyst','100–150h'),
    ('Solution Architect','40–60h'),
    ('QA','30–50h'),
    ('PM','25–40h'),
]
px = 0.5
for (pr,ph) in profiles:
    txt(sl, pr, px,11.15,4.5,0.4, size=8,bold=True,color=NAVY)
    txt(sl, ph, px,11.55,4.5,0.4, size=9,color=TEAL)
    px += 4.9

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MuleSoft Atendimento
# ═════════════════════════════════════════════════════════════════════════════
sl = blank()
header(sl,'⑤ MuleSoft Atendimento — Integração Agentforce + Legados',
       'Setup Greenfield Cloud · 7 endpoints · 15–20 APIs estimadas · Sem SWAGGER')

# Architecture strip
rect(sl, 0.3,1.85,24.8,1.5, fc=RGBColor(0xF0,0xFF,0xF4), lc=GREEN, lw=1.5)
txt(sl,'Setup (Greenfield Cloud)  →  Autenticação (GOV.BR / SSO)  →  Conectores Legado  →  APIs Agentforce  →  Monitoramento Anypoint',
    0.5,2.2,24,0.8, size=9.5,bold=True,color=GREEN,align=PP_ALIGN.CENTER)

# Endpoints
endpoints = ['Agente 1\nFAQ Geral','Agente 2\nRequisições','Agente 3\nOutorga\nMOSAICO','Legado A','Legado B','Legado C','Legado D']
ex = 0.4
for i,ep in enumerate(endpoints):
    ec = GREEN if i < 3 else GRAY_TXT
    rect(sl, ex,3.6,3.3,1.6, fc=GRAY_BG, lc=ec, lw=1.5)
    txt(sl, ep, ex+0.1,3.7,3.1,1.4, size=8.5,color=ec,align=PP_ALIGN.CENTER)
    if i < 6:
        txt(sl,'|', ex+3.3,4.15,0.3,0.5, size=12,color=GRAY_BOR,align=PP_ALIGN.CENTER)
    ex += 3.55

bullet_block(sl,'Capacidades MuleSoft',[
    'Setup Greenfield Cloud (sem on-premise)',
    'Autenticação e segurança (OAuth2, GOV.BR)',
    'Transformação e normalização entre legados',
    'Error handling e retry automático',
    'Auditoria de chamadas',
    'Monitoramento via Anypoint Platform',
], 0.4,5.45,12.0,5.5, tc=GREEN)

bullet_block(sl,'Riscos',[
    '🔴 ALTO: 4 plataformas legado sem documentação de APIs',
    '🔴 ALTO: Sem SWAGGER → estimativa de 15–20 contratos como premissa',
    '🔴 ALTO: Volume real de APIs/métodos desconhecido',
    '🟡 MED: SIT/UAT disponíveis para as integrações',
    '🟡 MED: Contrato de assinatura conhecido e disponível',
], 12.6,5.45,12.5,5.5, tc=RED, bg=RED_BG)

rect(sl, 0.4,11.1,24.8,1.8, fc=GRAY_BG, lc=NAVY, lw=1.5)
txt(sl,'Perfis Necessários  ·  MuleSoft Atendimento', 0.6,11.15,24,0.4, size=8.5,bold=True,color=NAVY)
profiles = [('MuleSoft Integration Architect','120–180h'),('MuleSoft Developer','300–480h'),
            ('Solution Architect','80–100h'),('QA','80–120h'),('PM','40–60h')]
px = 0.5
for (pr,ph) in profiles:
    txt(sl, pr, px,11.6,4.5,0.4, size=8,bold=True,color=NAVY)
    txt(sl, ph, px,12.0,4.5,0.4, size=9,color=GREEN)
    px += 4.9

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — MuleSoft TFF
# ═════════════════════════════════════════════════════════════════════════════
sl = blank()
header(sl,'⑥ MuleSoft TFF — Consolidação de Arrecadação (~10M Registros)',
       'Processo Batch · 4 Fontes · 5 Etapas de Consolidação · Prazo regulatório 31/03/2027')

# Fontes → MuleSoft → Destino
sources = [('MOSAICO','Legacy',ORANGE),('SITARWEB','Legacy',ORANGE),
           ('DB_TELECOM','SQL Server',SF_BLUE),('SMS/FISTEL','MongoDB',GREEN)]
sx = 0.3
for (sn,st,sc) in sources:
    rect(sl, sx,1.85,5.5,1.5, fc=GRAY_BG, lc=sc, lw=1.5)
    txt(sl, sn, sx+0.1,2.0,5.3,0.55, size=10,bold=True,color=sc)
    txt(sl, st, sx+0.1,2.55,5.3,0.6, size=8.5,color=GRAY_TXT)
    sx += 5.7
txt(sl,'→', 23.0,2.3,1.5,0.8, size=18,bold=True,color=SF_BLUE,align=PP_ALIGN.CENTER)

# MuleSoft center
rect(sl, 0.3,3.6,24.8,3.2, fc=RGBColor(0xFF,0xF7,0xF0), lc=ORANGE, lw=2)
rect(sl, 0.3,3.6,24.8,0.55, fc=ORANGE)
txt(sl,'MuleSoft Anypoint  ·  Orquestrador do Processo Batch', 0.5,3.63,24,0.5, size=10,bold=True,color=WHITE)
steps_tff = ['1. Normalização\nde esquemas','2. Transformação\nde dados','3. Deduplicação\ne identidade',
             '4. Roteamento\ne validação','5. Entrega no\nSQL Arrecadação']
stx = 0.5
for i,st in enumerate(steps_tff):
    rect(sl, stx,4.3,4.55,2.3, fc=WHITE, lc=ORANGE, lw=1)
    txt(sl, st, stx+0.1,4.45,4.35,2.0, size=9,color=ORANGE,align=PP_ALIGN.CENTER)
    if i < 4:
        txt(sl,'→', stx+4.55,5.1,0.5,0.5, size=12,bold=True,color=ORANGE,align=PP_ALIGN.CENTER)
    stx += 5.05

# Destino
rect(sl, 0.3,7.05,24.8,0.6, fc=RGBColor(0xE6,0xF5,0xEC), lc=GREEN, lw=1.5)
txt(sl,'Destino: SQL Server de Arrecadação  ·  Dados consolidados prontos para motor de cálculo externo',
    0.5,7.1,24,0.5, size=8.5,color=GREEN)

bullet_block(sl,'Premissas',[
    'Nova instância MuleSoft Cloud (Greenfield)',
    'Motor de cálculo fiscal permanece externo (fora do escopo)',
    'Bases TFF acessíveis antes do kick-off',
    'Regras de cálculo documentadas antes do início',
    'Processo batch executado algumas vezes por mês',
], 0.3,7.85,12.0,4.8, tc=GREEN)

bullet_block(sl,'Riscos & Alertas',[
    '🔴 CRÍTICO: Prazo regulatório 31/03/2027 não negociável',
    '🔴 ALTO: ~10M registros com identidade fragmentada',
    '🟡 MED: Variabilidade nas regras TFF não mapeada',
    '🟡 MED: SIAFI, contestação e histórico fora do escopo',
    '⚡ MOSAICO adicionado como 4ª fonte (diferente do slide TFF anterior)',
], 12.5,7.85,12.6,4.8, tc=RED, bg=RED_BG)

rect(sl, 0.3,12.82,24.8,1.35, fc=GRAY_BG, lc=NAVY, lw=1.5)
txt(sl,'Perfis  ·  MuleSoft TFF', 0.5,12.87,24,0.38, size=8.5,bold=True,color=NAVY)
profiles = [('MuleSoft Integration Architect','120–160h'),('MuleSoft Developer','280–400h'),
            ('Data Engineer','120–180h'),('QA','80–100h'),('PM','40–60h')]
px = 0.5
for (pr,ph) in profiles:
    txt(sl, pr, px,13.28,4.5,0.35, size=7.5,bold=True,color=NAVY)
    txt(sl, ph, px,13.63,4.5,0.35, size=8.5,color=RED)
    px += 4.9

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Estimativa Consolidada
# ═════════════════════════════════════════════════════════════════════════════
sl = blank()
header(sl,'Estimativa Consolidada — Horas e Perfis por Módulo',
       'ROM Ballpark · Faixas Min–Mid–Max · Premissado e revalidar com delivery')

# Table
cols_w = [5.5,2.1,2.1,2.1,4.0,8.6]
cols_h = ['Módulo','Mín (h)','Mid (h)','Máx (h)','Risco','Perfis Chave']
rows = [
    ['① Agentforce',    '860','1.080','1.330','🔴 ALTO','Agentforce Spec · MuleSoft Dev · SA'],
    ['② CLM Contratos', '650', '780',  '910','🟡 MED','CLM Consultant · MuleSoft Dev · SA'],
    ['③ Marketing Cloud','450', '570',  '670','🟡 MED', 'MC Consultant · MC Developer · SA'],
    ['④ Tableau',        '345', '460',  '520','🟢 BAIXO','Tableau Dev · Data Analyst · SA'],
    ['⑤ MuleSoft Atend.','620', '830', '1.040','🔴 ALTO','Mule Arch · Mule Dev · QA · SA'],
    ['⑥ MuleSoft TFF',  '640', '800',  '900','🔴 CRÍTICO','Mule Arch · Data Eng · Mule Dev · SA'],
    ['PM & Arquitetura\n(transversal)','240','300','360','—','PM · Solution Architect'],
]
totals = ['TOTAL','3.805','4.820','5.730','—','—']

# header row
ty = 1.65
tx = 0.3
for i,(cw,ch) in enumerate(zip(cols_w,cols_h)):
    rect(sl, tx,ty,cw,0.55, fc=NAVY)
    txt(sl, ch, tx+0.05,ty+0.02,cw-0.1,0.5, size=8.5,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    tx += cw

# data rows
risk_colors = {'🔴 ALTO':RED,'🔴 CRÍTICO':RED,'🟡 MED':AMBER,'🟢 BAIXO':GREEN,'—':GRAY_TXT}
ty += 0.55
for ri,row in enumerate(rows):
    fc = GRAY_BG if ri%2==0 else WHITE
    tx = 0.3
    rh = 0.78
    for i,(cw,cell) in enumerate(zip(cols_w,row)):
        rect(sl, tx,ty,cw,rh, fc=fc, lc=GRAY_BOR, lw=0.5)
        rc = risk_colors.get(cell, DARK) if i==4 else DARK
        bold = (i==0)
        fs = 8 if i==5 else 9
        txt(sl, cell, tx+0.08,ty+0.04,cw-0.12,rh-0.1, size=fs,bold=bold,color=rc,align=PP_ALIGN.CENTER)
        tx += cw
    ty += rh

# totals
tx = 0.3
for i,(cw,cell) in enumerate(zip(cols_w,totals)):
    rect(sl, tx,ty,cw,0.72, fc=NAVY)
    txt(sl, cell, tx+0.05,ty+0.04,cw-0.1,0.64, size=10,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    tx += cw

# Footer note
rect(sl, 0.3,ty+0.8,24.8,1.15, fc=AMB_BG, lc=AMBER, lw=1.5)
txt(sl,'⚠  Estimativa em nível ROM Ballpark — faixas amplas refletem riscos de discovery. Recomendado revalidar com Delivery antes da entrega ao cliente.',
    0.5,ty+0.88,24,0.5, size=8.5,bold=True,color=AMBER)
txt(sl,'Premissa: reusable assets para Agentforce reduzem risco. Módulos MuleSoft com maior incerteza por falta de documentação de APIs.',
    0.5,ty+1.38,24,0.42, size=8,color=RGBColor(0x7A,0x4F,0x00))

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Premissas & Próximos Passos
# ═════════════════════════════════════════════════════════════════════════════
sl = blank()
header(sl,'Premissas Críticas & Próximos Passos',
       'Condições para validade do ROM · Ações para refinar a proposta')

bullet_block(sl,'Premissas que tornam o ROM válido',[
    'Reusable assets disponíveis para os 3 agentes Agentforce',
    'GOV.BR autenticação padrão Prodesp/Dataprev replicável',
    'MuleSoft Cloud Greenfield (sem on-premise) — ambos módulos',
    'APIs legado disponíveis em SIT/UAT com contratos conhecidos',
    'Motor de cálculo TFF permanece externo (não estimado)',
    'Fórmulas Tableau conhecidas — P&D apenas para refinamento',
    'CLM Greenfield · 5 templates dinâmicos · 2 fluxos aprovação multi-nível',
    'CLM sem redline/track changes · versionamento nativo apenas',
    'Assinatura digital ICP-Brasil fora do escopo CLM',
    'Retenção 7 anos via Shield (CLM Fase 5)',
    'PO ANATEL dedicado com poder de decisão',
], 0.4,1.85,12.0,7.8, tc=GREEN, bg=GRN_BG)

bullet_block(sl,'O que aumentaria o escopo (não precificado)',[
    'Catálogo de serviços ausente → regras no agente (risco arquitetural)',
    'Mais de 20 contratos MuleSoft descobertos no discovery',
    'Integração Marinha do Brasil / DECEA / ANAC',
    'Assinatura digital ICP-Brasil ou DocuSign no CLM',
    'Transbordo humano (Service Console) no Agentforce',
    'Histórico de cálculos TFF anteriores a migrar',
    'Data Cloud para resolução de identidade TFF',
    'Canal único confirmado → Agente Orquestrador (4º agente)',
], 12.6,1.85,12.5,7.8, tc=RED, bg=RED_BG)

# Próximos passos
rect(sl, 0.4,9.85,24.8,0.5, fc=NAVY)
txt(sl,'Próximos Passos', 0.6,9.88,24,0.44, size=10,bold=True,color=WHITE)

steps_next = [
    ('Imediato','Validar ROM com Rogerio, Rodolpho e delivery antes de enviar à DTP'),
    ('Curto prazo','Alinhar precisão necessária com Mari/Fernanda para proposta comercial'),
    ('Discovery','Solicitar SWAGGER e documentação de APIs dos 4 legados MuleSoft'),
    ('Catálogo','Confirmar existência de catálogo de serviços para Agentforce'),
    ('Canal único','Confirmar se canal de atendimento é único → define necessidade do 4º agente'),
]
sx = 0.4
sy = 10.5
for (st,sd) in steps_next:
    rect(sl, sx,sy,4.75,2.55, fc=GRAY_BG, lc=SF_BLUE, lw=1.5)
    rect(sl, sx,sy,4.75,0.5, fc=SF_BLUE)
    txt(sl, st, sx,sy+0.02,4.75,0.46, size=8.5,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    txt(sl, sd, sx+0.12,sy+0.62,4.5,1.8, size=8.5,color=GRAY_TXT)
    sx += 4.98

out = '/Users/nfilho/claude/ANATEL_MapaMental_Estimativa.pptx'
prs.save(out)
print(f'Salvo: {out}')
