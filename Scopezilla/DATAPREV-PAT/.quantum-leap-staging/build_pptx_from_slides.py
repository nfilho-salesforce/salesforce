#!/usr/bin/env python3
"""Assemble a 16:9 PPTX from the per-slide PNG renders of presentation-deck.html.

Each slide was rendered at 1280x720 (the deck's native design resolution) via the
Playwright MCP browser, so dropping each PNG full-bleed onto a 13.333in x 7.5in
slide preserves the Salesforce/Scopezilla visual identity pixel-for-pixel — no
font substitution, no CSS re-implementation, no fidelity loss.

Usage: build_pptx_from_slides.py <png-dir> <out.pptx>
  <png-dir> must contain slide-01.png .. slide-NN.png (zero-padded, in order).
"""
import sys, pathlib, re
from pptx import Presentation
from pptx.util import Inches

def main():
    png_dir = pathlib.Path(sys.argv[1])
    out = pathlib.Path(sys.argv[2])

    pngs = sorted(
        png_dir.glob("slide-*.png"),
        key=lambda p: int(re.search(r"slide-(\d+)", p.name).group(1)),
    )
    if not pngs:
        print(f"ERROR: no slide-*.png found in {png_dir}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # fully blank layout

    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(png), 0, 0, width=prs.slide_width, height=prs.slide_height
        )

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"wrote {out} — {len(pngs)} slides from {png_dir}")

if __name__ == "__main__":
    main()
