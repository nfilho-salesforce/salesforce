# -*- coding: utf-8 -*-
"""DATA ÁGIL — ROM v5.0 · 15 jornadas · R$5M · Ago–Dez 2026 (paleta gov.br)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ─── Paleta gov.br (igual à página HTML) ───────────────────────────
BLUE_DARK  = RGBColor(0x07, 0x1D, 0x41)
BLUE_MID   = RGBColor(0x0C, 0x32, 0x6F)
BLUE_ACC   = RGBColor(0x13, 0x51, 0xB4)
BLUE_LIGHT = RGBColor(0xED, 0xF5, 0xFF)
TEAL       = RGBColor(0x37, 0x7E, 0xC1)
GREEN      = RGBColor(0x16, 0x88, 0x21)
GREEN_D    = RGBColor(0x16, 0x88, 0x21)
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)
YELLOW     = RGBColor(0xF7, 0xC9, 0x48)
GOV_YELLOW = RGBColor(0xFB, 0xBA, 0x00)
PURPLE     = RGBColor(0x7B, 0x2F, 0xBE)
SLACK      = RGBColor(0x4A, 0x15, 0x4B)
SEI_GOLD   = RGBColor(0xC7, 0x77, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x63, 0x63, 0x63)
GRAY_BG    = RGBColor(0xF8, 0xF8, 0xF8)
GRAY_BORD  = RGBColor(0xE6, 0xE6, 0xE6)
TEXT       = RGBColor(0x26, 0x26, 0x26)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

ONDA_COLOR = {1: GREEN_D, 2: BLUE_ACC, 3: PURPLE}

# ─── Helpers ────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None, radius=False):
    shp = slide.shapes.add_shape(5 if radius else 1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
        shp.line.width = line_w or Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def _set(run, size, bold, color, italic=False):
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    run.font.name = "Calibri"

def add_text(slide, text, l, t, w, h, size=14, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, italic=False, anchor=None):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame; tf.word_wrap = True
    if anchor: tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = align
    _set(p.add_run(), size, bold, color, italic) if False else None
    r = p.add_run(); r.text = text; _set(r, size, bold, color, italic)
    return txb

def add_bullets(slide, items, l, t, w, h, size=12, color=TEXT, gap=4):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        bold = False; txt = it
        if isinstance(it, tuple):
            txt, bold = it
        r = p.add_run(); r.text = "•  " + txt; _set(r, size, bold, color)
    return txb

def header(slide, kicker, title, accent=BLUE_ACC):
    add_rect(slide, 0, 0, 13.33, 1.15, fill=BLUE_DARK)
    add_rect(slide, 0, 1.15, 13.33, 0.06, fill=accent)
    add_text(slide, kicker.upper(), 0.55, 0.18, 12, 0.3, size=11, bold=True, color=YELLOW)
    add_text(slide, title, 0.55, 0.46, 12.2, 0.6, size=23, bold=True, color=WHITE)

def footer(slide, n):
    add_text(slide, "DATAPREV · DATA ÁGIL — ROM v5.0 · Salesforce Professional Services",
             0.55, 7.12, 10, 0.3, size=8.5, color=GRAY)
    add_text(slide, f"{n}", 12.5, 7.12, 0.6, 0.3, size=8.5, color=GRAY, align=PP_ALIGN.RIGHT)

def kpi_card(slide, l, t, w, big, small, color, h=1.25):
    add_rect(slide, l, t, w, h, fill=WHITE, line=GRAY_BORD)
    add_rect(slide, l, t, 0.08, h, fill=color)
    add_text(slide, big, l+0.15, t+0.18, w-0.25, 0.55, size=24, bold=True, color=color)
    add_text(slide, small, l+0.15, t+0.75, w-0.25, 0.45, size=10.5, color=GRAY)

def new():
    return prs.slides.add_slide(BLANK)

def cell(table, r, c, text, size=10, bold=False, color=TEXT, fill=None, align=PP_ALIGN.LEFT):
    cl = table.cell(r, c)
    cl.margin_left = Pt(5); cl.margin_right = Pt(5); cl.margin_top = Pt(2); cl.margin_bottom = Pt(2)
    cl.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill is not None:
        cl.fill.solid(); cl.fill.fore_color.rgb = fill
    tf = cl.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r0 = p.add_run(); r0.text = text; _set(r0, size, bold, color)

# ─── Dados: 15 jornadas ─────────────────────────────────────────────
J = [
    ("J1","Consulta Financeira","Externo+Interno · Protheus",1,"Leitura","R$ 263.158","Posição financeira do órgão via Slack",402,3),
    ("J2","Status de Chamado","Externo+Interno · Service Desk",1,"Leitura","R$ 263.158","Status/SLA + notificação de mudança",402,3),
    ("SEI-J1","Alerta de Prazo/tácito","Interno · SEI",1,"Leitura","R$ 350.877","Alerta antes do decurso de prazo",537,4),
    ("SEI-J2","Consulta linguagem natural","Interno · SEI",1,"Leitura","R$ 350.877","Achar processo por pergunta natural",537,4),
    ("SEI-J3","Notif. recebido/tramitado","Interno · SEI",1,"Leitura","R$ 263.158","Aviso quando processo chega à unidade",402,3),
    ("J3","Briefing de Projeto","Interno · Clarity",2,"Leitura","R$ 263.158","Briefing executivo em tempo real",402,3),
    ("J7","FAQ Interno via Conexão","Interno · CRM/RAG",2,"Leitura","R$ 263.158","Dúvidas de RH respondidas com fonte",402,3),
    ("SEI-J4","Meu painel de processos","Interno · SEI",2,"Leitura","R$ 263.158","Visão consolidada dos processos",402,3),
    ("SEI-J5","Digest de unidade","Interno · SEI",2,"Leitura","R$ 263.158","Resumo de fluxo/gargalos da unidade",402,3),
    ("SEI-J6","'Qual tipo uso?' + RAG","Interno · SEI",2,"Leitura","R$ 350.877","Recomenda tipo com norma citada",537,4),
    ("SEI-J7","Status assinatura + deep-link","Interno · SEI",2,"Leitura","R$ 263.158","Quem falta assinar + link p/ assinar",402,3),
    ("J4","Agendamento por Voz","Interno · MS Graph",3,"Escrita","R$ 438.596","Agenda reunião por áudio no Slack",671,5),
    ("SEI-J8","Ciência de documento","Interno · SEI",3,"Escrita leve","R$ 350.877","Registrar ciência + trilha",537,4),
    ("SEI-J9","Tramitar via aprovação","Interno · SEI",3,"Escrita+Gov","R$ 526.316","Tramitar com fluxo de aprovação",805,6),
    ("SEI-J10","Abrir processo","Interno · SEI",3,"Escrita+Gov","R$ 526.316","Abertura assistida c/ validação de tipo",805,6),
]
N = 1

# ═══ SLIDE 1 · CAPA ═══
s = new()
add_rect(s, 0, 0, 13.33, 7.5, fill=BLUE_DARK)
add_rect(s, 0, 5.05, 13.33, 0.09, fill=YELLOW)
add_text(s, "DATAPREV", 0.7, 1.5, 12, 0.5, size=15, bold=True, color=YELLOW)
add_text(s, "DATA ÁGIL", 0.7, 2.0, 12, 1.0, size=52, bold=True, color=WHITE)
add_text(s, "Autosserviço e automação via Slack, Agentforce e MuleSoft", 0.7, 3.15, 12, 0.5, size=18, color=RGBColor(0xBF,0xD8,0xF0))
add_text(s, "ROM v5.0  ·  15 jornadas (5 agentes + 10 SEI)  ·  Teto fixo R$ 5,0 M  ·  Ago–Dez 2026",
         0.7, 3.75, 12, 0.5, size=14, bold=True, color=WHITE)
add_text(s, "Salesforce Professional Services LATAM  ·  29 Jul 2026  ·  Faixa indicativa para planejamento",
         0.7, 5.3, 12, 0.5, size=11, color=GRAY)
N += 1

# ═══ SLIDE 2 · O PROBLEMA ═══
s = new(); header(s, "Diagnóstico", "O problema não é de processo — é de escala", ORANGE)
kpi_card(s, 0.55, 1.5, 2.9, "166×", "crescimento da base de clientes", ORANGE)
kpi_card(s, 3.65, 1.5, 2.9, "~30k", "chamados/mês (13–14k só INSS)", ORANGE)
kpi_card(s, 6.75, 1.5, 2.9, "80/20", "poucos temas concentram a demanda", ORANGE)
kpi_card(s, 9.85, 1.5, 2.9, "10/15", "jornadas nascem no SEI", SEI_GOLD)
add_rect(s, 0.55, 3.05, 12.23, 0.04, fill=GRAY_BORD)
add_text(s, "A informação existe — o acesso é reativo e manual", 0.55, 3.2, 12, 0.4, size=15, bold=True, color=BLUE_DARK)
add_bullets(s, [
    ("Financeiro / Protheus: clientes perguntam a dívida todo dia — hoje via ofício ou reunião.", False),
    ("Executivos convocados sem tempo de preparar ponto de controle (Clarity).", False),
    ("RH / Conexão: a resposta está publicada, mas o empregado não a alcança.", False),
    ("SEI: ninguém é avisado de prazo tácito, chegada de processo ou assinatura pendente.", False),
    ("Sem push no SEI → polling obrigatório; conteúdo sensível → LGPD (só metadados no Slack).", True),
], 0.55, 3.65, 12.2, 3.0, size=13, gap=8)
footer(s, N); N += 1

# ═══ SLIDE 3 · POR QUE NÃO WHATSAPP ═══
s = new(); header(s, "Diagnóstico", "Por que não WhatsApp", YELLOW)
add_bullets(s, [
    ("Custo por mensagem da Meta escala com o volume — imprevisível para governo.", True),
    "Governança e LGPD: dados sensíveis em plataforma de terceiros fora do controle da Dataprev.",
    "Sem integração nativa com Agentforce/MuleSoft nem com o ecossistema interno.",
    "Slack entrega o mesmo canal conversacional com governança, dual-workspace e trilha.",
    "Colaboradores e gestores já vivem no Slack — adoção com menos atrito.",
], 0.55, 1.5, 12.2, 4.5, size=14, gap=12)
footer(s, N); N += 1

# ═══ SLIDE 4 · ARQUITETURA ═══
s = new(); header(s, "Solução", "Arquitetura — 3 pilares + SEI", TEAL)
def pillar(l, color, icon, title, desc):
    add_rect(s, l, 1.6, 3.95, 3.3, fill=WHITE, line=GRAY_BORD)
    add_rect(s, l, 1.6, 3.95, 0.6, fill=color)
    add_text(s, icon+"  "+title, l+0.15, 1.72, 3.7, 0.4, size=13, bold=True, color=WHITE)
    add_text(s, desc, l+0.2, 2.4, 3.6, 2.3, size=11.5, color=TEXT)
pillar(0.55, SLACK, "💬", "Pilar 3 · Slack + bot",
       "Front conversacional. Dois workspaces: interno (colaboradores + gestores) e externo (clientes), separados pelo compartilhamento de canais públicos.")
pillar(4.7, BLUE_ACC, "🧠", "Pilar 2 · Agentforce",
       "Agentes especialistas — um por sistema de origem: Financeiro, Suporte, Projetos, Interno/RH e Agente SEI. Tópicos, actions e RAG quando aplicável.")
pillar(8.85, TEAL, "🔌", "Pilar 1 · MuleSoft",
       "Integração + exposição de MCP server dos legados. Para o SEI: agendador de polling + cache de estado (não há push/webhook).")
add_rect(s, 0.55, 5.15, 12.23, 1.4, fill=BLUE_LIGHT, line=BLUE_ACC)
add_text(s, "SEI no fluxo", 0.75, 5.28, 4, 0.35, size=12, bold=True, color=BLUE_DARK)
add_text(s, "Leitura por polling agendado (alertas, consultas, painéis). Escrita (ciência, tramitação, abertura) só após a Fase 0 / G1002, "
            "com confirmação humana e trilha. Conteúdo sensível nunca no Slack — só metadados + deep-link (LGPD).",
         0.75, 5.62, 11.8, 0.9, size=11.5, color=TEXT)
footer(s, N); N += 1

# ═══ SLIDE 5 · DOIS PÚBLICOS ═══
s = new(); header(s, "Solução", "Dois públicos · ambiente Slack separado", TEAL)
add_rect(s, 0.55, 1.6, 6.0, 3.3, fill=WHITE, line=GRAY_BORD); add_rect(s, 0.55, 1.6, 0.09, 3.3, fill=BLUE_ACC)
add_text(s, "👥 Interno · workspace Dataprev", 0.8, 1.75, 5.6, 0.4, size=14, bold=True, color=BLUE_DARK)
add_bullets(s, ["Acesso às 15 jornadas, inclusive as 10 do SEI.",
                "Vínculo usuário Dataprev ↔ SEI (painel e alertas).",
                "Jornadas internas: Projetos, FAQ/RH, Agendamento por voz."],
            0.8, 2.25, 5.5, 2.5, size=12, gap=8)
add_rect(s, 6.75, 1.6, 6.0, 3.3, fill=WHITE, line=GRAY_BORD); add_rect(s, 6.75, 1.6, 0.09, 3.3, fill=GREEN_D)
add_text(s, "🏛️ Externo · workspace de clientes", 7.0, 1.75, 5.6, 0.4, size=14, bold=True, color=BLUE_DARK)
add_bullets(s, ["Só leitura do próprio órgão: Consulta Financeira (J1) e Status de Chamado (J2).",
                "Sem acesso a jornadas SEI internas.",
                "Autenticação por perfil; dados restritos por ID de contrato/órgão."],
            7.0, 2.25, 5.5, 2.5, size=12, gap=8)
add_rect(s, 0.55, 5.15, 12.23, 1.3, fill=RGBColor(0xF3,0xE8,0xF5), line=SLACK)
add_text(s, "💬 Dois workspaces distintos — canais públicos são compartilhados dentro do workspace; a separação isola dados de cada cliente e o ambiente interno da Dataprev.",
         0.75, 5.45, 11.8, 0.8, size=12, bold=True, color=SLACK)
footer(s, N); N += 1

# ═══ SLIDE 6 · VISÃO DAS 15 JORNADAS ═══
s = new(); header(s, "Escopo", "As 15 jornadas de escopo", GREEN_D)
rows = 16; cols = 5
tb = s.shapes.add_table(rows, cols, Inches(0.55), Inches(1.45), Inches(12.23), Inches(5.4)).table
tb.columns[0].width = Inches(1.4); tb.columns[1].width = Inches(4.2)
tb.columns[2].width = Inches(3.5); tb.columns[3].width = Inches(1.8); tb.columns[4].width = Inches(1.33)
hd = ["ID","Jornada","Escopo","Tipo","R$ c/imp"]
for c, h in enumerate(hd):
    cell(tb, 0, c, h, size=10.5, bold=True, color=WHITE, fill=BLUE_MID)
for i, j in enumerate(J):
    idc, title, sub, onda, rw, val, esc, hrs, peso = j
    fill = GRAY_BG if i % 2 else WHITE
    cell(tb, i+1, 0, idc, size=9.5, bold=True, color=ONDA_COLOR[onda], fill=fill)
    cell(tb, i+1, 1, title, size=9.5, fill=fill)
    cell(tb, i+1, 2, esc, size=9, color=GRAY, fill=fill)
    cell(tb, i+1, 3, rw, size=9, fill=fill)
    cell(tb, i+1, 4, val, size=9.5, bold=True, color=BLUE_DARK, fill=fill, align=PP_ALIGN.RIGHT)
footer(s, N); N += 1

# ═══ SLIDES 7-9 · ONDAS ═══
def onda_slide(onda, titulo, janela, accent):
    global N
    s = new(); header(s, "Escopo · Ondas", titulo, accent)
    js = [j for j in J if j[3] == onda]
    add_text(s, janela, 0.55, 1.35, 12, 0.35, size=12, bold=True, color=accent)
    y = 1.85; cols = 2; cw = 6.05; ch = 1.55; gapx = 0.13; gapy = 0.15
    for i, j in enumerate(js):
        idc, title, sub, o, rw, val, esc, hrs, peso = j
        col = i % cols; row = i // cols
        l = 0.55 + col*(cw+gapx); t = y + row*(ch+gapy)
        add_rect(s, l, t, cw, ch, fill=WHITE, line=GRAY_BORD)
        add_rect(s, l, t, 0.08, ch, fill=accent)
        add_text(s, f"{idc} · {title}", l+0.2, t+0.12, cw-0.35, 0.4, size=12, bold=True, color=BLUE_DARK)
        add_text(s, sub, l+0.2, t+0.5, cw-2.0, 0.3, size=9.5, color=GRAY)
        add_text(s, esc, l+0.2, t+0.82, cw-2.0, 0.6, size=10, color=TEXT)
        add_text(s, val, l+cw-1.75, t+0.12, 1.6, 0.35, size=11, bold=True, color=accent, align=PP_ALIGN.RIGHT)
        add_text(s, rw, l+cw-1.75, t+0.5, 1.6, 0.3, size=9, color=GRAY, align=PP_ALIGN.RIGHT)
    footer(s, N); N += 1

onda_slide(1, "Onda 1 · Quick Wins (5)", "Ago–Set · go-live fim de Set · leitura, alto valor / baixo esforço", GREEN_D)
onda_slide(2, "Onda 2 · Expansão (6)", "Out–Nov · leitura + RAG + painéis", BLUE_ACC)
onda_slide(3, "Onda 3 · Escrita, voz e governança (4)", "Nov–Dez + hypercare · após a Fase 0 (G1002)", PURPLE)

# ═══ SLIDE 10 · SISTEMAS & VOLUMETRIA ═══
s = new(); header(s, "Solução", "Sistemas & Volumetria", TEAL)
sist = [("💰 Protheus","Financeiro — posição/faturas (leitura por ID)"),
        ("🎫 Service Desk","Chamados — status/SLA (~30k/mês)"),
        ("📊 Clarity","PPM — status e marcos de projeto"),
        ("🏠 Conexão / MS Graph","RH/FAQ (RAG) e Agenda/Teams (escrita)"),
        ("📑 SEI · mod-wssei v2","REST JWT ~150+ endpoints · polling · LGPD")]
for i,(t,d) in enumerate(sist):
    y = 1.55 + i*0.72
    color = SEI_GOLD if i==4 else BLUE_ACC
    add_rect(s, 0.55, y, 12.23, 0.62, fill=(RGBColor(0xFF,0xF6,0xE0) if i==4 else WHITE), line=GRAY_BORD)
    add_text(s, t, 0.75, y+0.1, 4.0, 0.4, size=12.5, bold=True, color=color)
    add_text(s, d, 4.8, y+0.13, 7.7, 0.4, size=11.5, color=TEXT)
add_rect(s, 0.55, 5.35, 12.23, 1.15, fill=RGBColor(0xFF,0xF6,0xE0), line=SEI_GOLD)
add_text(s, "Volumetria pendente (gap G1102)", 0.75, 5.45, 6, 0.35, size=12, bold=True, color=SEI_GOLD)
add_text(s, "Volume de processos/tramitações por unidade será fornecido na Fase 0 para dimensionar a frequência de polling (latência × custo) e calibrar os KPIs.",
         0.75, 5.8, 11.8, 0.6, size=11.5, color=TEXT)
footer(s, N); N += 1

# ═══ SLIDE 11 · ESTIMATIVA & BOLSÃO ═══
s = new(); header(s, "Estimativa & ROM", "Bolsão fixo de R$ 5,0 milhões", BLUE_ACC)
kpi_card(s, 0.55, 1.55, 3.9, "R$ 5,0 M", "teto com imposto (fixo)", GREEN_D)
kpi_card(s, 4.65, 1.55, 3.9, "R$ 4,67 M", "sem imposto (÷ 0,9345)", BLUE_ACC)
kpi_card(s, 8.75, 1.55, 4.0, "~7.650 h", "blended ~R$ 653/h c/imp", BLUE_DARK)
add_rect(s, 0.55, 3.1, 12.23, 3.35, fill=BLUE_LIGHT, line=BLUE_ACC)
add_text(s, "O bolsão por dentro", 0.75, 3.25, 8, 0.4, size=15, bold=True, color=BLUE_DARK)
add_bullets(s, [
    ("Teto fixo: o escopo cabe por priorização de ondas — não por corte de qualidade.", True),
    "~45% de cada valor é fundação compartilhada (MuleSoft core + MCP, setup Slack dual-workspace, framework de agentes, Fase 0, PM, QA, Change, UX baseline).",
    "Valor por jornada visível — pesos por complexidade: leitura=3, RAG/polling=4, voz+escrita Graph=5, escrita SEI + governança=6.",
    "Ganho de IA (≥25%) é o que torna 22 semanas e o teto de R$ 5 M viáveis para as 15 jornadas.",
], 0.75, 3.75, 11.8, 2.6, size=12.5, gap=10)
footer(s, N); N += 1

# ═══ SLIDE 12 · CARDÁPIO R$5M ═══
s = new(); header(s, "Estimativa & ROM", "Cardápio — R$ 5 M jornada por jornada", BLUE_ACC)
rows = 17
tb = s.shapes.add_table(rows, 5, Inches(0.55), Inches(1.4), Inches(12.23), Inches(5.5)).table
tb.columns[0].width = Inches(4.6); tb.columns[1].width = Inches(1.5)
tb.columns[2].width = Inches(2.0); tb.columns[3].width = Inches(2.4); tb.columns[4].width = Inches(1.73)
for c,h in enumerate(["Jornada","Onda","Horas","R$ c/imp","Peso"]):
    cell(tb, 0, c, h, size=10.5, bold=True, color=WHITE, fill=BLUE_MID)
for i, j in enumerate(J):
    idc, title, sub, onda, rw, val, esc, hrs, peso = j
    fill = GRAY_BG if i % 2 else WHITE
    cell(tb, i+1, 0, f"{idc} · {title}", size=10, fill=fill)
    cell(tb, i+1, 1, f"Onda {onda}", size=9.5, color=ONDA_COLOR[onda], bold=True, fill=fill, align=PP_ALIGN.CENTER)
    cell(tb, i+1, 2, f"{hrs}h", size=10, fill=fill, align=PP_ALIGN.RIGHT)
    cell(tb, i+1, 3, val, size=10, bold=True, color=BLUE_DARK, fill=fill, align=PP_ALIGN.RIGHT)
    cell(tb, i+1, 4, str(peso), size=10, fill=fill, align=PP_ALIGN.CENTER)
cell(tb, 16, 0, "TOTAL · 15 jornadas", size=10.5, bold=True, color=WHITE, fill=BLUE_DARK)
cell(tb, 16, 1, "", fill=BLUE_DARK)
cell(tb, 16, 2, "~7.650h", size=10.5, bold=True, color=WHITE, fill=BLUE_DARK, align=PP_ALIGN.RIGHT)
cell(tb, 16, 3, "R$ 5.000.000", size=10.5, bold=True, color=YELLOW, fill=BLUE_DARK, align=PP_ALIGN.RIGHT)
cell(tb, 16, 4, "57", size=10.5, bold=True, color=WHITE, fill=BLUE_DARK, align=PP_ALIGN.CENTER)
footer(s, N); N += 1

# ═══ SLIDE 13 · ONDAS TIMELINE ═══
s = new(); header(s, "Estimativa & ROM", "Ondas de entrega — Ago a Dez 2026", TEAL)
add_rect(s, 0.55, 1.5, 12.23, 0.9, fill=RGBColor(0xFD,0xF0,0xE0), line=ORANGE)
add_text(s, "🚧 Fase 0 (Agosto · obrigatória)", 0.75, 1.6, 6, 0.35, size=12, bold=True, color=ORANGE)
add_text(s, "Destrava o bloqueador de governança G1002 (perfis Protheus/SEI, workspace externo, LGPD Art. 48 / TCU). Gate da escrita SEI.",
         0.75, 1.95, 11.8, 0.4, size=11, color=TEXT)
ondas = [(GREEN_D,"Onda 1 · Quick Wins","Ago–Set · go-live fim Set","J1, J2, SEI-J1, SEI-J2, SEI-J3 (5) — leitura, alto valor/baixo esforço"),
         (BLUE_ACC,"Onda 2 · Expansão","Out–Nov","J3, J7, SEI-J4, SEI-J5, SEI-J6, SEI-J7 (6) — leitura + RAG + painéis"),
         (PURPLE,"Onda 3 · Escrita/Transação","Nov–Dez + hypercare","J4, SEI-J8, SEI-J9, SEI-J10 (4) — escrita, voz, governança")]
for i,(c,t,jan,desc) in enumerate(ondas):
    y = 2.65 + i*1.25
    add_rect(s, 0.55, y, 12.23, 1.1, fill=WHITE, line=GRAY_BORD)
    add_rect(s, 0.55, y, 0.1, 1.1, fill=c)
    add_text(s, t, 0.8, y+0.12, 5, 0.4, size=13, bold=True, color=c)
    add_text(s, jan, 8.8, y+0.14, 3.8, 0.35, size=11, bold=True, color=GRAY, align=PP_ALIGN.RIGHT)
    add_text(s, desc, 0.8, y+0.55, 11.7, 0.45, size=11.5, color=TEXT)
add_text(s, "Prioridade = valor ÷ esforço. Leitura primeiro; escrita e governança por último, após a Fase 0.",
         0.55, 6.55, 12, 0.35, size=10.5, italic=True, color=GRAY)
footer(s, N); N += 1

# ═══ SLIDE 13b · CRONOGRAMA GRÁFICO POR JORNADA (Gantt) ═══
# Tupla: (build_start, build_end, uat_start, prod_golive, finetune_end) em nº de semana (1..22)
SCHED = {
    "J1": (3,5,6,8,10), "J2": (3,5,6,8,10),
    "SEI-J1": (4,5,6,8,10), "SEI-J2": (4,6,7,8,10), "SEI-J3": (4,6,7,8,10),
    "J3": (9,11,12,13,15), "J7": (9,11,12,13,15),
    "SEI-J4": (9,11,12,13,15), "SEI-J5": (10,12,13,15,17),
    "SEI-J6": (10,13,14,15,17), "SEI-J7": (10,13,14,15,17),
    "J4": (14,17,18,19,21), "SEI-J8": (15,17,18,19,21),
    "SEI-J9": (16,18,19,20,22), "SEI-J10": (16,18,19,20,22),
}
ONDA_LIGHT = {1: RGBColor(0xDC,0xFC,0xE7), 2: RGBColor(0xDB,0xEA,0xFE), 3: RGBColor(0xF3,0xE8,0xFF)}
GOLD = RGBColor(0xB8,0x86,0x0B)
MONTHS = [("AGO",1,4),("SET",5,8),("OUT",9,13),("NOV",14,17),("DEZ",18,22)]

s = new(); header(s, "Estimativa & ROM", "Cronograma por jornada — do build ao go-live", TEAL)
TL_X0 = 3.30; TL_W = 9.55; WK = TL_W/22.0
def wx(w):   # borda esquerda da semana w (1-indexada)
    return TL_X0 + (w-1)*WK
AX_Y = 1.32; WKAX_Y = 1.60; ROW0 = 1.88; RH = 0.245
# eixo de meses
for name, a, b in MONTHS:
    add_rect(s, wx(a), AX_Y, (b-a+1)*WK, 0.26, fill=GRAY_BG, line=GRAY_BORD)
    add_text(s, name, wx(a), AX_Y+0.02, (b-a+1)*WK, 0.22, size=9, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
# eixo de semanas (rótulos a cada 2) + gridlines de mês
for w in range(1, 23):
    if w % 2 == 1:
        add_text(s, f"S{w}", wx(w), WKAX_Y, WK*2, 0.2, size=6.5, color=GRAY, align=PP_ALIGN.CENTER)
def timeline_bar(y, jid, onda):
    b0,b1,u0,p,f1 = SCHED[jid]
    c = ONDA_COLOR[onda]; lt = ONDA_LIGHT[onda]
    # build
    add_rect(s, wx(b0), y+0.03, (b1-b0+1)*WK, 0.15, fill=c)
    # UAT
    add_rect(s, wx(u0), y+0.03, (p-u0)*WK if p>u0 else WK*0.6, 0.15, fill=lt, line=c, line_w=Pt(0.75))
    # go-live (diamante = quadradinho rotacionado)
    d = add_rect(s, wx(p)+WK*0.5-0.07, y+0.03, 0.14, 0.14, fill=GOLD); d.rotation = 45
    # fine-tuning / hypercare
    if f1 >= p+1:
        add_rect(s, wx(p+1), y+0.055, (f1-p)*WK, 0.10, fill=lt, line=c, line_w=Pt(0.5))
def group_header(y, txt, c):
    add_rect(s, 0.55, y, 12.23, 0.22, fill=c)
    add_text(s, txt, 0.68, y+0.005, 12, 0.21, size=8.5, bold=True, color=WHITE)
y = ROW0
# Fase 0
add_text(s, "🚧 Fase 0 · Governança G1002", 0.55, y+0.02, 2.7, 0.22, size=8, bold=True, color=ORANGE)
add_rect(s, wx(1), y+0.03, 4*WK, 0.15, fill=ORANGE)
y += RH
for onda, lbl, c in [(1,"Onda 1 · Quick Wins · leitura (go-live S8)", GREEN_D),
                     (2,"Onda 2 · Expansão · leitura + RAG + painéis (go-live S13/S15)", BLUE_ACC),
                     (3,"Onda 3 · Escrita, voz e governança · após Fase 0 (go-live S19/S20)", PURPLE)]:
    group_header(y, lbl, c); y += RH
    for j in J:
        if j[3] != onda:
            continue
        jid = j[0]
        add_text(s, f"{jid} · {j[1]}", 0.55, y+0.015, 2.7, 0.22, size=7.5, bold=True, color=BLUE_DARK)
        timeline_bar(y, jid, onda)
        y += RH
# legenda
ly = y + 0.05
add_rect(s, 0.55, ly, 0.4, 0.13, fill=TEAL); add_text(s, "Build", 1.0, ly-0.02, 1.2, 0.2, size=8, color=TEXT)
add_rect(s, 2.2, ly, 0.4, 0.13, fill=RGBColor(0xDB,0xEA,0xFE), line=TEAL); add_text(s, "Entrega p/ UAT", 2.65, ly-0.02, 1.8, 0.2, size=8, color=TEXT)
dg = add_rect(s, 4.6, ly-0.01, 0.14, 0.14, fill=GOLD); dg.rotation = 45
add_text(s, "Go-live Produção", 4.85, ly-0.02, 2.0, 0.2, size=8, color=TEXT)
add_rect(s, 7.0, ly+0.02, 0.4, 0.09, fill=RGBColor(0xDB,0xEA,0xFE), line=TEAL, line_w=Pt(0.5)); add_text(s, "Fine-tuning / Hypercare", 7.45, ly-0.02, 3.0, 0.2, size=8, color=TEXT)
footer(s, N); N += 1

# ═══ SLIDE 14 · PERFIS & HORAS ═══
s = new(); header(s, "Estimativa & ROM", "Perfis & horas — roster Ago–Dez", GREEN_D)
roster = [("MuleSoft Technical Architect (Sr)","1","800","R$ 614.144"),
          ("MuleSoft Technical Consultant ×2","1","1.560","R$ 974.953"),
          ("Agentforce Specialist / TC ×2","2","1.560","R$ 974.953"),
          ("Solution Architect (Slack)","3","560","R$ 413.358"),
          ("UX Conversacional / Experience","3","520","R$ 383.833"),
          ("QA Consultant ×2","—","1.293","R$ 693.552"),
          ("Program Manager","—","528→997*","R$ 735.926"),
          ("Change & Adoption Manager","—","400","R$ 295.256")]
tb = s.shapes.add_table(len(roster)+2, 4, Inches(0.55), Inches(1.5), Inches(9.0), Inches(4.6)).table
tb.columns[0].width = Inches(4.6); tb.columns[1].width = Inches(1.2); tb.columns[2].width = Inches(1.6); tb.columns[3].width = Inches(1.6)
for c,h in enumerate(["Perfil","Pilar","Horas","Custo c/imp"]):
    cell(tb, 0, c, h, size=10.5, bold=True, color=WHITE, fill=BLUE_MID)
for i,(p,pi,h,c) in enumerate(roster):
    fill = GRAY_BG if i%2 else WHITE
    cell(tb, i+1, 0, p, size=9.5, fill=fill); cell(tb, i+1, 1, pi, size=9.5, fill=fill, align=PP_ALIGN.CENTER)
    cell(tb, i+1, 2, h, size=9.5, fill=fill, align=PP_ALIGN.RIGHT); cell(tb, i+1, 3, c, size=9.5, bold=True, color=BLUE_DARK, fill=fill, align=PP_ALIGN.RIGHT)
cell(tb, 9, 0, "TOTAL", size=10, bold=True, color=WHITE, fill=BLUE_DARK); cell(tb,9,1,"",fill=BLUE_DARK)
cell(tb, 9, 2, "~7.650h", size=10, bold=True, color=WHITE, fill=BLUE_DARK, align=PP_ALIGN.RIGHT)
cell(tb, 9, 3, "~R$ 5,0 M", size=10, bold=True, color=YELLOW, fill=BLUE_DARK, align=PP_ALIGN.RIGHT)
add_rect(s, 9.75, 1.5, 3.03, 4.6, fill=BLUE_LIGHT, line=BLUE_ACC)
add_text(s, "Regras Dataprev ✓", 9.95, 1.65, 2.7, 0.4, size=12, bold=True, color=BLUE_DARK)
add_bullets(s, ["Nenhum recurso < 20h/sem","Ratio QA 1:2 TC/Dev","PM ≥ 15% do time (*997h)","Ganho de IA ≥ 25%"],
            9.95, 2.15, 2.7, 3.5, size=11, gap=10)
footer(s, N); N += 1

# ═══ SLIDE 15 · TRADICIONAL × IA ═══
s = new(); header(s, "Estimativa & ROM", "Tradicional × IA-Native", PURPLE)
add_rect(s, 0.55, 1.8, 5.9, 2.6, fill=WHITE, line=GRAY_BORD); add_rect(s,0.55,1.8,5.9,0.5,fill=GRAY)
add_text(s, "Sem ferramentas de IA", 0.75, 1.88, 5.5, 0.35, size=13, bold=True, color=WHITE)
add_text(s, "~30 sem", 0.75, 2.5, 5.5, 0.9, size=44, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, "baseline · 15 jornadas + 3 pilares", 0.75, 3.7, 5.5, 0.4, size=11, color=GRAY, align=PP_ALIGN.CENTER)
add_rect(s, 6.85, 1.8, 5.9, 2.6, fill=WHITE, line=GREEN_D, line_w=Pt(2)); add_rect(s,6.85,1.8,5.9,0.5,fill=GREEN_D)
add_text(s, "Com IA-Native", 7.05, 1.88, 5.5, 0.35, size=13, bold=True, color=WHITE)
add_text(s, "~22 sem", 7.05, 2.5, 5.5, 0.9, size=44, bold=True, color=GREEN_D, align=PP_ALIGN.CENTER)
add_text(s, "cabe em Ago–Dez 2026  ·  −27%", 7.05, 3.7, 5.5, 0.4, size=12, bold=True, color=GREEN_D, align=PP_ALIGN.CENTER)
add_rect(s, 0.55, 4.8, 12.23, 1.6, fill=BLUE_LIGHT, line=BLUE_ACC)
add_text(s, "Atividades comprimidas pela IA", 0.75, 4.92, 8, 0.4, size=13, bold=True, color=BLUE_DARK)
add_text(s, "Geração de fluxos de integração MuleSoft, prompts/tópicos dos agentes, testes automatizados, documentação técnica e artefatos. "
            "Sem IA, o mesmo escopo não caberia nem no bolsão nem na janela.",
         0.75, 5.4, 11.8, 0.9, size=12, color=TEXT)
footer(s, N); N += 1

# ═══ SLIDE 16 · JUSTIFICATIVA DE INVESTIMENTO (business case / ROI) ═══
s = new(); header(s, "Estimativa & ROM", "Justificativa de investimento — a lógica do ROI", GREEN_D)
# Faixa de enquadramento (a equação do retorno)
add_rect(s, 0.55, 1.4, 12.23, 0.9, fill=BLUE_LIGHT, line=BLUE_ACC)
add_text(s, "Investimento (teto R$ 5,0 M) → retorno por 4 alavancas de valor mensuráveis",
         0.75, 1.5, 11.8, 0.35, size=13, bold=True, color=BLUE_DARK)
add_text(s, "ROI = (ganho recorrente anual − investimento) ÷ investimento.  A metodologia de cálculo está pronta; "
            "os valores finais entram com a volumetria real da Dataprev (G1102 · Fase 0).",
         0.75, 1.85, 11.8, 0.4, size=11, color=TEXT)
# Tabela: alavanca · jornadas · KPI-chave · meta ilustrativa · driver de ROI
tb = s.shapes.add_table(5, 5, Inches(0.55), Inches(2.45), Inches(12.23), Inches(2.95)).table
tb.columns[0].width = Inches(2.5); tb.columns[1].width = Inches(1.9); tb.columns[2].width = Inches(3.0)
tb.columns[3].width = Inches(1.7); tb.columns[4].width = Inches(3.13)
for c, h in enumerate(["Alavanca de valor","Jornadas","KPI-chave","Meta ilustrativa","Como vira R$ (driver de ROI)"]):
    cell(tb, 0, c, h, size=10, bold=True, color=WHITE, fill=BLUE_MID)
levers = [
    (GREEN_D, "Autoatendimento & deflection", "J1 · J2 · J7",
     "% de consultas e chamados resolvidos sem intervenção humana", "40–70%",
     "consultas desviadas/mês × custo médio por atendimento"),
    (BLUE_ACC, "Produtividade do servidor", "J3 · SEI-J2\nSEI-J4 · SEI-J5",
     "tempo médio de consulta processual e de preparação de status", "−60 a −70%",
     "horas economizadas/mês × custo-hora do servidor"),
    (SEI_GOLD, "Risco & conformidade", "SEI-J1\nSEI-J9 · SEI-J10",
     "prazos perdidos por decurso tácito · % de ações com trilha de auditoria", "−80% perdas\n100% trilha",
     "multas, retrabalho e exposição evitados (valor de risco)"),
    (PURPLE, "Velocidade de ciclo & experiência", "J4 · SEI-J7\nSEI-J8",
     "tempo de ciclo (agendar / assinar / dar ciência) · adoção do Slack", "−50% tempo\nadoção ≥ 70%",
     "capacidade recuperada + adesão (uso sustentado)"),
]
for i,(col,alav,jor,kpi,meta,driver) in enumerate(levers):
    fill = GRAY_BG if i % 2 else WHITE
    cell(tb, i+1, 0, alav, size=9.5, bold=True, color=col, fill=fill)
    cell(tb, i+1, 1, jor, size=8.5, color=GRAY, fill=fill)
    cell(tb, i+1, 2, kpi, size=9, fill=fill)
    cell(tb, i+1, 3, meta, size=9.5, bold=True, color=GREEN_D, fill=fill, align=PP_ALIGN.CENTER)
    cell(tb, i+1, 4, driver, size=8.5, color=TEXT, fill=fill)
# Faixa: como levar à aprovação + disclaimer
add_rect(s, 0.55, 5.55, 12.23, 1.05, fill=RGBColor(0xFD,0xF0,0xE0), line=SEI_GOLD)
add_text(s, "Como levar à aprovação", 0.75, 5.63, 11.8, 0.3, size=11, bold=True, color=SEI_GOLD)
add_text(s, "1) inserir os volumes reais por processo (G1102) na coluna de metas  ·  2) converter em R$ pelos drivers de cada alavanca  ·  "
            "3) comparar o ganho anual recorrente com o investimento único → payback.  Metas ilustrativas — nenhum número comprometido nesta fase.",
         0.75, 5.95, 11.8, 0.6, size=10, color=TEXT)
footer(s, N); N += 1

# ═══ SLIDE 17 · KPIs ═══
s = new(); header(s, "Estimativa & ROM", "KPIs propostos — ROI por processo", ORANGE)
add_rect(s, 0.55, 1.4, 12.23, 0.75, fill=RGBColor(0xFD,0xF0,0xE0), line=YELLOW)
add_text(s, "Metodologia: metas ilustrativas — números finais dependem da volumetria (G1102), fornecida na Fase 0. Nenhum número comprometido nesta fase.",
         0.75, 1.5, 11.8, 0.6, size=11, bold=True, color=SEI_GOLD)
kpis = [("Financeiro (J1)","% consultas self-service","→ 70%"),
        ("Chamados (J2)","Deflection de status","→ 40%"),
        ("Briefing (J3)","Tempo de preparação","−60%"),
        ("FAQ/Conexão (J7)","Deflection RH/suporte","→ 50%"),
        ("SEI prazos (SEI-J1)","Perdas por decurso tácito","−80%"),
        ("SEI consulta (SEI-J2/J4)","Tempo de consulta processual","−70%"),
        ("SEI transação (J9/J10)","Tempo de tramitação/abertura","−50%"),
        ("Adoção Slack","Usuários ativos","≥ 70%"),
        ("Autonomia","Jornadas sem escalar a humano","≥ 65%")]
for i,(p,k,m) in enumerate(kpis):
    col = i % 3; row = i // 3
    l = 0.55 + col*4.08; t = 2.35 + row*1.35
    add_rect(s, l, t, 3.9, 1.2, fill=WHITE, line=GRAY_BORD)
    add_text(s, p, l+0.15, t+0.1, 3.6, 0.35, size=10.5, bold=True, color=BLUE_ACC)
    add_text(s, k, l+0.15, t+0.45, 3.6, 0.45, size=10, color=TEXT)
    add_text(s, m, l+0.15, t+0.85, 3.6, 0.35, size=13, bold=True, color=GREEN_D)
footer(s, N); N += 1

# ═══ SLIDE 17 · CLOUDS ═══
s = new(); header(s, "Execução", "Clouds necessárias — só o essencial", PURPLE)
clouds = [(SLACK,"💬 Slack (Grid)","Front conversacional · dual-workspace","por usuário/mês"),
          (BLUE_ACC,"🧠 Agentforce","Public Sector · agentes especialistas","Data Library cobre o RAG"),
          (TEAL,"🔌 MuleSoft","Anypoint Titanium · integração + MCP","suporta polling do SEI")]
for i,(c,t,d,ref) in enumerate(clouds):
    l = 0.55 + i*4.08
    add_rect(s, l, 1.7, 3.9, 3.0, fill=WHITE, line=GRAY_BORD); add_rect(s,l,1.7,3.9,0.55,fill=c)
    add_text(s, t, l+0.15, 1.78, 3.6, 0.4, size=13, bold=True, color=WHITE)
    add_text(s, d, l+0.2, 2.5, 3.5, 1.2, size=12, color=TEXT)
    add_text(s, ref, l+0.2, 3.9, 3.5, 0.6, size=10.5, italic=True, color=GRAY)
add_rect(s, 0.55, 5.0, 12.23, 1.4, fill=RGBColor(0xFF,0xFB,0xE6), line=GOV_YELLOW)
add_text(s, "⚠️ Data Cloud — condicional, fora do escopo base", 0.75, 5.12, 8, 0.4, size=12, bold=True, color=BLUE_DARK)
add_text(s, "RAG (J7, SEI-J6) usa a Data Library do Agentforce. Data Cloud só entra se a volumetria (G1102) exigir indexação em escala — decisão na Fase 0.",
         0.75, 5.55, 11.8, 0.7, size=11.5, color=TEXT)
footer(s, N); N += 1

# ═══ SLIDE 18 · HCC ═══
s = new(); header(s, "Execução", "HCC — Change Management & UX Conversacional", GREEN_D)
add_rect(s, 0.55, 1.6, 6.0, 4.5, fill=WHITE, line=GRAY_BORD); add_rect(s,0.55,1.6,6.0,0.55,fill=GREEN_D)
add_text(s, "🚀 Adoção do Slack", 0.75, 1.68, 5.6, 0.4, size=13, bold=True, color=WHITE)
add_bullets(s, ["Champions por unidade (multiplicadores).","Onboarding por onda (Set/Nov/Dez).",
                "Playbook de canais interno/externo (LGPD).","Medição de adoção e feedback loop.",
                "Do 'abro o sistema' para 'pergunto ao agente'."],
            0.8, 2.3, 5.5, 3.6, size=12, gap=12)
add_rect(s, 6.75, 1.6, 6.0, 4.5, fill=WHITE, line=GRAY_BORD); add_rect(s,6.75,1.6,6.0,0.55,fill=BLUE_ACC)
add_text(s, "🎨 UX/UI conversacional", 7.0, 1.68, 5.6, 0.4, size=13, bold=True, color=WHITE)
add_bullets(s, ["Tom e persona por agente.","Desenho de diálogo: intents, fallback, handoff.",
                "Block Kit: cards, botões, formulários.","Confirmação + trilha na escrita (SEI-J8/9/10).",
                "Voz→texto (J3, J4) e leitura mobile."],
            7.0, 2.3, 5.5, 3.6, size=12, gap=12)
add_text(s, "Perfis dedicados: Change & Adoption Manager (400h) + UX Conversacional / Experience Architect (520h) — atravessam as 3 ondas.",
         0.55, 6.25, 12.2, 0.5, size=11, italic=True, color=GRAY)
footer(s, N); N += 1

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "DATAAGIL_ROM_v5.pptx")
prs.save(OUT)
print("OK ->", OUT, "· slides:", len(prs.slides._sldIdLst))
