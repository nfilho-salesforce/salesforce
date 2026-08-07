"""
Gera o deck PPTX para CLARO Agente PLM — 12 slides em português.
Fonte: outputs/06-presentation-outline.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Salesforce brand palette ──────────────────────────────────────────────────
SF_NAVY   = RGBColor(0x03, 0x2D, 0x60)   # dark navy — backgrounds, headers
SF_BLUE   = RGBColor(0x01, 0x76, 0xD3)   # brand blue — accents, dividers
SF_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SF_LIGHT  = RGBColor(0xF3, 0xF3, 0xF3)   # light grey — content slide bg
SF_DARK   = RGBColor(0x16, 0x16, 0x16)   # near-black — body text
SF_MID    = RGBColor(0x44, 0x44, 0x44)   # grey — secondary text
SF_RED    = RGBColor(0xC2, 0x31, 0x2F)   # risk/alert
SF_GREEN  = RGBColor(0x2E, 0x84, 0x4E)   # positive / KPI

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_layout(prs):
    return prs.slide_layouts[6]   # blank

def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, x, y, w, h, text, font_size=18, bold=False,
                color=SF_DARK, align=PP_ALIGN.LEFT, wrap=True,
                italic=False, font_name="Calibri"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txb

def add_para(tf, text, font_size=16, bold=False, color=SF_DARK,
             align=PP_ALIGN.LEFT, italic=False, space_before=0, bullet=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    if bullet:
        p.level = 1
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p

def navy_header(slide, title, subtitle=None):
    """Full-width navy header band at top."""
    add_rect(slide, 0, 0, W, Inches(1.4), SF_NAVY)
    add_textbox(slide, Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.75),
                title, font_size=28, bold=True, color=SF_WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.4), Inches(0.88), Inches(12.5), Inches(0.4),
                    subtitle, font_size=14, color=SF_BLUE)
    # bottom accent stripe
    add_rect(slide, 0, Inches(1.4), W, Inches(0.05), SF_BLUE)

def footer_bar(slide, text="Salesforce Professional Services LATAM · Confidencial"):
    add_rect(slide, 0, Inches(7.1), W, Inches(0.4), SF_NAVY)
    add_textbox(slide, Inches(0.4), Inches(7.12), Inches(12.5), Inches(0.3),
                text, font_size=9, color=SF_WHITE, align=PP_ALIGN.LEFT)

def content_bg(slide):
    add_rect(slide, 0, 0, W, H, SF_LIGHT)

# ── Slide builders ────────────────────────────────────────────────────────────

def slide_title(prs):
    """Slide 1 — Title"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, W, H, SF_NAVY)
    add_rect(slide, 0, Inches(4.8), W, Inches(0.08), SF_BLUE)

    add_textbox(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(1.2),
                "CLARO Agente PLM", font_size=48, bold=True, color=SF_WHITE,
                align=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.6), Inches(2.5), Inches(12), Inches(0.8),
                "Autoria Inteligente de Catálogo com Agentforce",
                font_size=26, color=SF_BLUE, align=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.6), Inches(3.5), Inches(12), Inches(0.5),
                "Salesforce Professional Services LATAM  ·  POC  ·  Junho 2026",
                font_size=16, color=SF_WHITE, align=PP_ALIGN.LEFT)

    add_textbox(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(0.4),
                "Confidencial — Salesforce PS LATAM",
                font_size=11, color=SF_MID, align=PP_ALIGN.LEFT, italic=True)


def slide_problema(prs):
    """Slide 2 — O Problema"""
    slide = prs.slides.add_slide(blank_layout(prs))
    content_bg(slide)
    navy_header(slide, "O Problema: O Catálogo Que Não Acompanha o Negócio",
                subtitle='"127 regras. Cada mudança exige um deploy. O mercado não espera."')
    footer_bar(slide)

    bullets = [
        ("Pain 1 — Velocidade:",
         "Cada nova oferta, bundle ou ajuste de preço exige um sprint de engenharia e um deploy cross-org — processo de dias a semanas."),
        ("Pain 2 — Estabilidade:",
         "Validações volumosas fazem overflow de CPU e heap — o pipeline falha silenciosamente; erros chegam ao downstream sem ser detectados."),
        ("Pain 3 — Competitividade:",
         "Enquanto Vivo e TIM lançam combos em semanas, Claro aguarda aprovação de TI para alterar uma regra."),
    ]

    y = Inches(1.65)
    for label, body in bullets:
        add_rect(slide, Inches(0.4), y, Inches(0.06), Inches(0.9), SF_BLUE)
        add_textbox(slide, Inches(0.65), y, Inches(11.8), Inches(0.38),
                    label, font_size=15, bold=True, color=SF_NAVY)
        add_textbox(slide, Inches(0.65), y + Inches(0.38), Inches(11.8), Inches(0.55),
                    body, font_size=14, color=SF_DARK)
        y += Inches(1.25)

    add_textbox(slide, Inches(0.4), Inches(5.6), Inches(12.5), Inches(0.4),
                'Para Fabricio: "O catalogo e o coracao da oferta comercial da Claro. Se o coracao bate devagar, o negocio inteiro desacelera."',
                font_size=12, color=SF_MID, italic=True)


def slide_diagnostico(prs):
    """Slide 3 — O Diagnóstico"""
    slide = prs.slides.add_slide(blank_layout(prs))
    content_bg(slide)
    navy_header(slide, "O Diagnóstico: Onde o Sistema Falha")
    footer_bar(slide)

    # Table header
    cols = [Inches(3.2), Inches(3.2), Inches(4.5)]
    headers = ["Processo", "Impacto", "Causa raiz"]
    rows_data = [
        ("Alteração de regra", "Dias (sprint + deploy)", "BRE acoplado ao pipeline de engenharia"),
        ("Validação em lote", "Falhas silenciosas", "Processamento síncrono — overflow de heap"),
        ("Diagnóstico de erros", "Horas de investigação manual", "Sem observabilidade nativa; sem DLQ"),
        ("Auditoria LGPD / ANATEL", "Sem rastreabilidade", "BRE sem lineage de regras"),
    ]

    row_h = Inches(0.72)
    hdr_h = Inches(0.5)
    x_start = Inches(0.35)
    y_start = Inches(1.6)

    # Header row
    x = x_start
    for i, (hdr, col_w) in enumerate(zip(headers, cols)):
        add_rect(slide, x, y_start, col_w, hdr_h, SF_NAVY)
        add_textbox(slide, x + Inches(0.1), y_start + Inches(0.08), col_w - Inches(0.15), hdr_h - Inches(0.1),
                    hdr, font_size=14, bold=True, color=SF_WHITE)
        x += col_w

    # Data rows
    for r_idx, row in enumerate(rows_data):
        y = y_start + hdr_h + r_idx * row_h
        bg = SF_WHITE if r_idx % 2 == 0 else RGBColor(0xE8, 0xF4, 0xFD)
        x = x_start
        for i, (cell, col_w) in enumerate(zip(row, cols)):
            add_rect(slide, x, y, col_w, row_h, bg, SF_BLUE)
            txt_color = SF_RED if i == 1 else SF_DARK
            add_textbox(slide, x + Inches(0.1), y + Inches(0.1), col_w - Inches(0.15), row_h - Inches(0.12),
                        cell, font_size=13, color=txt_color)
            x += col_w


def slide_solucao(prs):
    """Slide 4 — A Solução"""
    slide = prs.slides.add_slide(blank_layout(prs))
    content_bg(slide)
    navy_header(slide, "A Solução em Uma Frase",
                subtitle='"Compilação por IA. Execução determinística. Zero redeploy."')
    footer_bar(slide)

    # Two lanes
    lane_w = Inches(5.9)
    lane_h = Inches(4.5)
    y_lane = Inches(1.7)

    # Lane 1
    add_rect(slide, Inches(0.35), y_lane, lane_w, lane_h, SF_NAVY)
    add_textbox(slide, Inches(0.55), y_lane + Inches(0.15), lane_w - Inches(0.3), Inches(0.5),
                "AUTORIA — Admin Agent", font_size=15, bold=True, color=SF_BLUE)
    steps1 = [
        "1. Analista escreve regra em português",
        "2. Admin Agent → Einstein LLM compila",
        "3. AST JSON gravado em Plm_Rule_Spec__c",
        "4. Regra live em segundos — sem deploy",
    ]
    for i, s in enumerate(steps1):
        add_textbox(slide, Inches(0.55), y_lane + Inches(0.75) + i * Inches(0.82),
                    lane_w - Inches(0.3), Inches(0.7),
                    s, font_size=13, color=SF_WHITE)

    # Lane 2
    add_rect(slide, Inches(7.05), y_lane, lane_w, lane_h, SF_BLUE)
    add_textbox(slide, Inches(7.25), y_lane + Inches(0.15), lane_w - Inches(0.3), Inches(0.5),
                "OPERAÇÃO — Ops Agent", font_size=15, bold=True, color=SF_WHITE)
    steps2 = [
        "1. Lote CSV enviado via LWC Wizard",
        "2. Ops Agent dispara ingestão assíncrona",
        "3. AST Walker avalia: <50ms/registro",
        "4. Relatório HTML gerado por IA",
    ]
    for i, s in enumerate(steps2):
        add_textbox(slide, Inches(7.25), y_lane + Inches(0.75) + i * Inches(0.82),
                    lane_w - Inches(0.3), Inches(0.7),
                    s, font_size=13, color=SF_WHITE)

    add_textbox(slide, Inches(0.35), Inches(6.5), Inches(12.5), Inches(0.35),
                "Sem redeploy. Sem desenvolvedor. Sem falhas silenciosas.",
                font_size=13, bold=True, color=SF_NAVY, align=PP_ALIGN.CENTER)


def slide_arquitetura(prs):
    """Slide 5 — Arquitetura"""
    slide = prs.slides.add_slide(blank_layout(prs))
    content_bg(slide)
    navy_header(slide, "Arquitetura de Alto Nível",
                subtitle='"Três camadas: experiência, inteligência, motor."')
    footer_bar(slide)

    layers = [
        ("EXPERIENCE",   "LWC Upload Wizard  ·  Agentforce Chat Console",                          SF_BLUE,  SF_WHITE),
        ("INTELLIGENCE", "Admin Agent (compile-time)  ·  Ops Agent (runtime)\nAtlas Reasoning Engine  ·  Einstein Trust Layer", SF_NAVY,  SF_WHITE),
        ("ENGINE",       "AST Walker (Apex puro, <50ms)  ·  DLQ + Finalizers\nPlatform Events  ·  Compile Snapshots (LGPD lineage)", RGBColor(0x06, 0x4E, 0x8C), SF_WHITE),
    ]

    y = Inches(1.65)
    for label, content, bg, fg in layers:
        box_h = Inches(1.55)
        add_rect(slide, Inches(0.35), y, Inches(2.0), box_h, SF_NAVY)
        add_textbox(slide, Inches(0.35), y + Inches(0.5), Inches(2.0), Inches(0.6),
                    label, font_size=13, bold=True, color=SF_BLUE, align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(2.35), y, Inches(10.6), box_h, bg)
        add_textbox(slide, Inches(2.55), y + Inches(0.35), Inches(10.2), box_h - Inches(0.4),
                    content, font_size=14, color=fg)
        y += box_h + Inches(0.12)

    add_textbox(slide, Inches(0.35), Inches(6.5), Inches(12.5), Inches(0.35),
                "LLM apenas na autoria. Em produção: Apex puro — determinístico, rápido, auditável. Sem IA no caminho crítico.",
                font_size=12, bold=True, color=SF_NAVY, align=PP_ALIGN.CENTER, italic=True)


def slide_kpis(prs):
    """Slide 6 — KPIs"""
    slide = prs.slides.add_slide(blank_layout(prs))
    content_bg(slide)
    navy_header(slide, "Os KPIs: O Que Muda")
    footer_bar(slide)

    cols_w = [Inches(4.5), Inches(3.5), Inches(4.5)]
    headers = ["KPI", "Hoje", "Com Agentforce PLM"]
    rows_data = [
        ("Tempo para alterar uma regra",    "Dias (sprint + deploy)",     "0 minutos — sem redeploy"),
        ("Tempo de avaliação por registro", "Minutos (lote síncrono)",    "< 50ms por registro"),
        ("Capacidade por lote",             "Instável em volumes altos",  "≤ 10.000 linhas — estável"),
        ("Falhas silenciosas",              "Sim — não rastreadas",        "Zero — 100% capturado em DLQ"),
        ("Investigação de erros",           "Horas (manual)",             "Segundos — relatório HTML por IA"),
        ("Auditoria LGPD",                  "Sem rastreabilidade",         "Lineage completo — Spec_Key + Snapshots"),
    ]

    row_h = Inches(0.65)
    hdr_h = Inches(0.48)
    x_start = Inches(0.35)
    y_start = Inches(1.6)

    x = x_start
    for hdr, col_w in zip(headers, cols_w):
        add_rect(slide, x, y_start, col_w, hdr_h, SF_NAVY)
        add_textbox(slide, x + Inches(0.08), y_start + Inches(0.08),
                    col_w - Inches(0.12), hdr_h - Inches(0.08),
                    hdr, font_size=13, bold=True, color=SF_WHITE)
        x += col_w

    for r_idx, row in enumerate(rows_data):
        y = y_start + hdr_h + r_idx * row_h
        bg = SF_WHITE if r_idx % 2 == 0 else RGBColor(0xE8, 0xF4, 0xFD)
        x = x_start
        for i, (cell, col_w) in enumerate(zip(row, cols_w)):
            add_rect(slide, x, y, col_w, row_h, bg, SF_BLUE)
            txt_color = SF_RED if i == 1 else (SF_GREEN if i == 2 else SF_DARK)
            bold = i == 2
            add_textbox(slide, x + Inches(0.08), y + Inches(0.08),
                        col_w - Inches(0.12), row_h - Inches(0.1),
                        cell, font_size=12, color=txt_color, bold=bold)
            x += col_w


def slide_beachhead(prs):
    """Slide 7 — Por Que Agora: O Beachhead"""
    slide = prs.slides.add_slide(blank_layout(prs))
    content_bg(slide)
    navy_header(slide, "Por Que Agora: O POC Como Beachhead",
                subtitle='"O PLM é o primeiro de quatro agentes."')
    footer_bar(slide)

    agents = [
        ("Agente 1", "Knowledge Base",     "Atendimento ao cliente via base de conhecimento",  SF_MID,  False),
        ("Agente 2", "PLM — ESTE POC",     "Referência arquitetural para todos os outros agentes", SF_NAVY, True),
        ("Agente 3", "Next Best Offer",    "Propensão de oferta via sistema NBO",               SF_MID,  False),
        ("Agente 4", "Lead Qualification", "Qualificação de leads",                             SF_MID,  False),
    ]

    y = Inches(1.7)
    for num, name, desc, color, highlight in agents:
        h = Inches(1.1) if highlight else Inches(0.9)
        bg = SF_NAVY if highlight else SF_WHITE
        add_rect(slide, Inches(0.35), y, Inches(12.6), h, bg, SF_BLUE)
        num_color = SF_BLUE if highlight else SF_BLUE
        add_textbox(slide, Inches(0.55), y + Inches(0.1), Inches(1.2), h - Inches(0.15),
                    num, font_size=13, bold=True, color=num_color)
        name_color = SF_WHITE if highlight else SF_NAVY
        add_textbox(slide, Inches(1.9), y + Inches(0.1), Inches(3.5), Inches(0.45),
                    name, font_size=15 if highlight else 13,
                    bold=highlight, color=name_color)
        desc_color = RGBColor(0xA8, 0xD8, 0xFF) if highlight else SF_MID
        add_textbox(slide, Inches(1.9), y + Inches(0.5), Inches(9.5), Inches(0.45),
                    desc, font_size=12, color=desc_color, italic=not highlight)
        if highlight:
            add_rect(slide, Inches(0.35), y, Inches(0.08), h, SF_BLUE)
        y += h + Inches(0.1)

    add_textbox(slide, Inches(0.35), Inches(6.5), Inches(12.5), Inches(0.35),
                "Aprovando este POC, a Claro aprova a arquitetura da plataforma completa.",
                font_size=13, bold=True, color=SF_NAVY, align=PP_ALIGN.CENTER, italic=True)


def slide_plano(prs):
    """Slide 8 — Plano de Entrega: 8 Semanas"""
    slide = prs.slides.add_slide(blank_layout(prs))
    content_bg(slide)
    navy_header(slide, "Plano de Entrega: 8 Semanas")
    footer_bar(slide)

    rows_data = [
        ("Semana",  "Fase",                       "O Que Acontece"),
        ("1–2",     "Descoberta & Arquitetura",    "Alinhamento técnico, dependências confirmadas, registro de testes acordado"),
        ("3",       "Sprint 1 — Fundação",         "Motor de ingestão CSV + segurança e configuração de plataforma"),
        ("4",       "Sprint 2 — Motor Core",       "AST Walker + compilador LLM + observabilidade e DLQ"),
        ("5",       "Sprint 3 — Inteligência & UX","Agentes Admin + Ops + narrativa diagnóstica + componentes LWC"),
        ("6–7",     "UAT & Fine-tuning",           "Testes de carga, run paralelo vs. BRE legado, revisão legal LGPD"),
        ("8",       "Go-Live & Hipercuidado",      "Deploy em produção, transferência de conhecimento, encerramento formal"),
    ]

    col_w = [Inches(1.2), Inches(3.8), Inches(7.6)]
    row_h = Inches(0.7)
    x_start = Inches(0.35)
    y = Inches(1.6)

    for r_idx, row in enumerate(rows_data):
        x = x_start
        is_hdr = r_idx == 0
        for i, (cell, cw) in enumerate(zip(row, col_w)):
            if is_hdr:
                bg = SF_NAVY
                tc = SF_WHITE
                bold = True
                fs = 13
            elif r_idx in (2, 3, 4):   # sprints
                bg = RGBColor(0xE8, 0xF4, 0xFD) if r_idx % 2 == 0 else SF_WHITE
                tc = SF_DARK
                bold = False
                fs = 12
            else:
                bg = SF_WHITE if r_idx % 2 == 0 else RGBColor(0xF0, 0xF0, 0xF0)
                tc = SF_DARK
                bold = False
                fs = 12
            add_rect(slide, x, y, cw, row_h, bg, SF_BLUE)
            add_textbox(slide, x + Inches(0.08), y + Inches(0.1),
                        cw - Inches(0.12), row_h - Inches(0.12),
                        cell, font_size=fs, bold=bold, color=tc)
            x += cw
        y += row_h

    add_textbox(slide, Inches(0.35), Inches(6.65), Inches(12.5), Inches(0.3),
                "⚠  6 pré-requisitos que a Claro precisa confirmar nas primeiras 2 semanas (sandbox, licenças, artigos de conhecimento, LGPD)",
                font_size=11, color=SF_NAVY, italic=True)


def slide_time(prs):
    """Slide 9 — O Time Salesforce PS"""
    slide = prs.slides.add_slide(blank_layout(prs))
    content_bg(slide)
    navy_header(slide, "O Time Salesforce PS",
                subtitle='"A equipe que construiu este padrão."')
    footer_bar(slide)

    roles = [
        ("Arquiteto Técnico",   "8 semanas — integral",              "Decisões de arquitetura, Agentforce, Einstein Trust Layer, Knowledge Transfer"),
        ("Consultor Técnico",   "8 semanas — integral",              "Desenvolvimento Apex, LWC, testes unitários, deploys via Salesforce CLI"),
        ("Especialista em QA",  "8 semanas — 1.5x",                  "Estratégia de testes, validação KPIs, testes de carga, UAT facilitado"),
        ("Gerente de Projeto",  "8 semanas — dedicado, faturável*",  "Rastreamento de milestones, coordenação de dependências, gestão de riscos"),
    ]

    col_w = [Inches(2.8), Inches(2.8), Inches(7.1)]
    headers = ["Papel", "Dedicação", "Responsabilidades-chave"]
    hdr_h = Inches(0.48)
    row_h = Inches(1.1)
    x_start = Inches(0.35)
    y = Inches(1.6)

    x = x_start
    for hdr, cw in zip(headers, col_w):
        add_rect(slide, x, y, cw, hdr_h, SF_NAVY)
        add_textbox(slide, x + Inches(0.08), y + Inches(0.08),
                    cw - Inches(0.12), hdr_h - Inches(0.1),
                    hdr, font_size=13, bold=True, color=SF_WHITE)
        x += cw
    y += hdr_h

    for r_idx, (role, ded, resp) in enumerate(roles):
        bg = SF_WHITE if r_idx % 2 == 0 else RGBColor(0xE8, 0xF4, 0xFD)
        x = x_start
        for cell, cw in zip([role, ded, resp], col_w):
            add_rect(slide, x, y, cw, row_h, bg, SF_BLUE)
            add_textbox(slide, x + Inches(0.08), y + Inches(0.12),
                        cw - Inches(0.12), row_h - Inches(0.15),
                        cell, font_size=12, color=SF_DARK)
            x += cw
        y += row_h

    add_textbox(slide, Inches(0.35), Inches(6.6), Inches(12.5), Inches(0.3),
                "* Gerente de Projeto é um recurso PS dedicado, faturável, pago diretamente pela Claro.",
                font_size=10, color=SF_MID, italic=True)


def slide_riscos(prs):
    """Slide 10 — Riscos e Mitigações"""
    slide = prs.slides.add_slide(blank_layout(prs))
    content_bg(slide)
    navy_header(slide, "Riscos e Mitigações")
    footer_bar(slide)

    rows_data = [
        ("Artigos de Conhecimento não prontos na semana 1",   "Alta",   "Gate formal: Sprint 3 não começa sem KB confirmada. PM escalona imediatamente."),
        ("Licença Agentforce não ativa no sandbox",           "Média",  "PS valida via ConnectApi no fim da semana 2. Blocker escalado ao CSM Salesforce."),
        ("Overflow em lotes > 10k linhas",                    "Média",  "Arquitetura assíncrona mitiga até 10k. Acima de 50k no backlog pós-POC."),
        ("LGPD — sign-off jurídico atrasado",                 "Média",  "Janela da semana 6 agendada na semana 2. Sem sign-off = sem go-live. Sem exceções."),
        ("Deploy cross-org bloqueado (CannotQuickDeployError)","Confirmado","Registro RunSpecifiedTests acordado com DevOps Claro nas semanas 1–2."),
    ]

    col_w = [Inches(5.0), Inches(1.6), Inches(6.4)]
    headers = ["Risco", "Probab.", "Mitigação"]
    hdr_h = Inches(0.45)
    row_h = Inches(0.9)
    x_start = Inches(0.35)
    y = Inches(1.6)

    x = x_start
    for hdr, cw in zip(headers, col_w):
        add_rect(slide, x, y, cw, hdr_h, SF_NAVY)
        add_textbox(slide, x + Inches(0.08), y + Inches(0.06),
                    cw - Inches(0.12), hdr_h - Inches(0.08),
                    hdr, font_size=13, bold=True, color=SF_WHITE)
        x += cw
    y += hdr_h

    prob_colors = {"Alta": SF_RED, "Média": RGBColor(0xD9, 0x7F, 0x00), "Confirmado": SF_RED}
    for r_idx, (risk, prob, mit) in enumerate(rows_data):
        bg = SF_WHITE if r_idx % 2 == 0 else RGBColor(0xF5, 0xF5, 0xF5)
        x = x_start
        for i, (cell, cw) in enumerate(zip([risk, prob, mit], col_w)):
            add_rect(slide, x, y, cw, row_h, bg, SF_BLUE)
            tc = prob_colors.get(cell, SF_DARK) if i == 1 else SF_DARK
            bold = i == 1
            add_textbox(slide, x + Inches(0.08), y + Inches(0.1),
                        cw - Inches(0.12), row_h - Inches(0.12),
                        cell, font_size=12, color=tc, bold=bold)
            x += cw
        y += row_h


def slide_prereqs(prs):
    """Slide 11 — Pré-Requisitos"""
    slide = prs.slides.add_slide(blank_layout(prs))
    content_bg(slide)
    navy_header(slide, "O Que Precisamos da Claro",
                subtitle='"Seis dependências. Todas resolvíveis em duas semanas."')
    footer_bar(slide)

    rows_data = [
        ("1", "Cadeia de sandboxes provisionada (Dev + SIT + Ibuy UAT)", "Luciano",            "Fim da semana 1"),
        ("2", "3–5 CSVs de produção para profiling de volume",            "Lucas",              "Fim da semana 1"),
        ("3", "Licença Agentforce Unlimited + créditos Einstein ativos",  "Luciano",            "Fim da semana 2"),
        ("4", "10–15 FAQs com Data Categories no sandbox",               "Fabrício / Analistas","Fim da semana 2"),
        ("5", "Registro RunSpecifiedTests acordado com PS",               "Luciano / DevOps",   "Semanas 1–2"),
        ("6", "Janela de revisão jurídica LGPD agendada na semana 6",    "Jurídico Claro",     "Agendada na semana 2"),
    ]

    col_w = [Inches(0.5), Inches(5.5), Inches(3.2), Inches(3.4)]
    headers = ["#", "Pré-requisito", "Responsável", "Prazo"]
    hdr_h = Inches(0.45)
    row_h = Inches(0.78)
    x_start = Inches(0.35)
    y = Inches(1.6)

    x = x_start
    for hdr, cw in zip(headers, col_w):
        add_rect(slide, x, y, cw, hdr_h, SF_NAVY)
        add_textbox(slide, x + Inches(0.06), y + Inches(0.06),
                    cw - Inches(0.08), hdr_h - Inches(0.08),
                    hdr, font_size=13, bold=True, color=SF_WHITE)
        x += cw
    y += hdr_h

    for r_idx, row in enumerate(rows_data):
        bg = SF_WHITE if r_idx % 2 == 0 else RGBColor(0xE8, 0xF4, 0xFD)
        x = x_start
        for cell, cw in zip(row, col_w):
            add_rect(slide, x, y, cw, row_h, bg, SF_BLUE)
            add_textbox(slide, x + Inches(0.06), y + Inches(0.1),
                        cw - Inches(0.08), row_h - Inches(0.12),
                        cell, font_size=12, color=SF_DARK)
            x += cw
        y += row_h

    add_textbox(slide, Inches(0.35), Inches(6.6), Inches(12.5), Inches(0.3),
                "O PM Salesforce rastreia todas as seis dependências desde o primeiro dia.",
                font_size=11, color=SF_NAVY, bold=True, italic=True)


def slide_proximos(prs):
    """Slide 12 — Próximos Passos"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, W, H, SF_NAVY)
    add_rect(slide, 0, Inches(1.4), W, Inches(0.06), SF_BLUE)
    footer_bar(slide)

    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.9),
                "Próximos Passos", font_size=32, bold=True, color=SF_WHITE)
    add_textbox(slide, Inches(0.5), Inches(0.95), Inches(12), Inches(0.35),
                '"Três ações para colocar o relógio em movimento."',
                font_size=15, color=SF_BLUE, italic=True)

    steps = [
        ("1", "Confirmar pré-requisitos de ambiente",
              "Luciano provisiona sandbox chain na semana 1"),
        ("2", "Agendar kickoff técnico",
              "TA + Lucas + Luciano para alinhar schema JSON e registro RunSpecifiedTests"),
        ("3", "Aprovar engajamento",
              "Confirmar modelo T&M e equipe PS LATAM alocada"),
    ]

    y = Inches(1.7)
    for num, action, detail in steps:
        add_rect(slide, Inches(0.4), y, Inches(0.8), Inches(1.1), SF_BLUE)
        add_textbox(slide, Inches(0.4), y + Inches(0.2), Inches(0.8), Inches(0.7),
                    num, font_size=32, bold=True, color=SF_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.45), y + Inches(0.05), Inches(11.2), Inches(0.45),
                    action, font_size=17, bold=True, color=SF_WHITE)
        add_textbox(slide, Inches(1.45), y + Inches(0.55), Inches(11.2), Inches(0.45),
                    detail, font_size=14, color=RGBColor(0xA8, 0xD8, 0xFF))
        y += Inches(1.3)

    add_textbox(slide, Inches(0.4), Inches(5.5), Inches(12.5), Inches(0.8),
                "Em 8 semanas, a Claro terá o catálogo que o negócio merece —\ngovernado por analistas, validado em milissegundos, auditável para o LGPD.",
                font_size=15, color=SF_WHITE, align=PP_ALIGN.CENTER, italic=True)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    slide_title(prs)
    slide_problema(prs)
    slide_diagnostico(prs)
    slide_solucao(prs)
    slide_arquitetura(prs)
    slide_kpis(prs)
    slide_beachhead(prs)
    slide_plano(prs)
    slide_time(prs)
    slide_riscos(prs)
    slide_prereqs(prs)
    slide_proximos(prs)

    out = "/Users/nfilho/claude/Scopezilla/CLARO-AGENT-PLM/outputs/CLARO_Agente_PLM_Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
