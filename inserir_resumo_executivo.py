from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy
from lxml import etree

SF_BLUE       = RGBColor(0x00, 0x96, 0xFF)
SF_DARK_BLUE  = RGBColor(0x03, 0x2D, 0x60)
SF_LIGHT_BLUE = RGBColor(0xD4, 0xEE, 0xFF)
SF_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
SF_GRAY       = RGBColor(0x7F, 0x8C, 0x8D)
SF_LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF9)
SF_DARK_GRAY  = RGBColor(0x32, 0x3E, 0x48)
SF_GREEN      = RGBColor(0x2E, 0x7D, 0x32)
SF_ORANGE     = RGBColor(0xFF, 0x6B, 0x00)
SF_TEAL       = RGBColor(0x00, 0x6D, 0x9A)

PATH = "/Users/nfilho/claude/DATAPREV_SGP_Proposta_Tecnica.pptx"
prs = Presentation(PATH)
BLANK = prs.slide_layouts[6]
TOTAL_SLIDES = len(prs.slides) + 1  # +1 porque vamos inserir

def add_rect(slide, l, t, w, h, fill=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=14, bold=False, color=SF_DARK_GRAY,
             align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def slide_footer(slide, num, total):
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill=SF_DARK_BLUE)
    add_text(slide, "DATAPREV  |  SISDIP / DFT  |  Proposta Técnica — Salesforce Professional Services",
             0.3, 7.12, 11, 0.3, size=9, color=SF_WHITE)
    add_text(slide, f"{num} / {total}", 12.2, 7.12, 1, 0.3,
             size=9, color=SF_WHITE, align=PP_ALIGN.RIGHT)

# ─── Criar o novo slide em memória ───────────────────────────────────────────
new_slide = prs.slides.add_slide(BLANK)

s = new_slide
add_rect(s, 0, 0, 13.33, 7.5, fill=SF_LIGHT_GRAY)

# topo
add_rect(s, 0, 0, 13.33, 0.08, fill=SF_BLUE)
add_text(s, "Resumo Executivo", 0.4, 0.15, 10, 0.55,
         size=22, bold=True, color=SF_DARK_BLUE)
add_text(s, "SISDIP / DFT — Dimensionamento da Força de Trabalho  ·  DATAPREV / MGI", 0.4, 0.65, 12, 0.35,
         size=13, color=SF_GRAY)
add_rect(s, 0.4, 1.05, 12.53, 0.03, fill=SF_BLUE)

slide_footer(s, 2, TOTAL_SLIDES)

# ── Coluna esquerda: contexto + solução ──────────────────────────────────────

# Contexto
add_rect(s, 0.4, 1.18, 8.1, 0.42, fill=SF_DARK_BLUE)
add_text(s, "CONTEXTO", 0.55, 1.2, 7.8, 0.36, size=11, bold=True, color=SF_WHITE)

add_rect(s, 0.4, 1.62, 8.1, 1.45, fill=SF_WHITE)
add_text(s,
    "A DATAPREV opera o SISDIP como nó central do ecossistema de Dimensionamento da Força de Trabalho "
    "(DFT) do governo federal. Hoje, dados críticos de servidores estão fragmentados em 8+ plataformas "
    "(PGD, SEI, PEI, Recruta, entre outros), tornando a análise de cargos, perfis e movimentações um "
    "processo lento, manual e suscetível a erros. O MGI precisa de decisões rápidas e embasadas sobre "
    "força de trabalho — e o SISDIP precisa evoluir para suportar essa demanda.",
    0.55, 1.68, 7.8, 1.3, size=11, color=SF_DARK_GRAY)

# Solução
add_rect(s, 0.4, 3.15, 8.1, 0.42, fill=SF_BLUE)
add_text(s, "SOLUÇÃO PROPOSTA", 0.55, 3.17, 7.8, 0.36, size=11, bold=True, color=SF_WHITE)

add_rect(s, 0.4, 3.59, 8.1, 2.55, fill=SF_WHITE)

solucao_items = [
    ("Integração unificada via MuleSoft",
     "Conecta SISDIP a PGD, SEI, PEI, Recruta e outros em tempo real, sem replicar dados."),
    ("Agentforce — 2 agentes de IA",
     "Agente de Análise de Cargo (sugere cargo por entrega em ≤60s) e Agente de Resumo Executivo (consolida DFT por órgão sob demanda)."),
    ("Experiência do analista e do gestor",
     "Experience Cloud para revisão/aprovação de perfis. Tableau para painéis executivos. Slack para alertas e colaboração."),
    ("Governança e LGPD nativos",
     "Shield, Einstein Trust Layer e Zero Copy garantem rastreabilidade, soberania de dados e conformidade sem esforço adicional."),
]
for i, (titulo, desc) in enumerate(solucao_items):
    ly = 3.65 + i * 0.6
    add_rect(s, 0.45, ly, 0.1, 0.42, fill=SF_BLUE)
    add_text(s, titulo, 0.65, ly, 3.4, 0.24, size=10, bold=True, color=SF_DARK_BLUE)
    add_text(s, desc,   0.65, ly + 0.24, 7.6, 0.32, size=10, color=SF_DARK_GRAY)

# ── Coluna direita: destaques ─────────────────────────────────────────────────
destaques = [
    (SF_BLUE,      "Plataforma já contratada",
     "Service Cloud, MuleSoft, Tableau e Shield estão no contrato DTP vigente — sem novo procurement."),
    (SF_GREEN,     "Entrega incremental",
     "5 fases com Quality Gates. Valor visível desde a Fase 1, com MVP funcional ao final da Fase 3."),
    (SF_ORANGE,    "IA governada e auditável",
     "Agentforce com Einstein Trust Layer: nenhum dado sensível trafega fora do ambiente controlado."),
    (SF_DARK_BLUE, "Escalável para todo o governo",
     "Arquitetura multi-órgão: mesma plataforma expande para novos ministérios sem retrabalho."),
]

for i, (cor, titulo, desc) in enumerate(destaques):
    ly = 1.18 + i * 1.5
    add_rect(s, 8.7, ly, 4.25, 1.35, fill=SF_WHITE)
    add_rect(s, 8.7, ly, 0.18, 1.35, fill=cor)
    add_text(s, titulo, 8.98, ly + 0.1, 3.85, 0.32, size=11, bold=True, color=cor)
    add_text(s, desc,   8.98, ly + 0.45, 3.85, 0.8,  size=10, color=SF_DARK_GRAY)

# ── Barra de CTA ──────────────────────────────────────────────────────────────
add_rect(s, 0.4, 6.25, 12.53, 0.58, fill=SF_DARK_BLUE)
add_text(s,
    "Próximo passo recomendado:  Workshop de Discovery com DATAPREV + MGI  "
    "para confirmar licenças, mapear APIs e definir cronograma de entrega.",
    0.55, 6.3, 12.2, 0.45, size=11, bold=False, color=SF_WHITE)

# ─── Mover o novo slide para a posição 2 (índice 1) ──────────────────────────
xml_slides = prs.slides._sldIdLst
# O novo slide foi adicionado ao final — mover para posição 1 (após capa)
slides_xml = list(xml_slides)
# O último elemento é o novo slide
new_entry = slides_xml[-1]
xml_slides.remove(new_entry)
xml_slides.insert(1, new_entry)

prs.save(PATH)
print(f"Slide Resumo Executivo inserido na posição 2.")
print(f"Total de slides: {len(prs.slides)}")
