from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ─── Paleta Salesforce ───────────────────────────────────────────────────────
SF_BLUE       = RGBColor(0x00, 0x96, 0xFF)
SF_DARK_BLUE  = RGBColor(0x03, 0x2D, 0x60)
SF_LIGHT_BLUE = RGBColor(0xD4, 0xEE, 0xFF)
SF_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
SF_GRAY       = RGBColor(0x7F, 0x8C, 0x8D)
SF_LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF9)
SF_DARK_GRAY  = RGBColor(0x32, 0x3E, 0x48)
SF_GREEN      = RGBColor(0x2E, 0x7D, 0x32)
SF_ORANGE     = RGBColor(0xFF, 0x6B, 0x00)
SF_TEAL       = RGBColor(0x00, 0x69, 0x5C)
SF_PURPLE     = RGBColor(0x6A, 0x1B, 0x9A)
SF_AMBER      = RGBColor(0xF9, 0xA8, 0x25)

FLUXO_IMG = "/Users/nfilho/claude/SEFIN_CE_Fluxo_v1_FINAL.png"
OUTPUT    = "/Users/nfilho/claude/SEFIN_CE_Fluxo_Premissas_Perguntas.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

TOTAL = 13

# ─── Helpers ─────────────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill=None, line_color=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=11, bold=False, italic=False,
        color=SF_DARK_GRAY, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return txb

def add_para(tf, text, size=10, bold=False, italic=False,
             color=SF_DARK_GRAY, align=PP_ALIGN.LEFT, sp_before=0):
    p = tf.add_paragraph(); p.alignment = align
    if sp_before: p.space_before = Pt(sp_before)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color

def header(slide, title, sub=None, accent=SF_BLUE):
    rect(slide, 0, 0, 13.33, 0.08, fill=accent)
    txt(slide, title, 0.35, 0.12, 10, 0.55, size=22, bold=True, color=SF_DARK_BLUE)
    if sub:
        txt(slide, sub, 0.35, 0.62, 12.5, 0.35, size=11, italic=True, color=SF_GRAY)
    rect(slide, 0.35, 0.98, 12.63, 0.025, fill=accent)

def footer(slide, num):
    rect(slide, 0, 7.1, 13.33, 0.4, fill=SF_DARK_BLUE)
    txt(slide, "DATAPREV  |  SEFIN-CE  ·  Régua de Cobrança via WhatsApp  |  Agentforce  ·  Marketing Cloud",
        0.3, 7.12, 11, 0.3, size=8.5, color=SF_WHITE)
    txt(slide, f"{num} / {TOTAL}", 12.2, 7.12, 1, 0.3,
        size=8.5, color=SF_WHITE, align=PP_ALIGN.RIGHT)

def bg(slide):
    rect(slide, 0, 0, 13.33, 7.5, fill=SF_LIGHT_GRAY)

def label_box(slide, text, l, t, w, h, bg_color, text_color=SF_WHITE,
              size=10, bold=True):
    rect(slide, l, t, w, h, fill=bg_color)
    txt(slide, text, l + 0.08, t + 0.05, w - 0.16, h - 0.1,
        size=size, bold=bold, color=text_color)

def badge(slide, text, l, t, size=10, color=SF_BLUE, bg=SF_LIGHT_BLUE):
    w = len(text) * 0.085 + 0.25
    rect(slide, l, t, w, 0.28, fill=bg)
    txt(slide, text, l + 0.07, t + 0.04, w - 0.1, 0.22,
        size=size, bold=True, color=color)

def pill_row(slide, num_str, color, title, desc, lx, ly, w_total=12.5):
    rect(slide, lx, ly, 0.5, 0.5, fill=color)
    txt(slide, num_str, lx, ly + 0.04, 0.5, 0.42,
        size=14, bold=True, color=SF_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, title, lx + 0.6, ly, w_total - 0.65, 0.28,
        size=11, bold=True, color=SF_DARK_BLUE)
    txt(slide, desc, lx + 0.6, ly + 0.28, w_total - 0.65, 0.28,
        size=9.5, color=SF_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# 1 — CAPA
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=SF_DARK_BLUE)
rect(s, 0, 0, 13.33, 0.1, fill=SF_BLUE)
rect(s, 0, 7.4, 13.33, 0.1, fill=SF_BLUE)
rect(s, 0, 0.1, 0.35, 7.3, fill=SF_BLUE)

# faixa decorativa clara
rect(s, 0.35, 2.1, 9.5, 0.04, fill=SF_BLUE)

txt(s, "SEFIN-CE / DATAPREV", 0.7, 0.5, 11, 0.6,
    size=16, color=SF_LIGHT_BLUE)
txt(s, "Régua de Cobrança Proativa", 0.7, 1.1, 12, 0.8,
    size=38, bold=True, color=SF_WHITE)
txt(s, "via WhatsApp — Agentforce", 0.7, 1.85, 12, 0.6,
    size=28, color=SF_BLUE)

txt(s, "Fluxo de Atendimento v1.0  ·  Premissas Técnicas  ·  Perguntas em Aberto",
    0.7, 2.7, 12, 0.45, size=14, bold=True, color=SF_LIGHT_BLUE)

# chips produtos
chips = [
    (SF_BLUE,   "Agentforce"),
    (SF_TEAL,   "Marketing Cloud"),
    (SF_PURPLE, "Flow / Apex"),
    (SF_GREEN,  "Digital Engagement"),
]
cx = 0.7
for col, label in chips:
    rect(s, cx, 3.3, len(label) * 0.1 + 0.45, 0.34, fill=col)
    txt(s, label, cx + 0.12, 3.34, len(label) * 0.1 + 0.28, 0.28,
        size=10, bold=True, color=SF_WHITE)
    cx += len(label) * 0.1 + 0.6

rect(s, 0.7, 3.85, 5.2, 0.04, fill=SF_GRAY)

txt(s, "Cliente:  DATAPREV para SEFIN-CE — Secretaria Municipal de Finanças de Fortaleza",
    0.7, 4.0, 12, 0.4, size=12, color=SF_LIGHT_BLUE)
txt(s, "Programa:  Modernização da Cobrança Tributária — IPTU · Taxa do Lixo (TMRSU)",
    0.7, 4.4, 12, 0.4, size=12, color=SF_GRAY)

txt(s, "Julho 2026", 0.7, 5.1, 4, 0.4, size=13, color=SF_GRAY)
txt(s, "Salesforce Professional Services LATAM", 0.7, 5.5, 8, 0.35,
    size=11, italic=True, color=SF_GRAY)
txt(s, "Confidencial — Uso Restrito ao Time de Projeto", 0.7, 7.1, 9, 0.3,
    size=9, italic=True, color=SF_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# 2 — AGENDA
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Agenda", "O que vamos apresentar"); footer(s, 2)

agenda = [
    ("01", SF_DARK_BLUE,  "Contexto do Projeto",
     "Volumes, participantes, produtos Salesforce confirmados"),
    ("02", SF_BLUE,       "Solução Técnica",
     "Stack Agentforce + Marketing Cloud + Flow/Apex"),
    ("03", SF_PURPLE,     "Fluxo de Atendimento v1.0",
     "Diagrama completo — todos os nós e decisões"),
    ("04", SF_TEAL,       "Detalhamento do Fluxo",
     "Entrada → Identificação → DAM → ISS → Ajuda → Encerramento"),
    ("05", SF_GREEN,      "Premissas Técnicas",
     "15 premissas consolidadas que regem o escopo"),
    ("06", SF_ORANGE,     "Perguntas em Aberto",
     "11 questões que precisam ser respondidas pelo cliente"),
    ("07", SF_AMBER,      "APIs Mapeadas",
     "EmitirDamUnico · ConsultaImovel · API de Registro (a entregar)"),
    ("08", SF_DARK_GRAY,  "Próximos Passos",
     "Ações, responsáveis e prazos acordados"),
]

cols = [(0.35, 6.0), (7.05, 6.0)]
for i, (num, color, title, desc) in enumerate(agenda):
    col = i % 2; row = i // 2
    lx, w = cols[col]
    ly = 1.15 + row * 1.38
    rect(s, lx, ly, 0.55, 0.55, fill=color)
    txt(s, num, lx, ly + 0.04, 0.55, 0.48,
        size=16, bold=True, color=SF_WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, lx + 0.68, ly, w - 0.75, 0.32,
        size=12, bold=True, color=SF_DARK_BLUE)
    txt(s, desc, lx + 0.68, ly + 0.32, w - 0.75, 0.32,
        size=9.5, color=SF_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# 3 — CONTEXTO
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Contexto do Projeto",
              "SEFIN-CE · Modernização da cobrança tributária via WhatsApp"); footer(s, 3)

# col esquerda — volumes
rect(s, 0.35, 1.12, 5.9, 5.65, fill=SF_WHITE)
rect(s, 0.35, 1.12, 5.9, 0.42, fill=SF_DARK_BLUE)
txt(s, "Volumes e Dados do Projeto", 0.5, 1.15, 5.6, 0.35,
    size=11, bold=True, color=SF_WHITE)

volumes = [
    ("Msgs WhatsApp proativas / ano", "4.860.000"),
    ("Conversas chatbot / ano",        "400.000"),
    ("Posições de atendimento (PAs)",  "10"),
    ("Prazo estimado",                 "5 meses"),
    ("Perfis PS",                      "Dev SC + Dev MC + PM"),
    ("Tributos cobertos",               "IPTU (cód. 10) · TMRSU (cód. 980)"),
    ("ISS",                             "Redirecionamento para site SEFIN-CE"),
    ("Valor estimado total",            "~R$ 1.300.000"),
]
for i, (lbl, val) in enumerate(volumes):
    ly = 1.7 + i * 0.59
    bg_c = SF_LIGHT_GRAY if i % 2 == 0 else SF_WHITE
    rect(s, 0.4, ly, 5.8, 0.55, fill=bg_c)
    txt(s, lbl, 0.52, ly + 0.08, 3.5, 0.38, size=9.5, color=SF_GRAY)
    txt(s, val, 4.0, ly + 0.06, 2.1, 0.42,
        size=10, bold=True, color=SF_DARK_BLUE, align=PP_ALIGN.RIGHT)

# col direita — participantes + produtos
rect(s, 6.85, 1.12, 6.13, 2.85, fill=SF_WHITE)
rect(s, 6.85, 1.12, 6.13, 0.42, fill=SF_BLUE)
txt(s, "Participantes", 7.0, 1.15, 5.9, 0.35,
    size=11, bold=True, color=SF_WHITE)

partic = [
    ("Nelson Stebulaitis Filho", "Solutions Manager — Salesforce PS LATAM"),
    ("Alex Siqueira",            "Account Executive — Salesforce"),
    ("Augusto Cesar Martins",    "Delivery Lead — DATAPREV"),
    ("Osvaldo Melo",             "Arquiteto de Solução — DATAPREV"),
]
for i, (nome, papel) in enumerate(partic):
    ly = 1.68 + i * 0.59
    rect(s, 6.9, ly, 6.03, 0.54, fill=SF_LIGHT_GRAY if i % 2 == 0 else SF_WHITE)
    txt(s, nome, 7.02, ly + 0.05, 3.5, 0.25, size=10, bold=True, color=SF_DARK_BLUE)
    txt(s, papel, 7.02, ly + 0.29, 5.8, 0.22, size=8.5, italic=True, color=SF_GRAY)

rect(s, 6.85, 4.15, 6.13, 2.62, fill=SF_WHITE)
rect(s, 6.85, 4.15, 6.13, 0.42, fill=SF_GREEN)
txt(s, "Produtos Salesforce Confirmados", 7.0, 4.18, 5.9, 0.35,
    size=11, bold=True, color=SF_WHITE)

produtos = [
    (SF_BLUE,   "Service Cloud — Digital Engagement",
     "Canal WhatsApp  ·  10 PAs"),
    (SF_PURPLE, "Agentforce",
     "Agente WhatsApp — IPTU · TMRSU · ISS · Ajuda"),
    (SF_TEAL,   "Marketing Cloud",
     "Journey Builder — 4,86M msgs/ano proativas"),
    (SF_DARK_GRAY, "Flow Orchestration + Apex",
     "Orquestração subfluxos  ·  geração PDF DAM"),
]
for i, (col, nome, desc) in enumerate(produtos):
    ly = 4.72 + i * 0.48
    rect(s, 6.9, ly, 0.12, 0.38, fill=col)
    txt(s, nome, 7.1, ly + 0.01, 3.5, 0.22, size=9.5, bold=True, color=SF_DARK_BLUE)
    txt(s, desc, 7.1, ly + 0.22, 5.8, 0.2, size=8.5, color=SF_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# 4 — SOLUÇÃO TÉCNICA (stack)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Solução Técnica",
              "Arquitetura integrada — Agentforce · Marketing Cloud · Flow/Apex"); footer(s, 4)

# camadas horizontais
camadas = [
    (SF_DARK_BLUE, "CANAIS",
     "WhatsApp (proativo via HSM + reativo) · 10 PAs operadores humanos",
     "Marketing Cloud Journey Builder  ·  Digital Engagement"),
    (SF_BLUE,      "AGENTE IA",
     "Agentforce conectado ao canal WhatsApp — atendimento natural, sem botões rígidos",
     "Menu: IPTU · Taxa Lixo · ISS · Preciso de Ajuda"),
    (SF_PURPLE,    "ORQUESTRAÇÃO",
     "Flow Orchestration + Apex — subfluxos por serviço, chamadas API, geração de PDF",
     "EmitirDamUnico · ConsultaImovel · API Registro de Contato"),
    (SF_TEAL,      "RÉGUA PROATIVA",
     "Marketing Cloud Journey Builder — 4,86M msgs/ano · gatilhos de vencimento ARREC",
     "D-15 · D-5 · D+1 · D+15 · Confirmação pagamento"),
    (SF_GREEN,     "DADOS & INTEGRAÇÃO",
     "APIs SEFIN-CE (ARREC) consumidas via Apex · Customer Data Cloud Starter (plataforma)",
     "Sem MuleSoft · Sem Data Cloud full · Zero-copy onde aplicável"),
]
for i, (col, tag, main, sub) in enumerate(camadas):
    ly = 1.18 + i * 1.12
    rect(s, 0.35, ly, 12.63, 0.98, fill=SF_WHITE)
    rect(s, 0.35, ly, 0.18, 0.98, fill=col)
    rect(s, 0.53, ly, 1.4, 0.98, fill=col)
    txt(s, tag, 0.53, ly + 0.3, 1.4, 0.38,
        size=9, bold=True, color=SF_WHITE, align=PP_ALIGN.CENTER)
    txt(s, main, 2.05, ly + 0.08, 10.8, 0.42,
        size=11, bold=True, color=SF_DARK_BLUE)
    txt(s, sub, 2.05, ly + 0.52, 10.8, 0.38,
        size=9.5, italic=True, color=SF_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# 5 — FLUXO v1.0 (IMAGEM)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=SF_DARK_BLUE)
rect(s, 0, 0, 13.33, 0.08, fill=SF_BLUE)

txt(s, "Fluxo de Atendimento v1.0", 0.4, 0.15, 10, 0.5,
    size=22, bold=True, color=SF_WHITE)
txt(s, "Validado em 02/07/2026  ·  Agentforce SEFIN-CE / DATAPREV",
    0.4, 0.62, 12, 0.35, size=11, italic=True, color=SF_LIGHT_BLUE)

# legenda cores
legend = [
    (RGBColor(0x00, 0x96, 0xFF), "Ação do Agente"),
    (RGBColor(0x2E, 0x7D, 0x32), "Decisão / Condicional"),
    (RGBColor(0x6A, 0x1B, 0x9A), "Chamada de API"),
    (RGBColor(0xE6, 0x51, 0x00), "Falha / Transbordo"),
    (RGBColor(0x00, 0x69, 0x5C), "Pesquisa + Log"),
    (RGBColor(0x7B, 0x1A, 0x1A), "Encerramento"),
]
lx = 0.4
for col, label in legend:
    rect(s, lx, 1.05, 0.22, 0.22, fill=col)
    txt(s, label, lx + 0.28, 1.04, 1.8, 0.25, size=8, color=SF_LIGHT_BLUE)
    lx += 2.15

import os
if os.path.exists(FLUXO_IMG):
    s.shapes.add_picture(FLUXO_IMG,
                         Inches(0.35), Inches(1.35),
                         Inches(12.63), Inches(5.75))
else:
    txt(s, f"[Imagem não encontrada: {FLUXO_IMG}]",
        0.4, 3.5, 12.6, 0.5, size=14, color=SF_ORANGE, align=PP_ALIGN.CENTER)

footer(s, 5)

# ═══════════════════════════════════════════════════════════════════════════════
# 6 — FLUXO — DETALHAMENTO (pt.1): ENTRADA → IDENTIFICAÇÃO → CONSULTA IMÓVEL
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Fluxo — Detalhamento (Parte 1 de 2)",
              "Entrada · Identificação · Consulta Imóvel via API"); footer(s, 6)

sections = [
    (SF_BLUE, "ENTRADA", [
        ("INÍCIO",          "Cidadão envia mensagem ao canal WhatsApp da SEFIN-CE"),
        ("SAUDAÇÃO",        "Agente envia boas-vindas ao canal oficial"),
        ("Veio de HSM?",    "SIM → apresenta contexto do disparo proativo recebido  |  NÃO → Menu direto"),
        ("MENU DE SERVIÇOS","1. Boleto IPTU   2. Boleto Taxa do Lixo   3. ISS   4. Preciso de ajuda / Outros"),
    ]),
    (SF_PURPLE, "IDENTIFICAÇÃO", [
        ("Já identificado?", "SIM + mesmo CPF → vai direto à API  |  NÃO → solicita CPF/CNPJ"),
        ("Validação CPF",    "Formato 11 ou 14 dígitos · até 2 tentativas inválidas · esgotou → Pesquisa → Fim"),
        ("Solicita NOME",    "Nome completo do cidadão"),
        ("Solicita DATA",    "Data de nascimento (PF) ou data de abertura na RFB (PJ) · 2 tentativas"),
    ]),
    (SF_TEAL, "CONSULTA IMÓVEL (API)", [
        ("API ConsultaImovel",   "Método por-documento: CPF/CNPJ + nome + data + exercício fiscal"),
        ("Serviço OK?",          "NÃO → 1 retentativa automática  |  falhou 2x → FALHA COM REGISTRO → ALGO MAIS?"),
        ("Possui inscrições?",   "NÃO → AJUDO COM ALGO MAIS? (cidadão sem imóvel vinculado)"),
        ("Inscrições em aberto?","NÃO → AJUDO COM ALGO MAIS?  |  SIM → apresenta TODAS (sem limite — P-13)"),
    ]),
]

ly_start = 1.13
for col, sec_title, rows in sections:
    # cabeçalho seção
    rect(s, 0.35, ly_start, 12.63, 0.38, fill=col)
    txt(s, f"  {sec_title}", 0.4, ly_start + 0.05, 12.5, 0.3,
        size=11, bold=True, color=SF_WHITE)
    ly_start += 0.4
    for j, (node, desc) in enumerate(rows):
        bg_c = SF_LIGHT_GRAY if j % 2 == 0 else SF_WHITE
        rect(s, 0.35, ly_start, 12.63, 0.48, fill=bg_c)
        rect(s, 0.35, ly_start, 0.06, 0.48, fill=col)
        txt(s, node, 0.52, ly_start + 0.06, 2.3, 0.36, size=9.5, bold=True, color=SF_DARK_BLUE)
        txt(s, desc, 2.9, ly_start + 0.08, 10.0, 0.34, size=9.5, color=SF_DARK_GRAY)
        ly_start += 0.5
    ly_start += 0.1

# ═══════════════════════════════════════════════════════════════════════════════
# 7 — FLUXO — DETALHAMENTO (pt.2): DAM → ISS → AJUDA → ENCERRAMENTO
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Fluxo — Detalhamento (Parte 2 de 2)",
              "Emissão DAM · ISS · Preciso de Ajuda · Encerramento Universal"); footer(s, 7)

sections2 = [
    (SF_GREEN, "EMISSÃO DAM — IPTU / TAXA DO LIXO", [
        ("Cidadão seleciona inscrição",
         "Escolhe uma das inscrições em aberto apresentadas"),
        ("DAM só ano vigente?",
         "SIM → confirma emissão  |  NÃO → escolhe Cota Única ou Parcelamento + ano"),
        ("CONFIRMA EMISSÃO?",
         "NÃO → AJUDO COM ALGO MAIS?  |  SIM → chama API EmitirDamUnico"),
        ("API EmitirDamUnico",
         "CPF/CNPJ + tipoDebito + tipoPagamento + ano/parcelas · retorna link PDF DAM"),
        ("Falha API",
         "1 retentativa · se falhar 2x → Informa falha + registra + AJUDO COM ALGO MAIS? (SEM pesquisa)"),
        ("Loop inscrições",
         "Possui outra inscrição sem DAM? → Quer emitir? SIM = repete fluxo completo · NÃO = ALGO MAIS?"),
    ]),
    (SF_AMBER, "ISS  ·  PRECISO DE AJUDA", [
        ("ISS",
         "Informa link do site SEFIN-CE · sem identificação do cidadão [Q-E: URL a confirmar]"),
        ("Preciso de Ajuda",
         "Coleta CPF/CNPJ + nome → consulta Knowledge Base [Q-F: KB a definir]"),
        ("KB respondeu?",
         "SIM → entrega resposta → AJUDO COM ALGO MAIS?  |  NÃO → verifica DHA"),
        ("DHA",
         "FORA horário → informa indisponibilidade → AJUDO COM ALGO MAIS?  |  DENTRO → [TRANSBORDO — Q-B]"),
    ]),
    (SF_TEAL, "ENCERRAMENTO UNIVERSAL", [
        ("AJUDO COM ALGO MAIS?",
         "SIM → volta ao Menu  |  NÃO → Pesquisa de Satisfação"),
        ("PESQUISA DE SATISFAÇÃO",
         "⭐ Muito ruim  ⭐⭐ Ruim  ⭐⭐⭐ Razoável  ⭐⭐⭐⭐ Bom  ⭐⭐⭐⭐⭐ Muito bom  —  ALL encerramentos (P-05)"),
        ("Nota ≤ 3?",
         "SIM → solicita descrição da insatisfação antes de registrar"),
        ("API Registro de Contato",
         "CPF/CNPJ · nome · serviço · resultado · nota · justificativa  [Q-C / Q-D — a definir]"),
    ]),
]

ly_s = 1.13
for col, sec_title, rows in sections2:
    rect(s, 0.35, ly_s, 12.63, 0.38, fill=col)
    txt(s, f"  {sec_title}", 0.4, ly_s + 0.05, 12.5, 0.3,
        size=11, bold=True, color=SF_WHITE)
    ly_s += 0.4
    for j, (node, desc) in enumerate(rows):
        bg_c = SF_LIGHT_GRAY if j % 2 == 0 else SF_WHITE
        rect(s, 0.35, ly_s, 12.63, 0.48, fill=bg_c)
        rect(s, 0.35, ly_s, 0.06, 0.48, fill=col)
        txt(s, node, 0.52, ly_s + 0.06, 2.3, 0.36, size=9.5, bold=True, color=SF_DARK_BLUE)
        txt(s, desc, 2.9, ly_s + 0.08, 10.0, 0.34, size=9.5, color=SF_DARK_GRAY)
        ly_s += 0.5
    ly_s += 0.08

# ═══════════════════════════════════════════════════════════════════════════════
# 8 — PREMISSAS (P-01 a P-08)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Premissas Técnicas  (1 de 2)",
              "Premissas P-01 a P-08 — regem o escopo e o desenvolvimento"); footer(s, 8)

premissas_1 = [
    ("P-01", SF_TEAL,
     "Encerramento universal com Pesquisa de Satisfação",
     "Todo encerramento → Pesquisa → API Registro → Despedida. "
     "EXCEÇÃO: falha de API após retentativa encerra sem pesquisa."),
    ("P-02", SF_ORANGE,
     "Transbordo humano fora do escopo (a definir)",
     "Não há mecanismo de transbordo humano confirmado. DHA dentro do horário sinaliza "
     "[TRANSBORDO — A DEFINIR] como nó orange no fluxo."),
    ("P-03", SF_BLUE,
     "DHA fora do horário = informa indisponibilidade",
     "Quando DHA = FORA, agente informa que atendimento humano só ocorre em horário comercial "
     "[Q-J: horário a confirmar] e encerra com pesquisa normalmente."),
    ("P-04", SF_GREEN,
     "API de Registro recebe dados completos",
     "Campos: CPF/CNPJ · nome · serviço executado · resultado · nota de satisfação · "
     "justificativa de insatisfação (quando ≤ 3 estrelas)."),
    ("P-05", SF_PURPLE,
     "Pesquisa de satisfação em TODOS os encerramentos",
     "Incluindo erros, abandonos e timeouts. EXCEÇÃO única: falha persistente de API "
     "(2ª tentativa falhou) → encerra sem pesquisa."),
    ("P-06", SF_TEAL,
     "Nota ≤ 3 estrelas = coleta justificativa",
     "Antes de chamar a API de Registro, o agente solicita uma descrição textual "
     "do motivo da insatisfação."),
    ("P-07", SF_DARK_BLUE,
     "Máximo 2 tentativas por campo de entrada",
     "CPF/CNPJ e Data permitem 2 tentativas inválidas. Na 3ª falha, o fluxo "
     "encerra com pesquisa de satisfação e registro."),
    ("P-08", SF_GREEN,
     "Sessão já identificada = pula identificação",
     "Cidadão identificado com o mesmo CPF/CNPJ na mesma sessão vai direto "
     "para a API ConsultaImovel sem repetir nome/data."),
]

for i, (code, col, title, desc) in enumerate(premissas_1):
    ly = 1.12 + i * 0.73
    rect(s, 0.35, ly, 0.65, 0.62, fill=col)
    txt(s, code, 0.35, ly + 0.12, 0.65, 0.38,
        size=10, bold=True, color=SF_WHITE, align=PP_ALIGN.CENTER)
    rect(s, 1.02, ly, 12.0, 0.62, fill=SF_WHITE)
    txt(s, title, 1.12, ly + 0.03, 11.8, 0.28, size=10.5, bold=True, color=SF_DARK_BLUE)
    txt(s, desc, 1.12, ly + 0.32, 11.8, 0.28, size=9, color=SF_DARK_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# 9 — PREMISSAS (P-09 a P-15)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Premissas Técnicas  (2 de 2)",
              "Premissas P-09 a P-15 — regem o escopo e o desenvolvimento"); footer(s, 9)

premissas_2 = [
    ("P-09", SF_AMBER,
     "ISS = apenas link do site, sem identificação do cidadão",
     "Ao escolher ISS no menu, o agente informa o link do site SEFIN-CE "
     "[Q-E: URL a confirmar] e segue para AJUDO COM ALGO MAIS?"),
    ("P-10", SF_ORANGE,
     "Falha de API = 1 retentativa, depois informa e registra",
     "Após 2 falhas consecutivas: agente informa a falha, diz que o contato foi registrado "
     "e que alguém entrará em contato para resolver. Depois: AJUDO COM ALGO MAIS?"),
    ("P-11", SF_BLUE,
     "CONFIRMA EMISSÃO? = NÃO → AJUDO COM ALGO MAIS? sem chamar API",
     "Se o cidadão cancelar a confirmação, nenhuma chamada é feita à API EmitirDamUnico. "
     "Fluxo segue direto para AJUDO COM ALGO MAIS?"),
    ("P-12", SF_PURPLE,
     "Loop de emissão para múltiplas inscrições",
     "O processo de seleção + confirmação + API EmitirDamUnico repete-se individualmente "
     "para cada inscrição que o cidadão desejar emitir."),
    ("P-13", SF_DARK_BLUE,
     "Sem limite de inscrições — apresenta todas",
     "O fluxo original limitava a 7 inscrições. Esta premissa ignora esse limite: "
     "o agente apresenta TODAS as inscrições retornadas pela API. [Q-K: verificar motivo original]"),
    ("P-14", SF_TEAL,
     "Régua proativa via Marketing Cloud Journey Builder",
     "Sugestão de régua: D-15, D-5, D+1, D+15. Gatilhos alimentados pelo sistema ARREC. "
     "A ser validada com o cliente [Q-G / Q-H]."),
    ("P-15", SF_GREEN,
     "'Preciso de Ajuda' → KB → se sem resposta → [TRANSBORDO]",
     "Identifica o cidadão (CPF/nome), consulta Knowledge Base [Q-F]. "
     "Se KB não responder, verifica DHA e sinaliza [TRANSBORDO — A DEFINIR]."),
]

for i, (code, col, title, desc) in enumerate(premissas_2):
    ly = 1.12 + i * 0.84
    rect(s, 0.35, ly, 0.65, 0.73, fill=col)
    txt(s, code, 0.35, ly + 0.18, 0.65, 0.38,
        size=10, bold=True, color=SF_WHITE, align=PP_ALIGN.CENTER)
    rect(s, 1.02, ly, 12.0, 0.73, fill=SF_WHITE)
    txt(s, title, 1.12, ly + 0.05, 11.8, 0.28, size=10.5, bold=True, color=SF_DARK_BLUE)
    txt(s, desc, 1.12, ly + 0.38, 11.8, 0.3, size=9, color=SF_DARK_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# 10 — PERGUNTAS EM ABERTO (Q-A a Q-F)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Perguntas em Aberto para o Cliente  (1 de 2)",
              "Q-A a Q-F — respostas necessárias antes de iniciar o desenvolvimento"); footer(s, 10)

perguntas_1 = [
    ("Q-A", SF_ORANGE,
     "O que é HSM no contexto do fluxo original?",
     "Quais parâmetros e conteúdo das mensagens HSM proativas? Como o cidadão é identificado no disparo?",
     "Osvaldo / Augusto"),
    ("Q-B", SF_ORANGE,
     "Qual será o mecanismo de transbordo humano?",
     "Como se dará tecnicamente dado que não há Service Cloud/Digital Engagement no escopo? "
     "Há queue externa, telefone, outro canal?",
     "Augusto / SEFIN-CE"),
    ("Q-C", SF_DARK_BLUE,
     "Quais são os parâmetros da API de Registro de Contato?",
     "Endpoint, autenticação, campos esperados. Confirmar se recebe: CPF/CNPJ, nome, "
     "serviço, resultado, nota, justificativa.",
     "Osvaldo"),
    ("Q-D", SF_DARK_BLUE,
     "O registro da pesquisa de satisfação usa a mesma API de log?",
     "Ou há endpoint separado para avaliação? A nota e justificativa fazem parte do mesmo "
     "payload ou chamada distinta?",
     "Osvaldo"),
    ("Q-E", SF_BLUE,
     "Qual é a URL do site SEFIN-CE para o fluxo ISS?",
     "URL completa que o agente deve informar ao cidadão quando selecionar ISS no menu.",
     "SEFIN-CE"),
    ("Q-F", SF_BLUE,
     "Há base de conhecimento (KB) disponível para o agente?",
     "Se sim: quem mantém, em qual formato (PDF, artigos, FAQ), qual sistema atual? "
     "Agente consulta KB para 'Preciso de Ajuda' antes de escalar.",
     "SEFIN-CE / Augusto"),
]

for i, (code, col, pergunta, detalhe, resp) in enumerate(perguntas_1):
    ly = 1.12 + i * 0.98
    # número badge
    rect(s, 0.35, ly, 0.68, 0.85, fill=col)
    txt(s, code, 0.35, ly + 0.24, 0.68, 0.38,
        size=11, bold=True, color=SF_WHITE, align=PP_ALIGN.CENTER)
    # corpo
    rect(s, 1.05, ly, 12.0, 0.85, fill=SF_WHITE)
    txt(s, pergunta, 1.15, ly + 0.04, 9.5, 0.3, size=10.5, bold=True, color=SF_DARK_BLUE)
    txt(s, detalhe, 1.15, ly + 0.35, 9.5, 0.38, size=9, color=SF_DARK_GRAY)
    # responsável
    rect(s, 10.68, ly + 0.12, 2.25, 0.28, fill=SF_LIGHT_GRAY)
    txt(s, resp, 10.76, ly + 0.14, 2.1, 0.22, size=8, color=SF_GRAY, bold=True)

# ═══════════════════════════════════════════════════════════════════════════════
# 11 — PERGUNTAS EM ABERTO (Q-G a Q-K)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Perguntas em Aberto para o Cliente  (2 de 2)",
              "Q-G a Q-K — respostas necessárias antes de iniciar o desenvolvimento"); footer(s, 11)

perguntas_2 = [
    ("Q-G", SF_TEAL,
     "Há régua de comunicação proativa já definida?",
     "Se não, o cliente aceita adotar a régua sugerida via Marketing Cloud Journey Builder "
     "(D-15, D-5, D+1, D+15)? Quem aprova o conteúdo das mensagens?",
     "SEFIN-CE"),
    ("Q-H", SF_TEAL,
     "Qual é a fonte de dados para os gatilhos da régua proativa?",
     "Como o Marketing Cloud acessa os dados de quem está em débito e quando vence? "
     "API ARREC, SFTP batch, Direct Connect ou outro?",
     "SEFIN-CE / TI"),
    ("Q-I", SF_GREEN,
     "O cidadão pode solicitar opt-out das mensagens proativas?",
     "Se sim: como é registrado o opt-out e como o Journey Builder respeita essa preferência "
     "nos disparos seguintes? Há campo no ARREC?",
     "SEFIN-CE"),
    ("Q-J", SF_PURPLE,
     "Qual é o horário comercial de atendimento humano?",
     "Dias da semana (ex.: seg-sex), horário de início e fim (ex.: 8h–18h). "
     "Feriados nacionais e municipais de Fortaleza devem ser considerados?",
     "SEFIN-CE"),
    ("Q-K", SF_DARK_GRAY,
     "O limite de 7 inscrições do fluxo original é limitação de qual sistema?",
     "Bot atual, API ConsultaImovel ou regra de negócio? A nova solução pode apresentar "
     "todas as inscrições sem limite (Premissa P-13 assume que SIM)?",
     "Osvaldo / SEFIN-CE"),
]

for i, (code, col, pergunta, detalhe, resp) in enumerate(perguntas_2):
    ly = 1.12 + i * 1.12
    rect(s, 0.35, ly, 0.68, 0.98, fill=col)
    txt(s, code, 0.35, ly + 0.28, 0.68, 0.38,
        size=11, bold=True, color=SF_WHITE, align=PP_ALIGN.CENTER)
    rect(s, 1.05, ly, 12.0, 0.98, fill=SF_WHITE)
    txt(s, pergunta, 1.15, ly + 0.06, 9.5, 0.3, size=10.5, bold=True, color=SF_DARK_BLUE)
    txt(s, detalhe, 1.15, ly + 0.4, 9.5, 0.46, size=9, color=SF_DARK_GRAY)
    rect(s, 10.68, ly + 0.14, 2.25, 0.28, fill=SF_LIGHT_GRAY)
    txt(s, resp, 10.76, ly + 0.16, 2.1, 0.22, size=8, color=SF_GRAY, bold=True)

# ═══════════════════════════════════════════════════════════════════════════════
# 12 — APIs MAPEADAS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "APIs Mapeadas — Sistema ARREC (SEFIN-CE)",
              "Dois serviços prontos para consumo · API de Registro a ser entregue pelo cliente"); footer(s, 12)

# API 1 — EmitirDamUnico
rect(s, 0.35, 1.12, 12.63, 0.38, fill=SF_PURPLE)
txt(s, "  API 1 — EmitirDamUnico  ·  Gera boleto DAM para pagamento de tributo",
    0.4, 1.15, 12.5, 0.3, size=11, bold=True, color=SF_WHITE)

params_dam = [
    ("tipoDebito",      "Integer(3)", "S", "TMRSU = 980  ·  IPTU = 10"),
    ("tipoPessoa",      "String(1)",  "S", "F = Física  ·  J = Jurídica"),
    ("cpfcnpj",         "String(14)", "S", "Somente números"),
    ("tipoPagamento",   "String(9)",  "S", "COTAUNICA  ou  PARCELADO"),
    ("AnoDebito",       "Integer(4)", "N", "Ano do débito — opcional"),
    ("periodoParcelas", "Date",       "N", "Obrigatório se PARCELADO"),
]
header_row_bg = SF_DARK_BLUE
col_ws = [3.5, 2.2, 1.0, 5.8]
lx_s = [0.35, 3.85, 6.05, 7.05]
row_labels = ["Parâmetro", "Tipo", "Obrig.", "Descrição"]
ly = 1.52
rect(s, 0.35, ly, 12.63, 0.34, fill=header_row_bg)
for ci, (lbl, w) in enumerate(zip(row_labels, col_ws)):
    txt(s, lbl, lx_s[ci] + 0.06, ly + 0.05, w - 0.08, 0.25,
        size=9, bold=True, color=SF_WHITE)
ly += 0.34
for ri, (p, typ, obl, desc) in enumerate(params_dam):
    bg_c = SF_LIGHT_GRAY if ri % 2 == 0 else SF_WHITE
    rect(s, 0.35, ly, 12.63, 0.38, fill=bg_c)
    vals = [p, typ, obl, desc]
    for ci, (val, w) in enumerate(zip(vals, col_ws)):
        c = SF_PURPLE if ci == 0 else SF_DARK_GRAY
        txt(s, val, lx_s[ci] + 0.06, ly + 0.07, w - 0.08, 0.28, size=9, bold=(ci==0), color=c)
    ly += 0.38

txt(s, "  Retorno: link PDF do DAM pronto para download e envio via WhatsApp",
    0.4, ly + 0.06, 12.5, 0.28, size=9.5, italic=True, color=SF_TEAL, bold=True)

# API 2 — ConsultaImovel
ly += 0.42
rect(s, 0.35, ly, 12.63, 0.38, fill=SF_TEAL)
txt(s, "  API 2 — ConsultaImovel por-documento  ·  Retorna imóveis e inscrições vinculadas ao CPF/CNPJ",
    0.4, ly + 0.05, 12.5, 0.28, size=11, bold=True, color=SF_WHITE)

params_cons = [
    ("Tipo",      "String(8)",    "S", "Física  ou  Jurídica"),
    ("Documento", "String(14)",   "S", "CPF ou CNPJ formatado"),
    ("Data",      "Date ISO8601", "S", "Data de nascimento (PF) ou abertura na RFB (PJ)"),
    ("Exercicio", "Number(4)",    "S", "Exercício fiscal a pesquisar (ano)"),
]
ly += 0.38
rect(s, 0.35, ly, 12.63, 0.3, fill=header_row_bg)
for ci, (lbl, w) in enumerate(zip(row_labels, col_ws)):
    txt(s, lbl, lx_s[ci] + 0.06, ly + 0.04, w - 0.08, 0.22,
        size=9, bold=True, color=SF_WHITE)
ly += 0.3
for ri, (p, typ, obl, desc) in enumerate(params_cons):
    bg_c = SF_LIGHT_GRAY if ri % 2 == 0 else SF_WHITE
    rect(s, 0.35, ly, 12.63, 0.36, fill=bg_c)
    vals = [p, typ, obl, desc]
    for ci, (val, w) in enumerate(zip(vals, col_ws)):
        c = SF_TEAL if ci == 0 else SF_DARK_GRAY
        txt(s, val, lx_s[ci] + 0.06, ly + 0.07, w - 0.08, 0.25, size=9, bold=(ci==0), color=c)
    ly += 0.36

txt(s, "  Retorno: inscrição municipal · endereço · cartografia · titular",
    0.4, ly + 0.06, 8, 0.25, size=9.5, italic=True, color=SF_TEAL, bold=True)

# API 3 — Registro (a entregar)
ly += 0.38
rect(s, 0.35, ly, 12.63, 0.34, fill=SF_ORANGE)
txt(s, "  API 3 — Registro de Contato  ·  A ser entregue pelo cliente [Q-C / Q-D]  "
       "·  Campos esperados: CPF/CNPJ · nome · serviço · resultado · nota · justificativa",
    0.4, ly + 0.04, 12.4, 0.26, size=9.5, bold=True, color=SF_WHITE)

# ═══════════════════════════════════════════════════════════════════════════════
# 13 — PRÓXIMOS PASSOS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Próximos Passos",
              "Ações, responsáveis e prazos para avançar com o projeto"); footer(s, 13)

passos = [
    (SF_GREEN,     "✅", "Osvaldo entregou documentação das APIs",
     "01/07/2026", "Concluído", "Osvaldo Melo"),
    (SF_GREEN,     "✅", "Fluxo de atendimento v1.0 validado",
     "02/07/2026", "Concluído", "Nelson"),
    (SF_ORANGE,    "🔴", "Nelson + Alex revisam escopo de 800h e valor PS real",
     "URGENTE", "Em andamento", "Nelson + Alex Siqueira"),
    (SF_ORANGE,    "🔴", "Alex comunica retorno a Augusto com valor e escopo revisado",
     "URGENTE", "Aguardando #3", "Alex Siqueira"),
    (SF_BLUE,      "⏳", "Enviar perguntas Q-A a Q-K para DATAPREV / SEFIN-CE",
     "Próxima semana", "Pendente", "Nelson / Augusto"),
    (SF_BLUE,      "⏳", "Cliente responde perguntas — validar premissas restantes",
     "A combinar", "Pendente", "Osvaldo / SEFIN-CE"),
    (SF_BLUE,      "⏳", "Preparar ROM definitivo para DATAPREV fechar com CFIN",
     "Pós validação", "Pendente", "Nelson"),
    (SF_DARK_GRAY, "🔭", "Avaliar expansão para PGM e CMF (mesmo modelo)",
     "Futuro", "Oportunidade", "Augusto + Alex"),
]

col_ws2 = [5.5, 1.7, 1.8, 2.8]
lx2 = [1.1, 6.65, 8.4, 10.25]
ly = 1.12
# header
rect(s, 0.35, ly, 12.63, 0.34, fill=SF_DARK_BLUE)
for ci, (lbl, w) in enumerate(zip(["Ação", "Prazo", "Status", "Responsável"], col_ws2)):
    txt(s, lbl, lx2[ci] + 0.05, ly + 0.05, w - 0.08, 0.26,
        size=9, bold=True, color=SF_WHITE)
ly += 0.34

for ri, (col, icon, acao, prazo, status, resp) in enumerate(passos):
    bg_c = SF_LIGHT_GRAY if ri % 2 == 0 else SF_WHITE
    rect(s, 0.35, ly, 12.63, 0.62, fill=bg_c)
    # barra colorida lateral
    rect(s, 0.35, ly, 0.12, 0.62, fill=col)
    # ícone / badge
    rect(s, 0.5, ly + 0.12, 0.55, 0.38, fill=col)
    txt(s, icon, 0.5, ly + 0.13, 0.55, 0.35,
        size=12, color=SF_WHITE, align=PP_ALIGN.CENTER)
    # ação
    txt(s, acao, lx2[0], ly + 0.14, col_ws2[0] - 0.05, 0.34,
        size=10, bold=True, color=SF_DARK_BLUE)
    # prazo
    txt(s, prazo, lx2[1], ly + 0.17, col_ws2[1] - 0.05, 0.28,
        size=9, color=SF_GRAY)
    # status badge
    s_col = SF_GREEN if status == "Concluído" else \
            SF_ORANGE if "andamento" in status or "URGENTE" == prazo else \
            SF_BLUE if "Pendente" in status or "Aguardando" in status else SF_DARK_GRAY
    rect(s, lx2[2], ly + 0.17, col_ws2[2] - 0.1, 0.28, fill=s_col)
    txt(s, status, lx2[2] + 0.07, ly + 0.19, col_ws2[2] - 0.2, 0.22,
        size=8, bold=True, color=SF_WHITE)
    # responsável
    txt(s, resp, lx2[3], ly + 0.17, col_ws2[3] - 0.05, 0.28,
        size=9, bold=True, color=SF_DARK_BLUE)
    ly += 0.64

prs.save(OUTPUT)
import os
print(f"PPTX gerado: {OUTPUT}")
print(f"Tamanho: {os.path.getsize(OUTPUT)/1024/1024:.1f} MB")
print(f"Total de slides: {TOTAL}")
