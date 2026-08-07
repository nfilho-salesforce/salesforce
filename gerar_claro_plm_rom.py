# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

src = '/Users/nfilho/Downloads/URT_UNIFIED_ROM_TEMPLATE_PT (3).pptx'
out = '/Users/nfilho/claude/CLARO_Agente_PLM_ROM.pptx'

prs = Presentation(src)

# ── Helpers ───────────────────────────────────────────────────────────────────
def replace_text(shape, old, new):
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)

def replace_all(slide, mapping):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for old, new in mapping.items():
                replace_text(shape, old, new)

def set_cell(tbl, row, col, text, bold=False, size=8):
    cell = tbl.cell(row, col)
    cell.text = ''
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)

def set_shape_text(shape, lines, size=9, bold_first=False):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            r.text = ''
    first = True
    for i, line in enumerate(lines):
        if i == 0 and tf.paragraphs:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = (bold_first and first)
        first = False

# ── SLIDE 1 — Capa ────────────────────────────────────────────────────────────
sl = prs.slides[0]
replace_all(sl, {
    'Inserir_Papel_Membro_Equipe_Salesforce_EM': 'Engagement Manager',
    'Inserir_Papel_Membro_Equipe_Salesforce_SBS': 'Technical Architect',
    'Inserir_Papel_Membro_Equipe_Salesforce_AP': 'Technical Consultant / QA',
    'Inserir_Papel_Membro_Equipe_Salesforce_Sponsor': 'Account Executive',
    'INSERIR MÊS E ANO': 'Junho 2026',
    '(Inserir Nome do Cliente)': 'Claro Brasil',
    'XXX Cloud': 'Agentforce + Einstein 1 Platform',
    'XXX + Implementação YYYY': 'POC Agente PLM — Validacao de Catalogo Autonoma',
})

# ── SLIDE 7 — Missao ──────────────────────────────────────────────────────────
sl = prs.slides[6]
replace_all(sl, {
    'Entendemos que sua missão principal é (Inserir Missão do Cliente)...':
    'Entendemos que a missao estrategica da Claro nesta iniciativa e posicionar-se na vanguarda do setor de telecomunicacoes por meio de um catalogo de produtos unificado, inteligente e autonomo — que erradica a complexidade das 127 regras manuais legadas do motor BRE e automatiza os fluxos de validacao do produto ao dinheiro sem dependencia de deploys manuais, liberando as equipes de TI e Negocios para inovar com agilidade e seguranca.'
})

# ── SLIDE 8 — Motores ────────────────────────────────────────────────────────
sl = prs.slides[7]
replace_all(sl, {
    '(Inserir Motor Principal #1)':
    'Excelencia Operacional: Eliminar redeploys para atualizacao de regras de catalogo — migrando de dias de espera tecnica para autoria self-service via IA em 0 minutos (KPI-TI-002).',
    '(Inserir Motor Principal #2)':
    'Mitigacao de Riscos de Infraestrutura: Resolver estouros de CPU e Heap Memory em validacoes massivas de CSV, suportando lotes de ate 10.000 linhas de forma async-first resiliente (KPI-OPS-001).',
    '(Inserir Motor Principal #3)':
    'Inovacao com IA Cognitiva: Introduzir agentes autonomos Agentforce (Admin e Ops) no back-office corporativo, substituindo 127 regras rigidas do catalogo BRE legado por especificacoes declarativas compiladas pelo Einstein LLM.',
})

# ── SLIDE 9 — Objetivos (tabela) ──────────────────────────────────────────────
sl = prs.slides[8]
for shape in sl.shapes:
    if shape.has_table:
        tbl = shape.table
        data = [
            ('Reduzir o tempo de avaliacao de regras de catalogo para menos de 50ms por registro',
             'Implementar o interpretador Apex AST Walker (PlmRuleSpecEvaluator) executando arvores logicas diretamente em memoria, contornando a ausencia de Platform Cache.'),
            ('Eliminar redeploys para atualizacao de regras — meta: 0 minutos de downtime',
             'Agente Admin via Agentforce converte regras em linguagem natural (DSL pt-BR) para especificacoes JSON ativas via Einstein Prompt Templates, sem pipeline de deploy.'),
            ('Suportar ingestao de ate 10.000 linhas por lote sincrono e strings de ate 6M caracteres',
             'Arquitetura Queueable encadeada com cursor byte-offset resumivel fragmenta o CSV de forma assincrona, prevenindo estouro de heap e timeout de CPU.'),
            ('Garantir resiliencia total em falhas de processamento batch com cobertura 100% via DLQ',
             'Transaction Finalizers + Dead Letter Queues + PimZombieReaperService capturam e recuperam falhas sem perda de dados ou transacoes zumbis.'),
        ]
        for r, (obj, how) in enumerate(data, 1):
            if r < len(tbl.rows):
                set_cell(tbl, r, 0, obj, bold=True, size=8)
                set_cell(tbl, r, 1, how, size=8)
        break

# ── SLIDE 10 — Desafios ───────────────────────────────────────────────────────
sl = prs.slides[9]
for shape in sl.shapes:
    if shape.has_text_frame and 'Dados de clientes fragmentados' in shape.text_frame.text:
        challenges = [
            'RIGIDEZ DO CATALOGO BRE LEGADO: 127 regras acopladas em codigo rigido exigem deploys manuais para qualquer alteracao — alta dependencia de TI e lentidao na inovacao de ofertas.',
            'SOBRECARGA DE CPU E HEAP MEMORY: Processamento sincrono de layouts CSV volumosos causa travamentos recorrentes nas orgs STORM_PLM e Ibuy, sem opcao de Platform Cache.',
            'DEPLOYS CROSS-ORG FRAGEIS: Erros sistematicos de CannotQuickDeployError bloqueiam entregas e exigem execucao manual de rotinas especificas de teste (RunSpecifiedTests).',
            'BASE DE CONHECIMENTO DESORGANIZADA: Knowledge Articles sem Data Categories adequadas comprometem a precisao das respostas do Agente de IA.',
            'AUSENCIA DE GERENTE DE PROJETO DEDICADO: Papel de PM compartilhado com time SWE (Luciano) e risco de gargalos em tomadas de decisao nos marcos criticos.',
        ]
        set_shape_text(shape, challenges, size=8)
        break

# ── SLIDE 13 — Visao Futuro ───────────────────────────────────────────────────
sl = prs.slides[12]
replace_all(sl, {
    'Imagine um futuro onde cada equipe é empoderada\ncom os dados certos no momento certo…':
    'Imagine um futuro onde um analista de produtos da Claro descreve em portugues uma nova restricao promocional diretamente no chat do Agente Admin — e em segundos a regra e compilada, validada e ativa. Zero tickets de TI, zero deploys.',
    'Onde cada interação com o cliente é inteligente,\npersonalizada e fluida...':
    'Onde o motor de validacao de catalogo processa lotes de 10.000 linhas de CSV em menos de 50 milissegundos por registro — sem estouro de CPU, sem travamentos, com relatorio narrativo de diagnostico gerado automaticamente pelo Agente Ops.',
    'Onde a tecnologia não é mais uma barreira,\nmas o motor para seu crescimento e inovação…':
    'Onde as 127 regras rigidas do catalogo BRE legado foram substituidas por especificacoes JSON declarativas gerenciadas pelo CoE — com visibilidade total de auditoria, resiliencia ativa via DLQ e governanca autonoma de ciclo de vida de produtos.',
})

# ── SLIDE 14 — Valor (tabela) ─────────────────────────────────────────────────
sl = prs.slides[13]
for shape in sl.shapes:
    if shape.has_table:
        tbl = shape.table
        data = [
            ('Agentforce Employee Agents (Admin + Ops)',
             'Autoria self-service de regras via linguagem natural e diagnostico automatico de lotes — eliminando dependencia de TI para atualizacoes de catalogo.',
             'KPI-TI-002: Tempo de atualizacao de regras -> 0 minutos (sem redeploy)'),
            ('Apex AST Walker Engine (PlmRuleSpecEvaluator)',
             'Interpretador determinisico de arvores logicas JSON em memoria pura — contornando a ausencia de Platform Cache nas orgs Claro.',
             'KPI-TI-001: Avaliacao de regras < 50ms por registro'),
            ('Async Enqueueable Framework + Cursor Byte-Offset',
             'Fragmentacao inteligente de CSVs volumosos via Queueables encadeados — prevenindo estouros de heap e CPU em cargas massivas.',
             'KPI-OPS-001: Suporte a lotes de ate 10.000 linhas sincronas / 6M caracteres'),
            ('Transaction Finalizers + Dead Letter Queue + ZombieReaper',
             'Captura e recuperacao 100% de falhas asincronas — eliminando transacoes zumbis e garantindo rastreabilidade completa.',
             '100% de cobertura de erros via DLQ — zero falhas ocultas'),
        ]
        for r, (cap, impact, kpi) in enumerate(data, 1):
            if r < len(tbl.rows):
                set_cell(tbl, r, 0, cap, bold=True, size=8)
                set_cell(tbl, r, 1, impact, size=8)
                set_cell(tbl, r, 2, kpi, size=8)
        break

# ── SLIDE 17 — Arquitetura ────────────────────────────────────────────────────
sl = prs.slides[16]
replace_all(sl, {
    'Inserir MAPA da Arquitetura':
    'CAMADA DE EXPERIENCIA\nLWC Wizard (DemandaCsvWizardController) | Chat Panels Agentforce\n\n|\n\nCAMADA DE PROCESSO E INTELIGENCIA\nAtlas Reasoning Engine | Queueable Enqueueable Engine | Agente Admin (compile-time) | Agente Ops (runtime)\n\n|\n\nCAMADA DE DADOS\nDemanda__c + Item_Demanda__c | Apex Static Cache Maps | Custom Metadata JSON Specs\n\n|\n\nCAMADA DE INTEGRACAO E GenAI\nConnectApi | Einstein Prompt Templates | Einstein Trust Layer\n\n|\n\nSISTEMAS EXTERNOS\nService Cloud Knowledge Base | Sistema NBO Externo (Agente 3 - condicional)\n\n|\n\nRESILIENCIA CORE\nDead Letter Queue | Transaction Finalizers | PimZombieReaperService'
})

# ── SLIDE 18 — Capacidades Salesforce ────────────────────────────────────────
sl = prs.slides[17]
replace_all(sl, {
    'Exemplo: Data Cloud: Para unificar todos os dados de clientes em uma única fonte de verdade.\nExemplo: Sales Cloud: Para automatizar o processo de lead a fechamento.\nExemplo: Service Cloud: Para oferecer um atendimento ao cliente personalizado e multicanal.\nExemplo: MuleSoft: Para integrar perfeitamente com seus sistemas existentes como.':
    'AGENTFORCE EMPLOYEE AGENTS: Dois agentes autonomos — Admin (compile-time: compilar/recompilar regras e verificar status via chat) e Ops (runtime: importar CSV, avaliar lotes, disparar diagnosticos HTML). Elimina dependencia de TI para atualizacoes de catalogo.\n\nEINSTEIN PROMPT TEMPLATES + ConnectApi: Compilacao automatica de regras escritas em portugues (DSL pt-BR) para especificacoes AST JSON estruturadas via Einstein Trust Layer. Zero redeploys para novas regras.\n\nAPEX AST WALKER (PlmRuleSpecEvaluator): Interpretador deterministico puro-Apex que percorre arvores logicas JSON em memoria com tempo < 50ms por registro — contornando a proibicao de Platform Cache nas orgs STORM_PLM e Ibuy.\n\nASYNC ENQUEUEABLE FRAMEWORK: Queueables encadeados com cursor byte-offset e CPU guard processam CSVs de ate 6M caracteres e 10.000 linhas sem estouro de heap ou timeout sincrono.\n\nRESILIENCIA ATIVA: Transaction Finalizers + Dead Letter Queues + PimZombieReaperService garantem captura e recuperacao 100% de falhas em execucoes batch, com observabilidade por snapshots e Platform Events.'
})

# ── SLIDE 19 — Decisoes Arquitetonicas ───────────────────────────────────────
sl = prs.slides[18]
replace_all(sl, {
    'Decisão: Usar o MuleSoft como uma plataforma de integração estratégica em vez de desenvolver integrações ponto a ponto. Justificativa: A complexidade de conectar-se a 16 DMSs heterogêneos mais um ecossistema de sistemas satélites torna a abordagem de integração ponto a ponto insustentável, frágil e cara de manter. Uma abordagem de conectividade liderada por API com o MuleSoft é um requisito técnico inegociável para gerenciar esse nível de complexidade, garantir governança de dados centralizada e possibilitar a agilidade futura requerida pela visão "Evolutiva" do projeto. Esta decisão transforma integrações de um passivo técnico para um ativo estratégico reutilizável.':
    'DECISAO 1 — AST Walker Apex puro (vs. query relacional)\nJustificativa: A proibicao de Platform Cache nas orgs STORM_PLM e Ibuy inviabiliza buscas relacionais repetidas em runtime. Avaliar a arvore logica diretamente em memoria via Maps estaticos duradouros garante < 50ms por registro sem consumir limites de CPU sincrono. Alternativa descartada: queries SOQL consecutivas (estouro de CPU).\n\nDECISAO 2 — Async-first via Queueables encadeados (vs. processamento sincrono)\nJustificativa: Strings de ate 6M caracteres de layouts CSV corporativos estourariam heap e timeout sincrono em qualquer arquitetura convencional. O encadeamento com cursor byte-offset fragmenta de forma resiliente e retomavel. Alternativa descartada: processamento sincrono no controller LWC (estouro fatal).\n\nDECISAO 3 — Einstein Prompt Templates + ConnectApi (vs. compilacao manual)\nJustificativa: Traduzir regras textuais DSL pt-BR para AST JSON via LLM elimina redeploys e democratiza a governanca de catalogo para analistas de negocios. Alternativa descartada: codificacao manual no BRE legado (dependencia cronica de TI).\n\nDECISAO 4 — DLQ-first com Finalizers e ZombieReaper (vs. tolerancia a falha silenciosa)\nJustificativa: Processos batch assincronos sem mecanismos de captura geram transacoes zumbis invisiveis. A arquitetura DLQ-first com rollup M-D e idempotencia SHA-256 garante rastreabilidade 100%. Alternativa descartada: sem tratamento de chain-death (risco de corrompimento de dados).'
})

# ── SLIDE 20 — Principios ────────────────────────────────────────────────────
sl = prs.slides[19]
replace_all(sl, {
    'Exemplos:\nCliques, Não Código: Daremos prioridade à configuração declarativa para acelerar a entrega, reduzir custos de manutenção e capacitar sua equipe a fazer atualizações futuras.\nAbordagem Primeiro o Padrão: Ao aderir à funcionalidade pronta para uso, garantimos que sua plataforma seja escalável, segura e pronta para se beneficiar de 3 atualizações do Salesforce por ano.\nArquitetura Liderada por API: Projetaremos integrações para serem ativos reutilizáveis, prevenindo dívidas técnicas e criando uma base ágil para projetos futuros.':
    '1. IA e SEGURANCA EM PRIMEIRO LUGAR: Todo trafego dos agentes passa obrigatoriamente pelo Einstein Trust Layer — dados confidenciais da Claro nao sao expostos ou retidos por LLMs externos.\n\n2. RESILIENCIA ATIVA COMO PADRAO: Nenhum processo batch roda sem Transaction Finalizers, DLQ e ZombieReaper ativos — observabilidade total e zero falhas ocultas.\n\n3. CONFIGURACAO DECLARATIVA + ENGENHARIA FOCADA: Prompt Templates e capacidades nativas do Einstein sao priorizados. Codigo Apex customizado e restrito estritamente ao motor AST Walker de alta performance.\n\n4. GOVERNANCA UNIFICADA DE DADOS: A Claro assume responsabilidade total pelo saneamento previo dos Knowledge Articles — o ecossistema de IA opera sobre taxonomia limpa e Data Categories configuradas.\n\n5. ENTREGA INCREMENTAL SEM SURPRESAS: Cronograma de 8 semanas fixo e inegociavel. UAT exaustivo com Parallel Run validando vereditos vs. BRE legado antes de qualquer cutover. Zero big bang.'
})

# ── SLIDE 22 — Resumo do Escopo ──────────────────────────────────────────────
sl = prs.slides[21]
replace_all(sl, {
    'Resumo do Escopo do Projeto...':
    'RESUMO DO ESCOPO\n\nIniciativa: POC PLM & Agentforce — Validacao Autonoma de Catalogo de Produtos\nCliente: Claro Brasil | Duracao: 8 semanas | Modelo: T&M | Sizing: SC (Super Complex)\n\nEsta POC implementa uma arquitetura async-first com dois agentes Agentforce autonomos (Admin e Ops), um interpretador deterministico Apex AST Walker e mecanismos de resiliencia ativa para substituir o catalogo BRE legado de 127 regras pela governanca inteligente via Einstein LLM.\n\nProdutos Salesforce: Agentforce (Employee Agents) · Einstein 1 Platform (Prompt Templates) · Core Platform (Apex, LWC, Queueable, DLQ)\n\nEscopo: 3 tipos de produto (Fone, BL, TV) · 3 canais de severidade (ERRO, AVISO, INFO) · 2 agentes autonomos · CSV ate 10.000 linhas sincronas / 6M caracteres'
})

# ── SLIDE 23 — Atividades e Entregaveis ──────────────────────────────────────
sl = prs.slides[22]
for shape in sl.shapes:
    if shape.has_text_frame:
        t = shape.text_frame.text.strip()
        if t == 'Atividades':
            set_shape_text(shape, [
                'Atividades por Fase:',
                'F0 — Discovery & Design (Sem. 1-2): Alinhamento de regras DSL, imersao tecnica com Lucas/Fabricio, definicao do schema JSON de especificacoes, congelamento de historias de usuario.',
                'F1 — Build & Sprints (Sem. 3-5): Codificacao do AST Walker Apex, LWC DemandaCsvWizard, Queueables encadeados, Agentforce Admin e Ops, Prompt Templates, DLQ e ZombieReaper.',
                'F2 — UAT & Fine Tuning (Sem. 6-7): Testes de concorrencia, massa de dados CSV real, validacao de vereditos vs. BRE legado (Parallel Run), ajustes de prompts e scripts de homologacao.',
                'F3 — Deploy & Hypercare (Sem. 8): Deploy via RunSpecifiedTests, Go-Live, hiperatendimento pos-producao e cutover formal.',
            ], size=8, bold_first=True)
        elif t == 'Entregáveis':
            set_shape_text(shape, [
                'Entregaveis:',
                'Agentes Agentforce configurados: Admin (4 topicos compile-time) + Ops (6 topicos runtime) com planner bundles, GenAI plugins e Prompt Templates ativos.',
                'Motor Apex AST Walker (PlmRuleSpecEvaluator): validacao < 50ms, 3 tipos de produto, 3 canais de severidade, DLQ-first.',
                'Componente LWC DemandaCsvWizard: importacao assincrona de CSV com cursor byte-offset e CPU guard.',
                'Infraestrutura de resiliencia: Transaction Finalizers, DLQ, Snapshots de compilacao, PimZombieReaperService.',
                'Relatorio narrativo HTML (Diagnostico de Demanda) gerado pelo Agente Ops.',
                'Permset Validacao_Engine_Access + documentacao tecnica de handoff.',
            ], size=8, bold_first=True)

# ── SLIDE 24 — Dentro / Fora do Escopo ───────────────────────────────────────
sl = prs.slides[23]
for shape in sl.shapes:
    if shape.has_text_frame:
        t = shape.text_frame.text.strip()
        if t == 'Dentro do Escopo':
            set_shape_text(shape, [
                'DENTRO DO ESCOPO:',
                '2 agentes Agentforce autonomos: Admin (compile-time) e Ops (runtime)',
                '3 tipos de produto validados: Fone, BL, TV',
                '3 canais de severidade: ERRO (bloqueia Pre_Aprovado), AVISO, INFO',
                'Motor AST Walker Apex puro com tempo < 50ms por registro',
                'Ingestao CSV assincrona: ate 10.000 linhas sincronas / 6M chars',
                'Compilacao de regras DSL pt-BR -> JSON via Einstein LLM',
                'Infraestrutura DLQ + Finalizers + ZombieReaper',
                'Relatorio narrativo HTML via Agente Ops',
                'Permset Validacao_Engine_Access + objetos Demanda__c e Item_Demanda__c',
                'Deteccao de drift por Source_Hash__c + recompilacao sem redeploy',
            ], size=8, bold_first=True)
        elif t == 'Fora do Escopo':
            set_shape_text(shape, [
                'FORA DO ESCOPO:',
                'Ingestao Bulk API 2.0 para CSVs > 6MB (backlog W3.1)',
                'Avaliador Batchable para demandas > 50.000 itens (backlog W3.2)',
                'Roll-ups Apex customizados para campos summary do Demanda__c (backlog W3.3)',
                'Sweep automatico de import-state CSV (Running -> Failed em heartbeat obsoleto) (backlog W3.5)',
                'Hardening avancado: split de permset Operator vs. Admin, caps regex (backlog W3.4/W3.8)',
                'Platform Cache como tier de cache (proibido em STORM_PLM e Ibuy)',
                'Multi-org / packaging / namespace gerenciado',
                'Integracao cross-system (MuleSoft, APIs externas, CDC) — exceto NBO condicional para Agente 3',
                'Criacao de novos Knowledge Articles do zero (responsabilidade da Claro)',
                'Tipos de produto alem de Fone, BL e TV',
            ], size=8, bold_first=True)

# ── SLIDE 25 — Integracoes ────────────────────────────────────────────────────
sl = prs.slides[24]
for shape in sl.shapes:
    if shape.has_table:
        tbl = shape.table
        integrations = [
            ('Einstein Platform (LLM Core)', 'Agentforce Admin <-> Einstein', 'ConnectApi Nativo', 'REST Apex Core API (Bidirecional)', 'JSON Spec payload de regras logicas DSL -> AST'),
            ('Sistema Legado NBO Externo', 'Agente 3 -> NBO (condicional)', 'Named Credentials', 'RESTful API HTTP JSON (Mao unica)', 'Parametros de elegibilidade e propensao de planos'),
            ('Service Cloud Knowledge Base', 'Agente Ops <- Knowledge Articles', 'Salesforce OOTB', 'Internal Knowledge API', 'FAQs saneadas com Data Categories configuradas'),
        ]
        for r, (sys, dir_, mw, api, obj) in enumerate(integrations, 1):
            if r < len(tbl.rows):
                set_cell(tbl, r, 0, sys, bold=True, size=8)
                set_cell(tbl, r, 1, dir_, size=8)
                set_cell(tbl, r, 2, mw, size=8)
                set_cell(tbl, r, 3, api, size=8)
                set_cell(tbl, r, 4, obj, size=8)
        break

# ── SLIDE 26 — Requisitos e Premissas ────────────────────────────────────────
sl = prs.slides[25]
for shape in sl.shapes:
    if shape.has_text_frame:
        t = shape.text_frame.text.strip()
        if t == 'Requisitos':
            set_shape_text(shape, [
                'REQUISITOS (dependencias do cliente):',
                'Sandboxes estaveis e funcionais disponibilizadas pela Claro para build e testes isolados.',
                'Licencas Agentforce Unlimited e creditos Einstein provisionados no Dia 1 do build.',
                'Fornecimento de 10-15 FAQs prioritarias saneadas com Data Categories antes do inicio do build.',
                'Arquivos CSV no layout canonico de Carga (UTF-8, contrato de aceite documentado).',
                'Participacao ativa de Lucas (SME tecnico) e analistas de negocios nos ciclos de UAT.',
                'Definicao formal do papel de PM (compartilhado com Luciano/SWE) na Semana 1.',
            ], size=8, bold_first=True)
        elif t == 'Pressupostos':
            set_shape_text(shape, [
                'PREMISSAS E RESTRICOES:',
                'Platform Cache e proibido em STORM_PLM e Ibuy — cache via Maps estaticos + Custom Metadata.',
                'Quick-deploy rejeitado cross-org (CannotQuickDeployError) — deploys via RunSpecifiedTests.',
                'Campos Master-Detail e required NAO podem ter fieldPermissions explicitados em permsets (auto-granted).',
                'Agente 3 (NBO) e condicional: depende de documentacao estavel da API do sistema externo.',
                'Cronograma de 8 semanas e fixo e inegociavel — extensoes requerem aprovacao formal.',
                'Einstein/Prompt Templates exigem applicationName=PromptTemplateGenerationsInvocable em todos os generateMessagesForPromptTemplate.',
            ], size=8, bold_first=True)

# ── SLIDE 27 — Roteiro ───────────────────────────────────────────────────────
sl = prs.slides[26]
replace_all(sl, {
    'Fase 1\nAgente de Suporte a Negócios para Análise de Viabilidade Comercial & Agente de Catálogo de Produtos para Consulta Técnica e Comercial de Produtos':
    'Fase 0 — Discovery & Design\nSemanas 1-2\nAlinhamento de regras DSL, imersao com Lucas/Fabricio, definicao schema JSON, congelamento de historias, arquitetura finalizada.',
    'Fase 2\nxxx':
    'Fase 1 — Build & Sprints\nSemanas 3-5\nAST Walker, LWC Wizard, Queueables, Agentforce Admin e Ops, Prompt Templates, DLQ, ZombieReaper.',
    'Fase 3\nxxx':
    'Fase 2 — UAT & Fine Tuning\nSemanas 6-7\nTestes de concorrencia, Parallel Run vs BRE legado, ajuste de prompts, validacao de criterios de aceite com usuarios Claro.',
    'Início': 'Jun/2026',
    'Objetivo': 'Fase 3 — Deploy & Hypercare\nSemana 8\nDeploy via RunSpecifiedTests, Go-Live, hiperatendimento, cutover formal e desativacao do BRE legado.',
})

# ── SLIDE 28 — Roteiro Detalhado ─────────────────────────────────────────────
sl = prs.slides[27]
replace_all(sl, {
    'Fase 1: Visão 360 Fundamentada (Semanas 1-8)\nResultado: Seus agentes de serviço terão uma visão unificada do histórico do cliente, permitindo um atendimento mais informado e personalizado desde o primeiro dia.':
    'Fase 0: Discovery & Design (Semanas 1-2)\nResultado: Arquitetura finalizada, schema JSON de especificacoes aprovado, catalogo de regras DSL mapeado, historias de usuario congeladas. Base solida para o build sem surpresas.',
    'Fase 2: Automação Proativa de Vendas (Semanas 9-16)\nResultado: Sua equipe de vendas se beneficiará do roteamento automático de leads e gestão de tarefas, liberando tempo para se concentrar em vendas.':
    'Fase 1: Build & Sprints (Semanas 3-5)\nResultado: Motor AST Walker codificado e testado, LWC de importacao CSV funcional, Queueables encadeados resilientes, Agentes Admin e Ops com topicos e acoes ativos, Prompt Templates configurados e DLQ operacional.',
    'Fase 3: Atendimento ao Vivo & Realização de Valor (Semana 17+)\nResultado: A solução completa está ao vivo, com acompanhamento contínuo do valor em relação aos KPIs definidos no início de nosso envolvimento.':
    'Fases 2 e 3: UAT, Fine Tuning & Go-Live (Semanas 6-8)\nResultado: Testes de estresse (CSV 6M chars), Parallel Run vs. BRE legado com paridade de vereditos validada, UAT assinado pelos POs Claro, deploy controlado via RunSpecifiedTests em producao, 127 regras BRE legadas desativadas formalmente.',
})

# ── SLIDE 30 — Cronograma ─────────────────────────────────────────────────────
sl = prs.slides[29]
replace_all(sl, {
    '[Insira um gráfico de cronograma de alto nível mostrando as principais fases e suas durações estimadas.]\nFase 1: Fundação (Semanas 1-8) Descoberta, Design, Configuração da Plataforma Central\nFase 2: Automação de Vendas (Semanas 9-16) Construir, Testar, Implantar Processo de Vendas\nFase 3: Serviço & Suporte (Semanas 17-24) Construir, Testar, Implantar Consola de Serviços':
    'CRONOGRAMA — CLARO AGENTE PLM POC (8 SEMANAS)\n\n|-- Semanas 1-2   | FASE 0: Discovery & Design\n|                | Alinhamento de regras DSL | Schema JSON | Historias de usuario congeladas\n|                | PRE-REQUISITO: 10-15 FAQs saneadas + sandboxes disponiveis\n|\n|-- Semanas 3-5   | FASE 1: Build & Sprints\n|                | AST Walker Apex | LWC Wizard CSV | Queueables | Agentforce Admin+Ops | DLQ\n|                | PRE-REQUISITO: Licencas Agentforce ativas no Dia 1\n|\n|-- Semanas 6-7   | FASE 2: UAT & Fine Tuning\n|                | Testes de estresse CSV | Parallel Run vs BRE legado | Ajuste de prompts | Aceite\n|\n|-- Semana 8      | FASE 3: Deploy & Hypercare\n                 | RunSpecifiedTests | Go-Live | Cutover formal | Desativacao BRE legado\n\nLACANA CRITICA: Status do saneamento de FAQs e documento API do NBO externo (Agente 3) devem ser confirmados na Semana 1.'
})

# ── SLIDE 41 — Equipe ─────────────────────────────────────────────────────────
sl = prs.slides[40]
replace_all(sl, {
    'Nossa Equipe de Projeto Combinada': 'Nossa Equipe de Projeto Combinada',
    'Equipo Projeto':
    'SALESFORCE PROFESSIONAL SERVICES\n\nEngagement Manager\n-> Coordenacao geral, governanca e interface com o cliente\n\nTechnical Architect (1x dedicado)\n-> Arquitetura AST Walker, Queueables, DLQ, devops cross-org\n\nTechnical Consultant (1x dedicado)\n-> Build dos agentes Agentforce, Prompt Templates, LWC Wizard\n\nQA Consultant (1.5x)\n-> Estrategia de testes, Parallel Run, UAT, RunSpecifiedTests\n\n-----------------------\nCLARO BRASIL (Cliente)\n\nLucas — SME Tecnico (Key Subject Matter Expert)\nLuciano — Lider SWE / PM compartilhado\nFabricio — Lider de Operacoes / Patrocinador Operacional\nAnalistas de Catalogos (Fone, BL, TV)\nEquipe de Governanca de Knowledge Articles'
})

# ── SLIDE 42 — RACI ───────────────────────────────────────────────────────────
sl = prs.slides[41]
for shape in sl.shapes:
    if shape.has_table:
        tbl = shape.table
        raci_data = [
            ('Definicao do Schema JSON de Especificacoes', 'Facilita & Documenta', 'Fornece Especialistas (Lucas SME) & Aprova'),
            ('Build Agentforce Admin + Ops (topicos, acoes)', 'Projeta & Configura', 'Valida no UAT & Fornece Feedback'),
            ('Saneamento de FAQs e Knowledge Articles', 'Orienta Taxonomia e Data Categories', 'Responsavel exclusivo pela execucao e entrega'),
            ('Testes de Estresse CSV e Parallel Run', 'Executa e Corrige Defeitos', 'Fornece arquivos CSV reais & Valida vereditos'),
            ('Deploy via RunSpecifiedTests (Cutover)', 'Executa Deploy Controlado', 'Aprova e Assina Aceite Formal de Go-Live'),
        ]
        for r, (act, sf, client) in enumerate(raci_data, 1):
            if r < len(tbl.rows):
                set_cell(tbl, r, 0, act, size=8)
                set_cell(tbl, r, 1, sf, size=8)
                set_cell(tbl, r, 2, client, size=8)
        break

# ── SLIDE 50 — Investimento ───────────────────────────────────────────────────
sl = prs.slides[49]
for shape in sl.shapes:
    if shape.has_table:
        tbl = shape.table
        set_cell(tbl, 1, 0, 'Honorarios Estimados de Servicos Profissionais', bold=True, size=9)
        set_cell(tbl, 1, 1, 'A confirmar via SOW formal | Sizing: SC | 8 semanas T&M | 1 TA + 1 TC + 1.5 QA', size=9)
        set_cell(tbl, 2, 0, 'Modelo de Contratacao', bold=True, size=9)
        set_cell(tbl, 2, 1, 'Tempo & Materiais (T&M)', size=9)
        set_cell(tbl, 3, 0, 'Despesas de Viagem & Deslocamento (T&E)', bold=True, size=9)
        set_cell(tbl, 3, 1, 'Nao incluidas nos honorarios — a definir conforme presenca on-site necessaria', size=9)
        break

replace_all(sl, {
    'Validade: Esta Estimativa Bruta de Magnitude é válida até.':
    'Validade: Esta Estimativa Bruta de Magnitude e valida por 30 dias a partir da data de entrega.',
})

# ── SLIDE 52 — Proximos Passos ────────────────────────────────────────────────
sl = prs.slides[51]
replace_all(sl, {
    '1. Revisão e Alinhamento do ROM\nAtividade: Revisar em conjunto esta apresentação para alinhar escopo, abordagem e investimento.\nResultado: Acordo mútuo para prosseguir.':
    '1. Revisao e Alinhamento do ROM\nAtividade: Revisar esta apresentacao com Lucas, Luciano e Fabricio para alinhar escopo tecnico, premissas de infraestrutura (orgs STORM_PLM/Ibuy) e investimento.\nResultado: Acordo mutuo para prosseguir e confirmacao do papel de PM.',
    '2. Finalização do Escopo & SOW\nAtividade: Refinar os detalhes do escopo e desenvolver a Declaração Formal de Trabalho (SOW).\nResultado: Um SOW assinado.':
    '2. Finalizacao do Escopo & SOW\nAtividade: Confirmar documentacao da API NBO (Agente 3), status do saneamento de FAQs e disponibilidade de sandboxes. Formalizar papel de PM e SOW.\nResultado: SOW assinado com cronograma de 8 semanas travado.',
    '3. Início do Projeto\nAtividade: Início formal do projeto com a equipe do projeto combinada.\nResultado: Nossa jornada de transformação começa.':
    '3. Inicio do Projeto\nAtividade: Kick-off formal com Lucas, Luciano, Fabricio e time PS. Inicio da Fase 0 — Discovery & Design com sessoes de alinhamento de regras DSL e definicao do schema JSON.\nResultado: Arquitetura finalizada e historias de usuario congeladas ao final da Semana 2.',
})

# ── Logo placeholders ─────────────────────────────────────────────────────────
for sl in prs.slides:
    replace_all(sl, {'Inserir LOGO do Cliente': 'Claro Brasil'})

prs.save(out)
print(f'Salvo: {out}')
print(f'Total slides: {len(prs.slides)}')
